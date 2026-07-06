#!/usr/bin/env python
"""FCC更新停止者の判定技術 — 社内 patent review 用サマリ deck v5 (.pptx).

v4 (build_v4_pptx.py) を土台に、非エンジニア向けに徹底的に作り込んだ改訂版:
  * 基礎知識パート新設（バッテリと燃料計 / FCC学習と凍結 / ビジネス上の問題）
  * 仕組み図解 3 枚新設（学習機会とEND起点72h監査 / dual-track非対称リセット /
    有界保持の因果台帳）
  * 全コンテンツスライドに「💡ひとことで」帯（非エンジニア向け要約）
  * 図解説スライドは v4 の 4 ブロック（何のグラフ/軸/主張/読み方）を、実画像の
    検分結果（軸ラベル・棒の有無・凡例の重なり等）に合わせて正確化
  * 正直な開示 2 枚（敵対的レビューの必須開示 / 限界と残課題）
  * 知財戦略オプション（6/18 特許性評価と v4 エビデンス後の両見解を併記）
  * 付録拡充（データと再現性 / 代替実施形態 / 先行技術差別化 / 補足図 4 枚）

数値は data/processed/fcc_patent_evidence_v4/_v4_results_summary.json（ビルド時
読込）と、一次データで照合済みの件数（online_latest_snapshot_v2.csv /
fcc_final_action_labels.csv）から引用。NOT a legal opinion / 先行技術 UNVERIFIED /
介入・FWバージョンデータ NOT AVAILABLE / 捏造なし。
"""
from __future__ import annotations

import json
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.oxml.ns import qn

from battery_usage import patent_common_v4 as pc

R = json.load(open(pc.V4_DIR / "_v4_results_summary.json", encoding="utf-8"))
A2, A3, B, C2, C3, D, DM, E_ = (R["A2"], R["A3"], R["B"], R["C2"], R["C3"],
                                R["D"], R["Dmin"], R["E"])
FIG = pc.FIG_DIR
OUT = pc.REPORTS / "fcc_patent_summary_slides_v5.pptx"

# ---- palette ------------------------------------------------------------- #
NAVY = RGBColor(0x1F, 0x37, 0x5B); BLUE = RGBColor(0x1F, 0x77, 0xB4)
STEEL = RGBColor(0x33, 0x66, 0x99); GREY = RGBColor(0x55, 0x55, 0x55)
RED = RGBColor(0xC0, 0x39, 0x2B); GREEN = RGBColor(0x2C, 0xA0, 0x2C)
DGREEN = RGBColor(0x37, 0x6B, 0x3A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7); ORANGE = RGBColor(0xB0, 0x5A, 0x00)
TEAL = RGBColor(0x2E, 0x6E, 0x6A); PURPLE = RGBColor(0x6A, 0x3D, 0x8A)
BROWN = RGBColor(0x6B, 0x4A, 0x1F); INK = RGBColor(0x22, 0x22, 0x22)
YELLOW_BG = RGBColor(0xFF, 0xF6, 0xDA); YELLOW_LN = RGBColor(0xD9, 0xA4, 0x2A)
LGREEN_BG = RGBColor(0xE4, 0xF0, 0xE4); LBLUE_BG = RGBColor(0xE2, 0xEC, 0xF7)
LRED_BG = RGBColor(0xF7, 0xE4, 0xE2); LGREY_BG = RGBColor(0xEC, 0xEC, 0xEC)

FONT = "Meiryo UI"
DISC = ("技術的特許性エビデンス（NOT a legal opinion）。先行技術は UNVERIFIED。"
        "特許化の粒度・範囲の最終判断は社内reviewer/弁理士。捏造（地上真実/介入/FW/因果結論）なし。")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def EM(v):
    """Coerce any Length/float arithmetic result to integer EMU (PowerPoint が
    float 座標を拒否するため)。"""
    return Emu(int(round(v)))


def _cjk(run):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", FONT)


def _style(run, size=14, bold=False, color=None):
    run.font.size = Pt(size); run.font.bold = bold; run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    _cjk(run)


def _tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(EM(left), EM(top), EM(width), EM(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size=14, bold=False, color=None, align=None):
    p.text = text
    for run in p.runs:
        _style(run, size, bold, color)
    if align is not None:
        p.alignment = align


def footer(slide):
    tf = _tb(slide, Inches(0.3), SH - Inches(0.38), SW - Inches(1.1), Inches(0.32))
    _set(tf.paragraphs[0], DISC, size=8, color=GREY)


def band(slide, color, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EM(SW), EM(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def section(title, accent=BLUE, tag=None):
    s = prs.slides.add_slide(BLANK)
    band(s, accent, Inches(0.86))
    tf = _tb(s, Inches(0.5), Inches(0.12), SW - Inches(3.2), Inches(0.64))
    _set(tf.paragraphs[0], title, size=21, bold=True, color=WHITE)
    if tag:
        tt = _tb(s, SW - Inches(3.2), Inches(0.2), Inches(2.9), Inches(0.5))
        _set(tt.paragraphs[0], tag, size=11.5, bold=True,
             color=RGBColor(0xDD, 0xE8, 0xF2), align=2)
    return s


def takeaway(slide, text, y=None):
    y = Inches(6.5) if y is None else y
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 EM(Inches(0.35)), EM(y),
                                 EM(SW - Inches(0.7)), EM(Inches(0.52)))
    shp.fill.solid(); shp.fill.fore_color.rgb = YELLOW_BG
    shp.line.color.rgb = YELLOW_LN; shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    _set(tf.paragraphs[0], "💡 " + text, size=11.5, bold=True, color=BROWN)


def box(slide, left, top, w, h, text, fill, fg=WHITE, size=12, sub=None,
        sub_size=9.5, line=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 EM(left), EM(top), EM(w), EM(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.color.rgb = line
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    _set(tf.paragraphs[0], text, size=size, bold=True, color=fg, align=2)
    if sub:
        p = tf.add_paragraph(); _set(p, sub, size=sub_size, color=fg, align=2)
    return shp


def line_seg(slide, x1, y1, x2, y2, color=NAVY, width=2.25, dash=None,
             arrow_head=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   EM(x1), EM(y1), EM(x2), EM(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    if dash is not None:
        c.line.dash_style = dash
    if arrow_head:
        ln = c.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return c


def label(slide, x, y, w, h, text, size=10.5, bold=False, color=INK, align=None):
    tf = _tb(slide, x, y, w, h)
    _set(tf.paragraphs[0], text, size=size, bold=bold, color=color, align=align)
    return tf


def bullets(tf, items, size=12.5, gap=4, color=INK, first=True, marker="•  "):
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set(p, (marker + it) if marker else it, size=size, color=color)
        p.space_after = Pt(gap)
    return tf


def _fill_cell(c, text, size=10, bold=False, color=INK, fill=None):
    c.text = text
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    for pa in c.text_frame.paragraphs:
        for run in pa.runs:
            _style(run, size, bold, color)


def make_table(slide, headers, data, x, y, colw, row_h=0.42, hdr_size=11,
               cell_size=10, styler=None):
    n = len(data) + 1
    tbl = slide.shapes.add_table(n, len(headers), EM(x), EM(y),
                                 EM(Inches(sum(colw))),
                                 EM(Inches(row_h) * n)).table
    for j, w in enumerate(colw):
        tbl.columns[j].width = EM(Inches(w))
    for j, h in enumerate(headers):
        _fill_cell(tbl.cell(0, j), h, size=hdr_size, bold=True, color=WHITE,
                   fill=NAVY)
    for i, rowv in enumerate(data, 1):
        for j, v in enumerate(rowv):
            size, bold, color = cell_size, False, INK
            if styler:
                size, bold, color = styler(i, j, v, cell_size)
            _fill_cell(tbl.cell(i, j), v, size=size, bold=bold, color=color,
                       fill=LIGHT if i % 2 else WHITE)
    return tbl


def status_styler(hl_cols):
    def f(i, j, v, base):
        bold, color = False, INK
        if j in hl_cols:
            bold = True
            if "STRONG" in v:
                color = GREEN
            elif v.strip() in ("HIGH",) or v.startswith("HIGH"):
                color = RED
            elif "PROSPECTIVE" in v or "WEAK" in v or v == "—":
                color = GREY
            elif "NOT" in v:
                color = RED
            else:
                color = BLUE
        return base, bold, color
    return f


def glossary(title, items, tag="用語定義"):
    s = section(title, STEEL, tag)
    tf = _tb(s, Inches(0.55), Inches(1.0), SW - Inches(1.1), Inches(5.3))
    first = True
    for term, defn in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(4)
        r = p.add_run(); r.text = f"{term}　"; _style(r, 12.5, True, NAVY)
        r2 = p.add_run(); r2.text = defn; _style(r2, 11.5, False, INK)
    return s


def concept(title, tag, strength, risk, bg, prob, chal, sol, novelty, tw,
            accent=TEAL):
    s = section(title, accent, tag)
    tf = _tb(s, Inches(0.5), Inches(0.94), SW - Inches(1.0), Inches(0.4))
    p = tf.paragraphs[0]; p.text = ""
    for txt, col in [("技術エビデンス強度: ", NAVY),
                     (strength + "    ",
                      GREEN if "STRONG" in strength
                      else (GREY if "PROSPECTIVE" in strength else BLUE)),
                     ("先行技術リスク(UNVERIFIED): ", NAVY),
                     (risk, RED if "HIGH" in risk else GREY)]:
        r = p.add_run(); r.text = txt; _style(r, 12, True, col)
    tf = _tb(s, Inches(0.5), Inches(1.36), SW - Inches(1.0), Inches(5.0))
    first = True
    for head, body in [("背景", bg), ("問題", prob), ("課題", chal),
                       ("解決法", sol), ("特許性 / 新規性", novelty)]:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        _set(p, head, size=13, bold=True, color=NAVY); p.space_before = Pt(4)
        q = tf.add_paragraph(); _set(q, body, size=11.8)
    takeaway(s, tw)
    footer(s)
    return s


def features(title, tag, rows, row_h=0.44):
    s = section(title, RGBColor(0x2E, 0x5A, 0x88), tag)
    make_table(s, ["変数 / 特徴量", "種別", "作成方法（特徴量エンジニアリング）"],
               rows, Inches(0.4), Inches(1.02), [3.4, 1.35, 7.75], row_h=row_h,
               styler=lambda i, j, v, b: (
                   b, j == 1,
                   (GREEN if v == "派生" else RED if "NOT" in v else GREY)
                   if j == 1 else INK))
    footer(s)
    return s


def figure(title, tag, img, what, axes, claim, interp, tw, accent=DGREEN,
           max_w=7.1, max_h=4.95):
    s = section(title, accent, tag)
    p = FIG / img
    img_bottom = Inches(1.1)
    if p.exists():
        from PIL import Image
        try:
            iw, ih = Image.open(p).size
        except Exception:
            iw, ih = 1600, 1000
        scale = min(Inches(max_w) / iw, Inches(max_h) / ih)
        pic_w, pic_h = EM(iw * scale), EM(ih * scale)
        s.shapes.add_picture(str(p), EM(Inches(0.4)), EM(Inches(1.08)),
                             width=pic_w, height=pic_h)
        img_bottom = Inches(1.08) + pic_h
        cap = _tb(s, Inches(0.4), img_bottom + Inches(0.02), Inches(7.1),
                  Inches(0.3))
        _set(cap.paragraphs[0], f"図: {img}（dpi=300・匿名化済み）", size=8.5,
             color=GREY)
    tf = _tb(s, Inches(7.7), Inches(1.0), Inches(5.3), Inches(5.45))
    first = True
    for head, body, col in [("【何のグラフか】", what, BLUE),
                            ("【軸・変数の定義】", axes, NAVY),
                            ("【主張（言いたいこと）】", claim, GREEN),
                            ("【読み方・解釈】", interp, ORANGE)]:
        p1 = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        _set(p1, head, size=12, bold=True, color=col); p1.space_before = Pt(5)
        p2 = tf.add_paragraph(); _set(p2, body, size=11, color=INK)
    takeaway(s, tw)
    footer(s)
    return s


RAW, ENG, NA = "raw", "派生", "NOT AVAIL"

# =========================================================================== #
# 1. Title
# =========================================================================== #
s = prs.slides.add_slide(BLANK); band(s, NAVY, SH)
tf = _tb(s, Inches(0.8), Inches(1.35), SW - Inches(1.6), Inches(2.2))
_set(tf.paragraphs[0], "テレメトリからの FCC更新停止（ゲージ凍結）者の判定技術",
     size=33, bold=True, color=WHITE)
p = tf.add_paragraph()
_set(p, "— と、その実現過程で生まれた特許候補アイデア — （社内 patent review 用 v5）",
     size=19, color=RGBColor(0xCF, 0xDD, 0xEE))
tf2 = _tb(s, Inches(0.8), Inches(3.55), SW - Inches(1.6), Inches(3.0))
for i, ln in enumerate([
        "本題: ノートPCバッテリのテレメトリ(RSOC/FCC/サイクル/時刻)だけで『満充電容量(FCC)を再学習しなくなった個体』を判定する",
        "母集団: 実バッテリ履歴 752ユーザー / 3,130,394 サンプル / 全期間 24,711 学習機会(エピソード)",
        "特許候補は 3つ(+将来2つ) に集約。どこをどの粒度で特許化するかの最終判断は社内reviewerに委ねる",
        "技術効果は proxy ラベルに依存しない独立指標で検証。法的結論は主張しない",
        "本版(v5)は非エンジニア向けの基礎解説・仕組み図解・正直な開示を追加した拡充版"]):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    _set(p, "•  " + ln, size=14.5, color=WHITE); p.space_after = Pt(5)
footer(s)

# =========================================================================== #
# 2. Executive summary
# =========================================================================== #
s = section("エグゼクティブサマリ — 1枚で全体像", NAVY, "Summary")
tf = _tb(s, Inches(0.55), Inches(1.02), SW - Inches(1.1), Inches(5.3))
first = True
for head, body in [
    ("課題", "バッテリ健全性表示(SoH)を駆動する満充電容量(FCC)が長期間更新されない『ゲージ凍結』は、"
             "静的検査では『浅い使い方で正常』『要再較正』『FW/HW起因』を判別できない。752台中114台が凍結、うち70台は使用で説明不能。"),
    ("解決", "『凍結を予測』は失敗(AUC≈0.54=ほぼコイン投げ)。代わりに『再学習の機会(深放電→再充電)に対して実際に応答したか』を"
             "END起点72h窓で機械的に監査する方式へ転換。テレメトリのみ・機種非依存・欠測/打ち切りに頑健・30日保持でも全期間と等価。"),
    ("実績(実フリート752台)", "全履歴監査でゲージ再較正候補18台/FW確認候補14台に自動振り分け(既存active層からの誤エスカレーション0件)。"
             "30日オンライン運用版はFW/ゲージのCore判定で proxy 精度1.0・能動誤報0件。5本柱の独立検証(A2/A3/B/E/D)で技術効果を実証。"),
    ("特許候補", "①機会条件付き無応答監査(中核) ②デュアルトラック非対称リセット ③有界保持の因果証拠台帳 ＋ 将来④クローズドループ介入検証 "
             "⑤機種非依存スクリーニング+version局在。技術エビデンスは①〜③すべてSTRONG。ただし新規性リスクは別軸(全てUNVERIFIED)。"),
    ("reviewerへの依頼", "(1) 候補①〜③の特許化の粒度・範囲の判断 (2) 候補②の着想日・社内外公開有無の確認 "
             "(3) 正式なFTO/特許性調査(弁理士)の発注判断 (4) 将来④の介入データ収集開始の判断"),
]:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    _set(p, "■ " + head, size=13.5, bold=True, color=NAVY); p.space_before = Pt(6)
    q = tf.add_paragraph(); _set(q, body, size=12, color=INK)
takeaway(s, "『動くべき機会に動かなかった』事実だけを、欠測に騙されず数える監査技術。エビデンスは強いが、特許性(新規性)の判断はここから先が本番。")
footer(s)

# =========================================================================== #
# 3. 構成と読み方
# =========================================================================== #
s = section("本資料の構成と読み方", NAVY, "Guide")
parts = [
    ("§1 基礎知識", "バッテリと燃料計 / FCC学習と凍結 / なぜ問題か", STEEL, "非エンジニアの方はここから"),
    ("§2 用語定義", "バッテリ・学習機会 / 検証手法・統計", STEEL, "つまずいたら戻る辞書"),
    ("§3 本題(Part1)", "目的と難しさ → 失敗と発想転換 → 手法 → 実フリート結果", NAVY, "判定技術そのもの"),
    ("§4 根拠", "5本柱の検証 A2/A3/B/E/D (図解説つき)", DGREEN, "なぜ信じてよいか"),
    ("§5 特許候補(Part2)", "候補①〜③ + 将来④⑤ (背景/問題/課題/解決/特許性)", PURPLE, "reviewer判断の対象"),
    ("§6 開示とまとめ", "正直な開示 / 知財戦略オプション / 依頼事項", RED, "リスクと残課題"),
    ("§7 付録", "特徴量 / データと再現性 / 代替実施形態 / 先行技術 / 補足図", GREY, "深掘り用"),
]
y = Inches(1.1)
for name, desc, col, note in parts:
    box(s, Inches(0.5), y, Inches(2.5), Inches(0.62), name, col, size=12.5)
    label(s, Inches(3.2), y + Inches(0.03), Inches(6.3), Inches(0.6), desc,
          size=12, color=INK)
    label(s, Inches(9.7), y + Inches(0.03), Inches(3.3), Inches(0.6), note,
          size=11, color=GREY)
    y += Inches(0.72)
tf = _tb(s, Inches(0.5), y + Inches(0.05), SW - Inches(1.0), Inches(0.85))
_set(tf.paragraphs[0],
     "読み方: 各スライド下部の黄色い帯（💡）は非エンジニア向けの一言要約。図スライドは4ブロック"
     "（何のグラフ/軸/主張/読み方）構成。急ぐ方はスライド 2→13→14→20→36 の5枚で概観できます。",
     size=11.5, color=INK)
q = tf.add_paragraph()
_set(q, "この資料の正直さの記号 — UNVERIFIED: 先行技術は未検証 / NOT AVAILABLE: そのデータは存在しない(捏造しない) / "
        "proxy: 正解ラベルは本番システムの推定値(地上真実ではない)。", size=11.5, bold=True, color=RED)
footer(s)

# =========================================================================== #
# 4. 基礎① バッテリと燃料計
# =========================================================================== #
s = section("基礎① ノートPCバッテリと『燃料計』 — FCC・RSOC・SoH とは", STEEL, "§1 基礎知識")
make_table(
    s, ["用語", "意味", "身近なたとえ"],
    [["FCC（満充電容量）", "電池が『いま満タンでどれだけ入るか』の学習値(mWh)。劣化とともに減っていく",
      "何度も計り直して覚えておく『いまの満タン量』"],
     ["RSOC（相対残量）", "いまの残量% = 残容量 ÷ FCC × 100", "ガソリンメーターの針(満タン比の%)"],
     ["SoH（健全性）", "健康度% = FCC × 100 ÷ 設計容量(新品時)", "新品と比べて容量が何%残っているか"],
     ["FCC学習（再較正）", "『満充電付近→深く放電→再び満充電』の往復のときだけ、ゲージICがFCCを計り直せる",
      "体重計のゼロ点合わせ。乗らなければ校正できない"],
     ["ゲージ凍結", "FCCが長期間更新されない状態。SoH表示も止まる", "校正が止まり、表示が古いまま動かない"]],
    Inches(0.45), Inches(1.05), [2.6, 6.1, 3.7], row_h=0.62, cell_size=11)
tf = _tb(s, Inches(0.5), Inches(4.95), SW - Inches(1.0), Inches(1.4))
bullets(tf, [
    "電池の実容量は直接測れない。ゲージIC(燃料計チップ)が使用中の充放電から『推定』している",
    "FCCは整数mWhで量子化(最小刻み10mWh)されており、階段状にしか動かない",
    "SoH表示・残り時間表示・保証判断は、すべてこのFCC学習値に依存している"], size=12)
takeaway(s, "PCの電池残量・健康度表示は『実測』ではなく『燃料計チップの学習値』。学習が止まると表示だけが古いまま取り残される。")
footer(s)

# =========================================================================== #
# 5. 基礎② FCC学習とゲージ凍結（階段図）
# =========================================================================== #
s = section("基礎② FCC学習のしくみと『ゲージ凍結』", STEEL, "§1 基礎知識")
# --- staircase panel (left) ---
px0, px1 = Inches(0.7), Inches(6.7)
py0, py1 = Inches(1.6), Inches(5.7)
line_seg(s, px0, py1, px1 + Inches(0.2), py1, color=GREY, width=1.5, arrow_head=True)   # time axis
line_seg(s, px0, py1, px0, py0 - Inches(0.15), color=GREY, width=1.5, arrow_head=True)  # y axis
label(s, px1 - Inches(0.6), py1 + Inches(0.05), Inches(1.4), Inches(0.3), "時間 →", size=10, color=GREY)
label(s, px0 - Inches(0.25), py0 - Inches(0.45), Inches(2.2), Inches(0.3), "容量 (mWh)", size=10, color=GREY)
# true capacity: dashed diagonal
line_seg(s, px0 + Inches(0.2), py0 + Inches(0.35), px1, py0 + Inches(2.2),
         color=GREY, width=1.75, dash=MSO_LINE.DASH)
label(s, px1 - Inches(2.1), py0 + Inches(2.25), Inches(2.6), Inches(0.3),
      "実容量（劣化で緩やかに低下）", size=9.5, color=GREY)
# healthy gauge staircase (blue)
xs = [0.2, 1.3, 1.3, 2.6, 2.6, 3.9, 3.9, 5.2, 5.2, 6.0]
ys = [0.55, 0.55, 0.95, 0.95, 1.35, 1.35, 1.75, 1.75, 2.15, 2.15]
for i in range(len(xs) - 1):
    line_seg(s, px0 + Inches(xs[i]), py0 + Inches(ys[i]),
             px0 + Inches(xs[i + 1]), py0 + Inches(ys[i + 1]), color=BLUE, width=2.5)
label(s, px0 + Inches(3.15), py0 + Inches(0.62), Inches(3.4), Inches(0.3),
      "健全: 学習機会のたびFCCを更新（階段）", size=9.5, bold=True, color=BLUE)
# frozen gauge (red flat)
line_seg(s, px0 + Inches(0.2), py0 + Inches(0.55), px1 - Inches(0.5), py0 + Inches(0.55),
         color=RED, width=2.5)
label(s, px0 + Inches(1.2), py0 + Inches(0.12), Inches(4.6), Inches(0.35),
      "凍結: FCCが動かず SoH表示が実態から乖離", size=9.5, bold=True, color=RED)
# --- right text ---
tf = _tb(s, Inches(7.15), Inches(1.1), Inches(5.75), Inches(5.2))
first = True
for head, body in [
    ("学習機会（エピソード）", "RSOCが 高(≈80%)→低(≈20%)→高(≈80%) と往復する『深い放電→再充電』1サイクル。"
     "燃料計はこのときだけFCCを計り直せる。"),
    ("健全なゲージ", "学習機会のたびにFCCを更新し、実容量の低下を階段状に追いかける（左図・青）。"),
    ("凍結したゲージ", "FCCが更新されず横ばい（左図・赤）。SoH・残り時間・保証判断の元データが古いまま固定される。"),
    ("ここが難所", "『凍結』は故障とは限らない。浅い使い方で機会が無いだけかもしれない。"
     "凍結という結果だけ見ても原因(正常/要較正/FW疑い)は分からない。")]:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    _set(p, "■ " + head, size=12.5, bold=True, color=NAVY); p.space_before = Pt(5)
    q = tf.add_paragraph(); _set(q, body, size=11.5)
takeaway(s, "燃料計は『深く使って満充電に戻す』ときだけ計り直せる。動かない理由が『チャンスが無い』のか『チャンスに応えない』のかが本質的な問い。")
footer(s)

# =========================================================================== #
# 6. 基礎③ なぜ問題か
# =========================================================================== #
s = section("基礎③ なぜ問題か — 誤保守のコストとフリートの実態", STEEL, "§1 基礎知識")
tf = _tb(s, Inches(0.55), Inches(1.05), Inches(6.1), Inches(4.0))
_set(tf.paragraphs[0], "SoH表示が古いままだと（定性・金額効果は未算定）", size=13.5, bold=True, color=NAVY)
bullets(tf, [
    "劣化した電池を見逃す → 突然のシャットダウン・稼働時間の短縮・ユーザー満足度の低下",
    "健全な電池を誤って交換 → 部品・工数・回収物流の無駄",
    "保証・リース返却の判断がSoH表示に依存 → 表示が古いと判断そのものが歪む",
    "フリート管理(法人PC)では数千〜数万台規模で同じ歪みが累積する"], size=12, first=False)
box(s, Inches(7.0), Inches(1.15), Inches(5.8), Inches(1.05),
    "実フリート 752台の実態（実バッテリ履歴あり）", NAVY,
    sub="全1,808コレクション中、実履歴のある752台を全数解析（3,130,394サンプル）")
box(s, Inches(7.0), Inches(2.4), Inches(2.8), Inches(1.5), "凍結 114台\n(15.2%)", ORANGE,
    sub="60日以上FCC更新なし\n(stale 59 + very_stale 55)")
box(s, Inches(10.0), Inches(2.4), Inches(2.8), Inches(1.5), "使用で説明可\n44台", GREEN,
    sub="常時AC 22 / 低サイクル 16\n/ 浅放電 6")
box(s, Inches(7.0), Inches(4.1), Inches(5.8), Inches(1.45),
    "説明不能（HW/FW疑い）70台", RED,
    sub="activeより多くサイクルし(96.9 vs 65.7 cyc/yr)・深く放電(min RSOC=1)するのに凍結。"
        "行動・熱・ロギングでは説明不能。※『疑い』であり確定診断ではない(FWバージョンデータ不在)")
tf = _tb(s, Inches(0.55), Inches(5.2), Inches(6.1), Inches(1.1))
_set(tf.paragraphs[0],
     "→ 凍結は珍しくない(約15%)。しかも半分以上は使い方では説明できない。"
     "『どの個体に・どの対処を』を自動で仕分ける技術が必要。", size=12.5, bold=True, color=NAVY)
takeaway(s, "『表示が動かない』個体が752台中114台。原因はユーザー起因〜FW疑いまで様々で、対処も違う。見分ける仕組みがないと保守が空回りする。")
footer(s)

# =========================================================================== #
# 7-8. 用語定義
# =========================================================================== #
s = glossary("用語定義 (1/2) — バッテリと学習機会", [
    ("FCC（満充電容量）", "燃料計が学習する『今の満充電容量』(mWh)。劣化で減る。充放電で段階的に再学習・更新。"),
    ("FCC更新停止 / ゲージ凍結", "FCCが長期間更新されない状態。SoH(=FCC×100/DesignCapacity)が動かなくなる。本技術の判定対象。"),
    ("RSOC（相対残量, %）", "= remainingCapacity/FCC×100。今の充電残量%。high=満充電付近、low=空付近。"),
    ("学習機会 / エピソード", "燃料計がFCCを再学習できる『深い放電→再充電』1サイクル。RSOCが high→low→high と動いた区間。"),
    ("START / LOW / END", "エピソードの 開始(高RSOC) / 谷(最深放電) / 終了(再び高RSOC=再充電完了)。END=機会の完了時刻。"),
    ("有効ステップ（≥50mWh）", "意味あるFCC更新=『学習応答』とみなす最小変化。量子化最小は10mWh。"),
    ("micro-step（微小ステップ）", "50mWh未満の微小なFCC変化。ノイズとは呼ばない(microはeffectiveより反転しにくい)が、学習応答とは区別する。"),
    ("応答ステータス（END後72h窓）", "responded=有効step有 / no_response=窓を全部観測したが無し / censored=窓の観測が途中で終了=保留(無応答に数えない)。"),
    ("品質ティア", "エピソードの観測品質 HIGH_OK / MEDIUM_GAP / LOW_LARGE_GAP。睡眠ギャップ等が大きい機会は証拠として使わない。"),
])
takeaway(s, "『テスト(機会)を受けたか』『答案(応答)を出したか』『試験中に停電(欠測)しなかったか』を区別するための言葉たち。")
footer(s)

s = glossary("用語定義 (2/2) — 検証手法と統計", [
    ("対照 / ヌル分布(null)", "効果が無い場合に期待される値の分布。実データがその95%区間の外なら『偶然では説明できない』と判断。"),
    ("負の対照検定(A2)", "機会END時刻をわざとズラした『ニセ機会』5種と比較し、応答が真の機会に特異かを検証。治験のプラセボに相当。"),
    ("アンカー比較(A3)", "応答の起点を START/LOW/END で比較し『因果汚染(=再充電完了前の変化を応答と誤計上)』を定量化。"),
    ("応答ハザード(B)", "機会ENDからの経過時間に対する累積応答率(CIF)を生存解析で推定。"),
    ("欠測ストレス(E) / 保持グリッド(D)", "E=欠測・打ち切りを人工注入し誤判定耐性を検証 / D=生データ保持を短くしても全期間版と等価か検証。"),
    ("user-bootstrap", "ユーザ単位で再標本して信頼区間を計算。同一ユーザの複数機会を独立扱いしない(ばらつきの過小評価を防ぐ)。"),
    ("dual-track / 非対称リセット", "any(全変化)とeffective(≥50mWh)を別系統で追跡。微小変化はany系統のみリセットしeffective証拠は保持。"),
    ("stateful / stateless", "永続状態(台帳)あり/なし。statefulは保持窓をまたぐ証拠を状態で回収できる。"),
    ("AUC", "判別力の点数。1.0=完璧、0.5=コイン投げと同じ。0.54は『ほぼ偶然』の意味。"),
    ("proxyラベル", "地上真実(実際の故障確定)が無いため、別ロジック(本番システム)の出力を仮の正解として使ったもの。"),
])
takeaway(s, "検証はすべて『偶然・欠測・データ制約に騙されていないか』を突く設計。UNVERIFIED/NOT AVAILABLE/proxyの3語は正直さの目印。")
footer(s)

# =========================================================================== #
# 9. Part1 目的と難しさ
# =========================================================================== #
s = section("本来の目的と、なぜ難しいか", NAVY, "§3 Part1 ─ 本題")
tf = _tb(s, Inches(0.6), Inches(1.1), Inches(6.0), Inches(5.0))
_set(tf.paragraphs[0], "目的", size=15, bold=True, color=BLUE)
q = tf.add_paragraph()
_set(q, "テレメトリ(RSOC・FCC・サイクル・時刻)だけから、FCCを再学習しなくなった個体(ゲージ凍結)をフリート規模で正しく判定する。",
     size=12.5)
p = tf.add_paragraph(); _set(p, "なぜ難しいか", size=15, bold=True, color=BLUE)
p.space_before = Pt(10)
bullets(tf, [
    "静的検査では区別不能: FCCが動かない理由が『浅い充放電で正常』か『要再較正』か『FW/HW起因』か判別できない",
    "欠測・睡眠ギャップ・記録打ち切り: データの穴を『無応答』と誤判定しやすい",
    "生データ保持が有界: 直近30日しか残せない等の制約で、過去の証拠が失われる",
    "機種非依存が必須: 機種名でハードコードすると過学習・汎化不能"], size=12, first=False)
box(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(1.9),
    "FCCが動かない = ?", BLUE, sub="静的検査では下の3つを区別できない")
for i, (txt, col) in enumerate([("正常\n(浅充放電)", GREEN), ("要再較正\n(ゲージ)", ORANGE),
                                ("FW/HW\n起因", RED)]):
    box(s, Inches(7.0) + Inches(2.0) * i, Inches(3.7), Inches(1.85), Inches(1.3),
        txt, col, size=12)
takeaway(s, "『FCCが動かない』という結果は同じでも原因は3通り。原因ごとに正しい対処(何もしない/較正を促す/FW調査)が違う。")
footer(s)

# =========================================================================== #
# 10. 失敗した素朴アプローチ → 発想転換
# =========================================================================== #
s = section("失敗した素朴アプローチ → 発想の転換", NAVY, "§3 Part1 ─ 本題")
tf = _tb(s, Inches(0.6), Inches(1.05), Inches(8.2), Inches(2.5))
_set(tf.paragraphs[0], "『使用挙動から凍結を予測』しようとしたが、解けなかった:",
     size=14, bold=True, color=NAVY)
bullets(tf, [
    "教師あり予測(33特徴量, 公平比較領域 obs≥180d): AUC 0.535/0.540 — ほぼランダム",
    "FCC履歴を除いた規範ML(異常スコア): AUC 0.5584 — near-random",
    "最重要特徴 min_rsoc は『深放電ほど凍結』という反usage方向(交絡) → 『予測』では解けない"],
    size=12.5, first=False)
box(s, Inches(9.1), Inches(1.15), Inches(3.75), Inches(2.1),
    "AUCとは", STEEL, size=12,
    sub="判別力の点数。1.0=完璧、\n0.5=コイン投げと同じ。\n0.54は『ほぼ偶然』のレベル。", sub_size=10.5)
box(s, Inches(1.0), Inches(3.85), Inches(4.8), Inches(1.7),
    "予測アプローチ（失敗）", GREY,
    sub="使用挙動 → 凍結を機械学習で予測\nAUC≈0.54(公平領域) / 0.56(規範)")
line_seg(s, Inches(6.0), Inches(4.7), Inches(7.3), Inches(4.7), width=3,
         arrow_head=True)
box(s, Inches(7.4), Inches(3.85), Inches(5.3), Inches(1.7),
    "機械的監査アプローチ（提案）", GREEN,
    sub="『学習機会があったのに応答したか?』を\n機会END起点で因果的に監査(決定論カウンタ)")
tf = _tb(s, Inches(0.6), Inches(5.75), SW - Inches(1.2), Inches(0.6))
_set(tf.paragraphs[0],
     "発想転換: 『凍結を当てる』のではなく『再学習する“機会”に対して実際に応答したかを機械的に監査する』。本技術の核。",
     size=13, bold=True, color=GREEN)
takeaway(s, "壊れそうかをAIに予言させるのは失敗(コイン投げ並み)。『チャンスに応えたか』を出席簿のように数える方式に切り替えたら解けた。")
footer(s)

# =========================================================================== #
# 11. 提案手法の全体像
# =========================================================================== #
s = section("提案手法の全体像 — テレメトリから判定までの流れ", NAVY, "§3 Part1 ─ 本題")
tf = _tb(s, Inches(0.4), Inches(1.02), Inches(12.5), Inches(1.15))
_set(tf.paragraphs[0],
     "入力: テレメトリ(RSOC / FCC / cycleCount / timestamp) のみ。ハードウェア識別子(機種名・ベンダー・FRU)は判定に使わない(機種非依存)。",
     size=12.5, bold=True, color=NAVY)
q = tf.add_paragraph()
_set(q, "出力: 各ユーザを NORMAL / ゲージ再較正候補 / FW確認候補 / 監視 / 判定保留(データ品質) に振り分け。①〜⑤の各段が後述の特許候補に対応。",
     size=12)
steps = [
    ("① 学習機会の抽出", "RSOC high→low→high\n(深放電→再充電)"),
    ("② END起点で応答監査", "再充電完了後72h窓の\n有効ステップ≥50mWh"),
    ("③ 欠測/打ち切り耐性", "段階品質ティア+\ncensored除外"),
    ("④ 二分岐トリアージ", "機会反復×無応答→FW\n機会皆無→ゲージ再較正"),
    ("⑤ 有界保持で証拠保全", "最小状態の因果台帳\n(30日保持で全期間等価)"),
]
bw = Inches(2.36); gap = Inches(0.12); x0 = Inches(0.35); y = Inches(2.65)
for i, (t, sub) in enumerate(steps):
    x = x0 + (bw + gap) * i
    box(s, x, y, bw, Inches(1.7), t,
        [BLUE, TEAL, ORANGE, PURPLE, NAVY][i], sub=sub, size=12)
    if i < 4:
        line_seg(s, x + bw, y + Inches(0.85), x + bw + gap, y + Inches(0.85),
                 width=2.5, arrow_head=True)
tf = _tb(s, Inches(0.4), Inches(4.65), Inches(12.5), Inches(1.6))
bullets(tf, [
    "①+③+④ → 特許候補①（機会条件付き無応答監査・中核判定法）",
    "②(微小変化の扱い) → 特許候補②（デュアルトラック非対称リセット状態機械）",
    "⑤(保持制約下の証拠保全) → 特許候補③（有界保持の因果証拠台帳・最小状態）"],
    size=12.5, color=GREY)
takeaway(s, "5段のパイプライン=1つの判定技術。その中の『核になる工夫』3つが、そのまま特許候補①②③になっている。")
footer(s)

# =========================================================================== #
# 12. 仕組み図解A — 学習機会と END 起点 72h 監査
# =========================================================================== #
s = section("仕組み図解A — 学習機会と『END起点72h監査』", TEAL, "§3 Part1 ─ 図解")
# panel axes
ax0, ay0 = Inches(0.7), Inches(1.45)   # panel top-left
ax1, ay1 = Inches(7.35), Inches(5.75)  # panel bottom-right
line_seg(s, ax0, ay1, ax1, ay1, color=GREY, width=1.5, arrow_head=True)
line_seg(s, ax0, ay1, ax0, ay0, color=GREY, width=1.5, arrow_head=True)
label(s, ax1 - Inches(0.75), ay1 + Inches(0.03), Inches(1.2), Inches(0.3), "時間 →", size=10, color=GREY)
label(s, ax0 - Inches(0.3), ay0 - Inches(0.32), Inches(1.6), Inches(0.3), "RSOC (%)", size=10, color=GREY)
# threshold dashed lines: 80% at y=2.35, 20% at y=4.95
y80, y20 = Inches(2.35), Inches(4.95)
line_seg(s, ax0, y80, ax1 - Inches(0.15), y80, color=GREY, width=1, dash=MSO_LINE.DASH)
line_seg(s, ax0, y20, ax1 - Inches(0.15), y20, color=GREY, width=1, dash=MSO_LINE.DASH)
label(s, ax0 + Inches(0.02), y80 - Inches(0.28), Inches(1.0), Inches(0.25), "80%", size=9.5, color=GREY)
label(s, ax0 + Inches(0.02), y20 - Inches(0.28), Inches(1.0), Inches(0.25), "20%", size=9.5, color=GREY)
# audit window rect FIRST (behind curve): END at x=5.15, window to 6.75
wx0, wx1 = Inches(5.15), Inches(6.9)
wrect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, EM(wx0), EM(ay0 + Inches(0.25)),
                           EM(wx1 - wx0), EM(ay1 - ay0 - Inches(0.4)))
wrect.fill.solid(); wrect.fill.fore_color.rgb = LGREEN_BG
wrect.line.color.rgb = GREEN; wrect.line.dash_style = MSO_LINE.DASH
wrect.shadow.inherit = False
label(s, wx0 + Inches(0.06), ay0 + Inches(0.28), Inches(1.8), Inches(0.55),
      "監査窓\n[END, END+72h]", size=9.5, bold=True, color=DGREEN)
# RSOC curve
pts = [(0.35, 0.75), (1.55, 0.75), (2.95, 3.85), (3.35, 3.85), (4.45, 0.75), (6.3, 0.75)]
for i in range(len(pts) - 1):
    line_seg(s, ax0 + Inches(pts[i][0]), ay0 + Inches(pts[i][1]),
             ax0 + Inches(pts[i + 1][0]), ay0 + Inches(pts[i + 1][1]),
             color=BLUE, width=2.75)
# markers START / LOW / END
for xm, lab, col in [(1.86, "START", NAVY), (3.15, "LOW(谷)", NAVY), (4.45, "END=機会完了", DGREEN)]:
    line_seg(s, ax0 + Inches(xm), ay0 + Inches(0.35), ax0 + Inches(xm), ay1 - Inches(0.1),
             color=col, width=1.25, dash=MSO_LINE.ROUND_DOT)
    label(s, ax0 + Inches(xm - 0.5), ay1 - Inches(0.05), Inches(1.7), Inches(0.3),
          lab, size=9.5, bold=True, color=col)
# right outcome boxes
box(s, Inches(7.7), Inches(1.35), Inches(5.25), Inches(1.15),
    "responded（応答）", GREEN,
    sub="監査窓内に有効ステップ(≥50mWh)を観測 → ゲージは学習している(健全寄り)", sub_size=10)
box(s, Inches(7.7), Inches(2.7), Inches(5.25), Inches(1.15),
    "no_response（無応答）", RED,
    sub="窓[END, END+72h]を最後まで観測したのに有効ステップ無し → 無応答の証拠として計数", sub_size=10)
box(s, Inches(7.7), Inches(4.05), Inches(5.25), Inches(1.3),
    "censored（打ち切り=保留）", GREY,
    sub="窓の観測が途中で終わった(データがそこで尽きた等) → 無応答に数えない。\n期限は『窓を全部観測できたとき』だけ発火(未来情報も使わない)", sub_size=10)
label(s, Inches(7.7), Inches(5.5), Inches(5.25), Inches(0.8),
      "起点をSTARTにすると『放電途中のFCC変化』まで応答に誤計上される(56%汚染)。"
      "ENDなら構造的に0%(根拠②)。", size=10.5, color=INK)
takeaway(s, "採点は試験が終わって(END)から72時間だけ。途中退室(censored)は不合格に数えない — この2つの規律が誤判定を防ぐ。")
footer(s)

# =========================================================================== #
# 13. 実フリートでの判定結果
# =========================================================================== #
s = section("実フリート752台での判定結果 — 全履歴監査と30日オンライン運用", NAVY, "§3 Part1 ─ 結果")
label(s, Inches(0.5), Inches(0.98), Inches(6.0), Inches(0.35),
      "A. 全履歴監査（回顧的・v2.0-final、proxyラベルの供給源）", size=12.5, bold=True, color=NAVY)
make_table(
    s, ["最終ラベル", "件数", "割合"],
    [["NORMAL_OR_RESPONDING（正常/応答中）", "327", "43.5%"],
     ["REVIEW_INSUFFICIENT_DATA（データ不足で保留）", "338", "44.9%"],
     ["WATCH（低更新率・判定曖昧で経過観察）", "55", "7.3%"],
     ["ACTIONABLE_GAUGE_RESET（ゲージ再較正候補）", "18", "2.4%"],
     ["ACTIONABLE_FW_CHECK（FW確認候補）", "14", "1.9%"]],
    Inches(0.5), Inches(1.38), [4.35, 0.8, 0.85], row_h=0.42, cell_size=10,
    styler=lambda i, j, v, b: (b, j > 0, INK))
label(s, Inches(0.5), Inches(4.05), Inches(6.1), Inches(1.0),
      "候補抽出96台→二分岐。既存active(更新中)層からの誤エスカレーション0件。"
      "FW 14台・ゲージ18台は全て stale/very_stale 層由来。", size=10.5, color=INK)
label(s, Inches(6.95), Inches(0.98), Inches(6.0), Inches(0.35),
      "B. 30日オンライン運用版 v2（毎日スライド・9段階）", size=12.5, bold=True, color=NAVY)
make_table(
    s, ["運用ラベル(v2)", "件数"],
    [["REVIEW_DATA_QUALITY（データ品質で保留）", "325"],
     ["NORMAL_RESPONDING（正常応答）", "183"],
     ["WATCH_LARGE_GAP_OR_CENSORED（大ギャップ/打ち切り監視）", "128"],
     ["FW_WATCH_HIGH_ANOMALY（FW監視）", "43"],
     ["WATCH_LOW_EVIDENCE（低証拠監視）", "35"],
     ["GAUGE_SOFT_CALIBRATION（軽い較正促し）", "22"],
     ["GAUGE_REVIEW（ゲージ・要人手確認）", "7"],
     ["FW_CHECK_CORE（FW確認・最有力）", "5"],
     ["GAUGE_RESET_CORE（ゲージ再較正・最有力）", "4"]],
    Inches(6.95), Inches(1.38), [5.15, 0.75], row_h=0.4, cell_size=9.5,
    styler=lambda i, j, v, b: (
        b, i >= 8 or j == 1,
        RED if "FW" in v else ORANGE if "GAUGE" in v else INK if j == 0 else NAVY))
label(s, Inches(0.5), Inches(4.85), Inches(6.1), Inches(1.5),
      "オンライン版の対proxy性能: FW/ゲージCore判定 precision 1.0（proxy基準）、"
      "能動誤報0件（3種の活動基準すべて）。FWエンジニアリング上位50件は proxy FW 14台を全カバー(recall 1.0)。"
      "状態の持ち越しにより、30日窓単体比で無応答検出 22→41台(+29)。", size=10.5, color=INK)
takeaway(s, "精密検査(全履歴)と日常見守り(30日運用)の二本立て。約半数は『データ不足で保留』と正直に言う設計 — 誤指示より保留が安全。")
footer(s)

# =========================================================================== #
# 14. 根拠マップ
# =========================================================================== #
s = section("判定法が機能する根拠 — 5本柱の独立検証", DGREEN, "§4 根拠")
make_table(
    s, ["根拠", "検証", "主張", "ひとことで（非エンジニア向け）"],
    [["① A2", "負の対照（ニセ機会5種）", "応答は真の機会に特異(真0.39が5/5対照のヌル外)", "プラセボでは効かず、本物でだけ効く"],
     ["② A3", "応答起点の比較", "END起点の因果汚染0 vs START 0.56", "採点はテスト後から＝カンニング防止"],
     ["③ B", "応答ハザード（生存解析）", "真の機会後ほど速く応答(72h 0.39 vs 0.29)", "本物の機会の後ほど早く計り直しが起きる"],
     ["④ E", "欠測・打ち切り注入", "誤『無応答』を naive 643→提案 4.1 に抑制", "データの穴を冤罪にしない"],
     ["⑤ D", "保持グリッド 7〜90日", "30日保持で全期間と一致(recall 1/重複0)・容量4%", "30日分のメモで全履歴と同じ結論"]],
    Inches(0.4), Inches(1.1), [1.0, 2.6, 4.7, 4.2], row_h=0.6, cell_size=10.5,
    styler=lambda i, j, v, b: (b, j == 0, DGREEN if j == 0 else INK))
tf = _tb(s, Inches(0.45), Inches(4.6), SW - Inches(0.9), Inches(1.6))
bullets(tf, [
    "いずれも proxy ラベルに依存しない独立エンドポイントで評価（9項目すべて supported、patent_technical_effects_v4.csv）",
    "信頼区間はユーザー単位bootstrap（同一ユーザの複数機会を独立扱いしない）",
    "ベースライン再現ゲート16/16一致・成果物40件にSHA-256・決定的シード（再現性）"], size=11.5)
takeaway(s, "次の5枚で1本ずつ図解する。各図は【何のグラフ/軸/主張/読み方】の4ブロック構成。")
footer(s)

# =========================================================================== #
# 15-19. 根拠図
# =========================================================================== #
RES = "§4 根拠"
figure("根拠① 応答は『真の機会』に特異（負の対照検定 A2）", RES,
       "negative_control_true_vs_null.png",
       "『機会ENDの後に応答が増える』が偶然でないかを、END時刻を人工的にズラした5種類のニセ機会(ネガティブコントロール)と比較した検定。",
       "縦軸=END後72h以内に有効ステップが起きた確率。青棒=各ニセ機会の応答率(黒ヒゲ=95%区間)。赤の水平線=実データ(約0.39)。",
       f"実データ {A2['true_resp_prob_72h']:.2f} は、どのニセ機会(0.25〜0.31)よりも95%区間の上端を超えて高い"
       f"（{A2['n_controls_outside_null']}/{A2['n_controls_total']} 対照でヌル外、bootstrap方向一致 {A2['n_controls_directionally_supported']}/5）。"
       "応答は『真の機会END時刻』に特異で、時間経過・活動量・ユーザ差では説明できない。",
       "赤線が全ての青棒(ヒゲ含む)より上にあることを確認する。棒の高さ・真値の数値は図に印字されていないため軸目盛りから読む。"
       f"対象: {A2['n_anchors']}機会 / {A2['n_users_with_anchors']}ユーザー。",
       "薬の治験のプラセボと同じ発想。ニセの機会5種類では応答が増えず、本物の機会でだけ増えた(5戦5勝)。判定の土台が成立。")
figure("根拠② 応答の起点はEND(再充電完了)が正しい（A3）", RES,
       "response_anchor_contamination.png",
       "応答を測る『起点』を START/LOW/END のどこに置くべきかの比較。誤った起点だと『再充電が完了する前のFCC変化』まで応答に誤計上する(=因果汚染)。",
       "横軸=起点3種(start/low/end)。赤棒=因果汚染率(完了前の変化を誤計上した割合)、紫棒=重複帰属率(同じ変化を複数機会に数えた割合)。",
       f"START起点は {A3['start_contamination_frac_72h']:.2f}(56%)が汚染、LOWは {A3['low_contamination_frac_72h']:.2f}、"
       "END起点は汚染0（endの赤棒が無い=ゼロの意味）。※END=0は定義から従う構造的性質(設計の正しさの確認)。",
       "赤棒が低いほど良い。紫(重複)は起点によらず約22%で残るため、これは別の仕組み(episode_id+seen_ids、候補③)で防ぐ。"
       f"production(any-change)とのEND一致率 {A3['end_agreement_with_production']:.2f}。",
       "採点は試験終了後に始める、というルール。STARTから数えると半分以上が『試験中の答案』を誤って採点していることになる。")
figure("根拠③ 真の機会の後ほど速く応答する（応答ハザード B）", RES,
       "response_hazard_true_vs_pseudo.png",
       "機会ENDからの経過時間に対する累積応答率(生存解析のCIF)。真の機会 vs 観測条件を揃えたニセ時点(matched pseudo)の比較。",
       "横軸=END後の経過時間(0〜168h)。縦軸=その時刻までに有効応答(≥50mWh)した累積割合。青=真の機会、橙=pseudo。",
       f"真の曲線が常に上: 72hで {B['true_cif_72h_50mwh']:.2f} vs {B['pseudo_cif_72h_50mwh']:.2f}(差 +{B['true_minus_pseudo_72h']:.2f})。"
       f"有効応答の中央値 {B['median_response_h_50mwh']:.0f}h → 72h窓は『待ちすぎず・切り捨てすぎない』妥当な締め切り。",
       "立ち上がりの傾きに注目: 差は最初の約30時間で開き、その後約15ポイント差がほぼ一定に保たれる。"
       "応答が無いまま時間だけが経つ個体=更新停止候補、という読み方の裏付け。",
       "本物の機会の直後は、無関係な時点よりも明らかに早く・多く『計り直し』が起きる。だから機会の直後だけを見張ればよい。")
figure("根拠④ 欠測・打ち切りでも誤判定しない（欠測ストレス E）", RES,
       "missingness_false_escalation.png",
       f"欠測・連続ギャップ・末尾打ち切り・睡眠ギャップなど{E_['n_regimes']}種類の欠測パターンを人工注入し、"
       "各検出器が『誤って無応答と確定』した件数を比較したストレステスト。",
       "縦軸=注入パターン(17行: MCAR 5〜50% / ギャップ3〜48h×位置 / 打ち切り等)。横軸=誤確定無応答の平均件数。"
       "緑=naive(素朴)、橙=graded(段階品質のみ)、青=binary gate、赤=proposed(提案)。",
       f"naiveは全パターンで550〜670件を誤確定(平均 {E_['naive_mean_false_no_response']:.0f})。"
       f"提案(段階品質+censor考慮)は平均 {E_['proposed_mean_false_no_response']:.1f} 件(約99%減)、"
       f"かつ機会の回収率 {E_['proposed_episode_recovery']:.2f} を維持(除外しすぎない)。",
       "棒が短いほど良い。提案(赤)の棒はほぼ見えない=誤判定ほぼゼロの意味。binary gateは誤確定こそ少ないが"
       "機会を除外しすぎて見逃しが増える(付録の補足図参照)。凡例ボックスが上2行の緑棒の右端に重なる点に注意。",
       "『データの穴』を無視も冤罪もせず、判定保留(censored)として脇に置く。だから穴だらけの実運用データでも騙されない。")
figure("根拠⑤ 30日保持でも全期間と同じ判定（保持不変 D）", RES,
       "retention_invariance_heatmap.png",
       "生データの保持期間を7〜90日に制限したとき、全期間データを使った場合と同じ応答判定になるか(一致率)のマップ。",
       "行=方式(上=stateful:台帳あり / 下=stateless:台帳なし)。列=保持日数(7/14/21/30/45/60/90)。"
       "セルの数値=全履歴版との応答ステータス一致率(1.00=完全一致)。緑が濃いほど高い。",
       f"statefulは全列 1.00(保持7日でも完全一致)。statelessは7日保持で 0.62 まで劣化。"
       f"別途の同一エンジン検証で stateful は recall {D['stateful_verify_recall']:.0f}・重複 {D['stateful_verify_duplicates']}・"
       f"無応答数の誤差MAE≈{D['stateful_verify_no_response_mae']}、必要ストレージは全生データ比 {D['min_stateful_equivalent_storage_ratio']*100:.1f}%。",
       f"上段が全部緑=『台帳(最小状態)を持てば生データを捨てても結論が変わらない』。下段との差が台帳の価値。"
       f"グリッドは全{D['n_grid_configs']}構成(保持×応答窓×stride×整列)で検証。",
       "レシート(生データ)は30日で捨てても、家計簿(台帳)に転記してあれば集計は狂わない。家計簿の厚さは元の約4%。")

# =========================================================================== #
# 20. Part2 位置づけ + 集約表
# =========================================================================== #
s = section("過程で生まれた特許候補アイデア（位置づけ）", PURPLE, "§5 Part2")
tf = _tb(s, Inches(0.6), Inches(1.05), Inches(12.1), Inches(1.7))
_set(tf.paragraphs[0],
     "本技術(FCC更新停止者の判定)は単一目的の一連の手法。その実現過程で、単独でも特許になりうる中核アイデアが生まれた。",
     size=13.5, bold=True, color=NAVY)
q = tf.add_paragraph()
_set(q, "細分化しすぎると評価しづらいため、ここでは3つ(+将来2つ)に集約して提示する。どこをどの粒度で特許化するかの"
        "最終判断は社内reviewerに委ねる。各候補に技術エビデンス強度と新規性リスクを併記し判断材料を提供。",
     size=12, color=GREY)
make_table(
    s, ["特許候補（集約後）", "対応する手法段", "技術エビデンス", "新規性リスク(UNVERIFIED)"],
    [["① 機会条件付き無応答監査（中核判定法）", "①+③+④", "STRONG", "MEDIUM-HIGH"],
     ["② デュアルトラック非対称リセット状態機械", "②(微小変化)", "STRONG", "HIGH（着想日依存）"],
     ["③ 有界保持の因果証拠台帳（最小状態）", "⑤(保持制約)", "STRONG", "MEDIUM-HIGH"],
     ["(将来④) 診断依存クローズドループ介入検証", "介入後検証", "PROSPECTIVE", "—"],
     ["(将来⑤) 機種非依存スクリーニング+version局在", "原因局在", "MEDIUM/PROSPECTIVE", "MEDIUM"]],
    Inches(0.5), Inches(2.9), [5.4, 2.0, 2.4, 2.5], row_h=0.52, cell_size=10.5,
    styler=status_styler({2, 3}))
takeaway(s, "『技術エビデンスが強い』と『特許が取れる』は別物。強い証拠(左)と、先行技術に潰されるリスク(右)を分けて見るのがこの表の読み方。")
footer(s)

# =========================================================================== #
# 21. 候補① 概要
# =========================================================================== #
concept("特許候補① 機会条件付き無応答監査（中核判定法）",
        "§5 候補①（IC1+IC6相当）", "STRONG", "MEDIUM-HIGH",
        "FCC再学習は『深い充放電(学習機会)』でしか起きない。だが機会の有無と応答の有無は静的検査では分離できない。",
        "FCC凍結が『機会が無いだけ(正常/再較正)』なのか『機会はあるのに応答しない(FW/HW疑い)』なのか判別できない。",
        "テレメトリだけで、機会の有無と応答の有無を因果的に分離し、欠測・打ち切りに騙されずに判定する。",
        "RSOC high→low→high 機会を抽出 → END起点72h窓で有効ステップを responded/no_response/censored に分類 → "
        "段階品質ティアとcensor除外で欠測を無応答に数えない → 機会反復×無応答=FW候補 / 機会皆無=ゲージ再較正候補に二分岐。",
        "A2で機会-応答の特異性、A3でEND起点の因果汚染0、Eで欠測耐性を実証(根拠①②④)。"
        "【リスク】『非発生イベント監視』は広い先行技術(UNVERIFIED)→ 狭め/中位スコープ(80/20/80帯・72h窓・50mWh・censor除外を明示)での出願を推奨。",
        "『テスト(機会)を受けたのに白紙(無応答)』の個体だけを、欠測に騙されずに自動抽出する仕組み。判定技術の中核。")

# 22. 候補① 特徴量
s = features("特許候補① 必要な特徴量（と作成方法）", "§5 候補①", [
    ("RSOC / FCC / cycleCount / timestamp", RAW, "テレメトリ生データ（PWMログ、30分間隔）。識別子は使わない"),
    ("学習機会 episode (start/low/end)", ENG, "RSOC high→low→high 状態機械(80/20/80等): extract_high_low_high_episodes"),
    ("有効FCCステップ(≥50mWh)", ENG, "時刻順FCC差分の絶対値≥50mWh: fcc_step_indicator(FCC,50)"),
    ("応答ステータス response_status_72h", ENG, "END起点[end,end+72h]の有効step有無+窓完全性で responded/no_response/censored"),
    ("段階品質ティア quality_tier", ENG, "エピソード最大ギャップ→HIGH_OK/MEDIUM_GAP/LOW_LARGE_GAP: online_gap_quality"),
    ("二分岐スコア(FW/ゲージ)", ENG, "機会反復×無応答 vs 機会皆無 のカウンタから振り分け"),
    ("(学習除外) 機種名/ベンダー/FRU", "—", "ハードウェア識別子は判定に一切使わない(機種非依存)"),
])
takeaway(s, "生データは4列だけ。あとは全て派生量 — 『何を数えるか』の設計そのものが発明。")

# 23. 候補② 概要
concept("特許候補② デュアルトラック非対称リセット状態機械",
        "§5 候補②（IC2+C3相当）", "STRONG", "HIGH",
        f"燃料計は整数mWh量子化(最小{C3['quantization_unit_mwh']:.0f}mWh)。微小変化(micro)と意味ある再学習(effective≥50mWh)が"
        f"混在する(二峰分布: モード {C3['gmm_micro_mode_mwh']:.0f} / {C3['gmm_effective_mode_mwh']:.0f}mWh、micro率58.1%)。",
        "状態を1系統で持つと、microステップが『まだ未解決の無応答証拠』をリセットして消してしまい、FW疑いを見逃す。",
        "微小ゆらぎで証拠を失わず、かつ『ソフト較正で済む個体』と『ハードリセットが要る個体』を区別する(過剰なハード指示を避ける)。",
        "any系統(任意変化でリセット)とeffective系統(≥50mWhのみリセット)を分離。microはany系統だけリセットし、"
        "effective状態・未解決証拠は保持(非対称)。同時刻イベントは complete<reset<deadline の決定的順序で処理。",
        f"C2アブレーション: 対称リセットは保留{C2['d2_pending_erased']}件+確定無応答{C2['d2_no_response_erased']}件を"
        f"{C2['d2_users_evidence_erased']}ユーザーで消去、非対称は消去0。ハード指示 {C2['hard_prompts_d1_effective_only']}→"
        f"{C2['hard_prompts_d4_proposed']}件({C2['d4_gauge_soft']}件をソフト較正へ)。"
        "【リスク高】本設計は既にproduction実装済み(v4は特徴付け/検証であり着想ではない)+deadband系は一般的先行技術"
        "→ 新規性は着想日依存(法的論点)。クレームは非対称リセット規則を明示。",
        "小銭のゆらぎ(micro)で捜査ファイル(証拠)を閉じない、二重帳簿+非対称リセット。ただし着想日の確認が必須。")

# 24. 仕組み図解B — dual-track
s = section("仕組み図解B — dual-track 非対称リセットの動き", TEAL, "§5 候補② ─ 図解")
# timeline
tl_y = Inches(2.05)
line_seg(s, Inches(0.6), tl_y, Inches(7.3), tl_y, color=GREY, width=1.5, arrow_head=True)
label(s, Inches(6.55), tl_y + Inches(0.06), Inches(1.0), Inches(0.3), "時間 →", size=10, color=GREY)
events = [(1.7, "micro\n+20mWh", ORANGE, 0.28), (3.6, "有効ステップ\n+300mWh", GREEN, 0.5),
          (5.5, "micro\n−30mWh", ORANGE, 0.28)]
for ex, lab, col, hh in events:
    line_seg(s, Inches(ex), tl_y - Inches(hh), Inches(ex), tl_y, color=col, width=3)
    label(s, Inches(ex - 0.55), tl_y - Inches(hh) - Inches(0.48), Inches(1.5), Inches(0.45),
          lab, size=9, bold=True, color=col, align=2)
# any lane
lane1 = box(s, Inches(0.6), Inches(2.6), Inches(6.7), Inches(1.15),
            "any系統 — どんな変化でもリセット", LBLUE_BG, fg=NAVY, size=11.5,
            sub="『最後に何か動いてから何日/何サイクル』を追跡（微小ゆらぎ検知・軽い較正促し用）", sub_size=9.5, line=BLUE)
for ex, _, col, _h in events:
    label(s, Inches(ex - 0.25), Inches(3.37), Inches(0.8), Inches(0.3), "↻リセット", size=9, bold=True, color=BLUE)
# effective lane
box(s, Inches(0.6), Inches(3.95), Inches(6.7), Inches(1.3),
    "effective系統 — ≥50mWhのみリセット", LGREEN_BG, fg=DGREEN, size=11.5,
    sub="未解決の無応答証拠・機会カウンタ・pending期限を保持（FW疑い判定用）", sub_size=9.5, line=GREEN)
label(s, Inches(1.45 - 0.3), Inches(4.85), Inches(0.9), Inches(0.3), "保持", size=9.5, bold=True, color=DGREEN)
label(s, Inches(3.6 - 0.3), Inches(4.85), Inches(0.9), Inches(0.3), "↻リセット", size=9.5, bold=True, color=GREEN)
label(s, Inches(5.5 - 0.35), Inches(4.85), Inches(0.9), Inches(0.3), "保持", size=9.5, bold=True, color=DGREEN)
label(s, Inches(0.65), Inches(5.42), Inches(6.6), Inches(0.7),
      "microイベント(橙)は any系統だけをリセットし、effective系統の証拠は保持 = 『非対称』。", size=10.5, bold=True, color=INK)
# right: symmetric vs asymmetric comparison
box(s, Inches(7.7), Inches(1.35), Inches(5.25), Inches(1.85),
    "もし対称リセットだったら", GREY,
    sub=f"microのたびに証拠が消える:\n保留中の機会 {C2['d2_pending_erased']}件 + 確定無応答 {C2['d2_no_response_erased']}件が\n"
        f"{C2['d2_users_evidence_erased']}ユーザーで消失（C2アブレーション実測）", sub_size=10)
box(s, Inches(7.7), Inches(3.4), Inches(5.25), Inches(1.85),
    "非対称リセット（提案・production実装）", GREEN,
    sub=f"消去 0件。有効のみ系統との比較でも確定無応答 +{C2['evidence_preserved_vs_symmetric']}件を温存。\n"
        f"ハード指示 {C2['hard_prompts_d1_effective_only']}→{C2['hard_prompts_d4_proposed']}件"
        f"（{C2['d4_gauge_soft']}件は『軽い較正促し』に振替）", sub_size=10)
label(s, Inches(7.7), Inches(5.45), Inches(5.25), Inches(0.85),
      "2つの帳簿を分けることで『完全凍結(両方止まる)』と『微小ゆらぎのみ(anyだけ動く)』を区別し、"
      "対処の強さを変えられる。", size=10.5, color=INK)
takeaway(s, "帳簿を2冊持ち、決定的証拠(≥50mWh)のときだけ両方をリセット。メモ書き(micro)では捜査ファイルを閉じない。")
footer(s)

# 25. 候補② 根拠図 C2
figure("候補②の根拠: 微小変化で証拠が消える vs 守る（C2）", "§5 候補②",
       "dual_track_erased_evidence.png",
       "同一イベント列を6つのリセット方式(D0〜D5)でリプレイし、microステップによって『消されてしまう』証拠の件数を比較した直接アブレーション。",
       "横軸=リセット方式(D0 any単独/D1 effective単独/D2・D3 対称/D4 非対称=提案/D5 適応型)。"
       "縦軸=消去件数の積み上げ(橙=確定無応答の証拠、茶=保留中の機会)。",
       f"対称リセット(D0/D2)は 無応答{C2['d2_no_response_erased']}件+保留{C2['d2_pending_erased']}件 ≈ 計2,264件を"
       f"{C2['d2_users_evidence_erased']}ユーザーで消去。非対称(D4=提案)とD1/D3/D5は消去0(棒が無い=ゼロの意味)。",
       "棒が立っている方式=証拠を失う方式。数値は棒に印字されていないため軸から読む。"
       "『非対称リセットが証拠を守る』ことの最も直接的な実測証拠。",
       "『何でも消しゴム』方式だと281人分の捜査記録が消える。非対称ならゼロ — 図の空白(棒なし)がその証拠。")

# 26. 候補② 根拠図 C3
figure("候補②の根拠: 有効閾値50mWhはデータの谷が裏付け（C3）", "§5 候補②",
       "effective_threshold_mixture_fit.png",
       "FCCステップ量(絶対値・対数軸)の分布と閾値候補の位置。微小ゆらぎと意味ある再学習を分ける境界の妥当性を確認する図。",
       "横軸=log10|FCCステップ|(mWh)。1.0=10mWh、2.0=100mWh、3.0=1,000mWh。縦軸=密度。"
       "青破線=閾値候補(量子化10 / 谷11 / p50=30 / GMM谷35 / p75=200mWh)、黒実線=採用値50mWh。",
       f"分布は二峰: 左に量子化刻み由来の鋭い山(micro、モード≈{C3['gmm_micro_mode_mwh']:.0f}mWh)、"
       f"右に実質更新の山(effective、モード≈{C3['gmm_effective_mode_mwh']:.0f}mWh)。GMM谷={C3['gmm_valley_mwh']:.1f}mWh"
       "(bootstrap CI [26.3, 54.1])で、50mWhは谷の上・CI内。micro率58.1%。",
       "左の山=ほぼ意味のない微小変化、右の山=本物の容量変化、黒線がその谷間。"
       "※50mWhは『データから発見した最適値』ではなく事前指定の設計値で、データは『矛盾しない』ことの独立裏付け"
       "(青破線5本は同色で区別しづらい点に注意)。広い発明概念は適応閾値(付録C)。",
       "さざ波と本物の段差は、大きさの分布がきれいに2つの山に分かれる。境界線はその谷間に引いてある。")

# 27. 候補② 特徴量
s = features("特許候補② 必要な特徴量（と作成方法）", "§5 候補②", [
    ("FCCステップ(符号付き/絶対値)", ENG, "時刻順FCC差分の非ゼロ値: patent_dual_track.fcc_steps"),
    ("is_effective / is_micro", ENG, "絶対値 ≥ / < 50mWh(有効閾値): step_threshold_mwh"),
    ("量子化単位 quantization_unit", ENG, "観測された最小の正の|step|(=10mWh)"),
    ("any/effective 状態・未解決証拠カウンタ", ENG, "イベント列の状態機械リプレイ(micro/effで非対称リセット): online_step_state"),
    ("混合分布 micro/effective モード・谷", ENG, "log10(|step|)に2成分Gaussian Mixtureを適合: patent_effective_threshold"),
    ("持続/反転(persistence/reversal)", ENG, "各stepの次step時刻差・符号反転(24h反転率: micro 0.14 < eff 0.36)"),
    ("適応閾値(代替実施形態)", ENG, "max(k·quantization, α·DesignCapacity, noise percentile)"),
])
takeaway(s, "鍵は『しきい値の数字』ではなく『2系統を非対称に扱う規則』。閾値自体は代替実施形態で広くカバーする。")

# 28. 候補③ 概要
concept("特許候補③ 有界保持下の因果証拠台帳（最小十分状態）",
        "§5 候補③（IC5相当）", "STRONG", "MEDIUM-HIGH",
        "フリート監視では生テレメトリの保持期間が有界(例: 直近30日のみ)。それより古い生データは消える。",
        f"保持窓をまたぐ機会の証拠が、生データ破棄で失われる/二重に数えられる(台帳なしだと7日保持で一致率0.62、"
        f"recall {D['stateless_7d_recall_72h']:.2f}・重複率 {D['stateless_7d_dup_rate_72h']:.1f})。",
        "有界保持でも『全期間版と等価』な判定を、できるだけ小さい永続状態で実現する。",
        "部分FSM・pending期限キュー・確定済みID集合(seen_ids)・直近any/effective変化・gap/censorカウンタ等を永続化し、"
        "解決済みイベントを時刻順(同時刻は complete<reset<deadline)にリプレイ。期限は窓を全部観測できたときのみ発火(未来リーク無し)。",
        f"Dグリッド({D['n_grid_configs']}構成)で recall=1.0・重複0・誤差MAE≈{D['stateful_verify_no_response_mae']} を"
        f"ストレージ比 {D['min_stateful_equivalent_storage_ratio']*100:.1f}% で達成(根拠⑤)。"
        f"最小状態アブレーションで全{DM['n_necessary']}要素が必要と実証。"
        "【リスク】streaming+cachingは既知の組合せ(自明性)→ アブレーションで必要性を示した『最小状態構造』をクレーム核に。",
        "レシートは30日で捨てても、7項目の家計簿(台帳)があれば全履歴と同じ結論。『どの7項目か』を特定し必要性を実証したのが発明。")

# 29. 仕組み図解C — 有界保持台帳
s = section("仕組み図解C — 有界保持の因果証拠台帳", TEAL, "§5 候補③ ─ 図解")
# timeline with sliding window (window rect first so the arrow stays visible)
tl2 = Inches(1.9)
wrect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, EM(Inches(4.3)), EM(tl2 - Inches(0.55)),
                           EM(Inches(2.6)), EM(Inches(1.0)))
wrect.fill.solid(); wrect.fill.fore_color.rgb = LBLUE_BG
wrect.line.color.rgb = BLUE; wrect.shadow.inherit = False
line_seg(s, Inches(0.6), tl2, Inches(7.3), tl2, color=GREY, width=1.5, arrow_head=True)
label(s, Inches(6.95), tl2 + Inches(0.12), Inches(0.9), Inches(0.3), "時間 →", size=10, color=GREY)
label(s, Inches(4.4), tl2 - Inches(0.5), Inches(2.5), Inches(0.35),
      "生データ窓（直近30日のみ）", size=9.5, bold=True, color=BLUE)
label(s, Inches(0.75), tl2 - Inches(0.5), Inches(3.4), Inches(0.35),
      "これより古い生データは破棄される", size=9.5, color=GREY)
# episode straddling boundary
ep = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, EM(Inches(3.4)), EM(tl2 + Inches(0.25)),
                        EM(Inches(1.9)), EM(Inches(0.5)))
ep.fill.solid(); ep.fill.fore_color.rgb = LGREEN_BG; ep.line.color.rgb = GREEN
ep.shadow.inherit = False
tfep = ep.text_frame; tfep.word_wrap = True
_set(tfep.paragraphs[0], "窓をまたぐ学習機会", size=9.5, bold=True, color=DGREEN, align=2)
label(s, Inches(0.7), tl2 + Inches(0.95), Inches(6.5), Inches(0.8),
      "機会の開始は窓の外(破棄済み)、応答期限は窓の中 — 生データだけでは検出も判定もできない。"
      "台帳の部分FSM・pendingが これを回収する。", size=10.5, color=INK)
# ledger box
comps = [
    ("fsm", "機会検出の途中状態(高→低→高のどこまで進んだか)"),
    ("pending", "未解決機会の応答期限キュー"),
    ("seen_ids", "確定済み機会ID集合(二重計上防止)"),
    ("last_eff_ts / eff_cycle", "直近の有効変化の時刻・サイクル(経過日数/サイクルを駆動)"),
    ("gap_censor", "ギャップ・打ち切りのカウンタ(品質ゲート)"),
    ("ordering", "同時刻イベントの決定的順序 complete<reset<deadline"),
]
box(s, Inches(0.6), Inches(3.75), Inches(6.7), Inches(0.5),
    f"永続する最小状態（台帳）— 必要{DM['n_necessary']}要素(アブレーションで各要素の必要性を実証)",
    NAVY, size=11)
yy = Inches(4.35)
tf = _tb(s, Inches(0.7), yy, Inches(6.6), Inches(1.9))
first = True
for name, desc in comps:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    r = p.add_run(); r.text = f"{name}: "; _style(r, 10.5, True, NAVY)
    r2 = p.add_run(); r2.text = desc; _style(r2, 10.5, False, INK)
# right effect box
box(s, Inches(7.7), Inches(1.35), Inches(5.25), Inches(2.3),
    "効果（Dグリッドで実証）", DGREEN,
    sub=f"検出率 recall = 1.0（取りこぼしゼロ）\n重複計上 = 0\n無応答数の誤差 MAE ≈ {D['stateful_verify_no_response_mae']}\n"
        f"必要ストレージ = 生データの約 {D['min_stateful_equivalent_storage_ratio']*100:.1f}%",
    sub_size=11)
box(s, Inches(7.7), Inches(3.85), Inches(5.25), Inches(1.6),
    "台帳なし(stateless)だと", GREY,
    sub=f"7日保持: 一致率 0.62 / recall {D['stateless_7d_recall_72h']:.2f} / 重複率 {D['stateless_7d_dup_rate_72h']:.1f}\n"
        "保存量を増やしても完全一致(1.00)には届かない(補足図)", sub_size=10.5)
takeaway(s, "生データを捨てる制約は変えられない。だから『何をメモしておけば全履歴と等価か』を突き詰めた — その最小メモ(7要素)が発明の核。")
footer(s)

# 30. 候補③ 根拠図 D minimal state
figure("候補③の根拠: どの状態も削れない（最小性アブレーション）", "§5 候補③",
       "minimal_state_necessity.png",
       "台帳(永続状態)の構成要素を1つずつ外して同じ判定を再実行する破壊実験。どの要素が欠けても等価性が壊れることを確認する。",
       "横軸=外した要素(fsm / last_eff_ts / pending / seen_ids)。縦軸=不変量の破綻量。"
       "赤=recall低下(見逃しの増加)、橙=無応答判定の誤差(MAE)。",
       "fsm除去→recall約0.96を喪失(窓またぎ機会を取り逃す)+誤差0.85。last_eff_ts除去→誤差最大(≈1.8)。"
       f"pending除去→誤差≈1.3。アブレーション全体では全{DM['n_necessary']}要素が何らかの不変量を破る。",
       "棒が高い要素ほど外したときの壊れ方が大きい=必須。seen_idsはこの図の2指標では小さく見えるが、"
       "その必要性は『重複計上ゼロ』という別の不変量に現れる(重複はこの図に描かれていない)。"
       "この『これ以上削れない』ことの実証が、クレーム核(最小十分状態)の裏付け。",
       "家計簿のどのページを破っても集計が狂う=全ページ必要、を1ページずつ破って確かめた図。")

# 31. 候補③ 特徴量
s = features("特許候補③ 必要な特徴量（と作成方法）", "§5 候補③", [
    ("episode_id(正準ID)", ENG, "user|band|start|end(同一機会を一意識別、重複防止。公開時user部はhash)"),
    ("応答期限 / 直近変化 ts+cycle", ENG, "end+72h / 状態機械が逐次更新(経過日数・サイクルを駆動)"),
    ("確定済みID集合 seen_ids", ENG, "確定済み機会IDの集合(重複カウント防止)"),
    ("pending 期限キュー", ENG, "未解決機会の期限待ち(生データ破棄後も証拠保持)"),
    ("部分FSM 状態", ENG, "WAIT_HIGH/LOW/HIGH_AGAIN の途中状態(窓またぎ機会検出)"),
    ("gap/censor カウンタ", ENG, "品質ゲート用の欠測・打ち切り集計"),
    ("最小十分状態(全体)", ENG, f"上記の集合(全{DM['n_necessary']}要素)。アブレーションで必要性を実証: patent_state_minimality"),
])
takeaway(s, "台帳の中身は『結果の要約』ではなく『因果リプレイに必要な最小情報』。ここがstreaming+cachingの一般論との違い。")

# =========================================================================== #
# 32. 将来候補
# =========================================================================== #
s = section("将来の特許候補（PROSPECTIVE — データはまだ無い、捏造しない）", GREY, "§5 将来候補")
tf = _tb(s, Inches(0.55), Inches(1.02), Inches(12.3), Inches(2.6))
_set(tf.paragraphs[0], "将来④ 診断依存クローズドループ介入検証", size=14, bold=True, color=NAVY)
bullets(tf, [
    "内容: 診断ラベル別の介入(ゲージ→OEM承認の安全な較正のみ / FW→BIOS・EC・電池FW更新)後、『介入後最初の良質(HIGH_OK)な機会』"
    "から72h以内の有効FCC応答を主要評価項目として観測し、診断の正しさを閉ループで検証。介入は因果台帳に記録。",
    "準備済み: 介入データスキーマ11フィールド(ハッシュ化ID・介入種別・version前後・介入後初機会・72h応答) + 実験プロトコル + 検出力シミュレーション",
    "検出力の事前試算(ベースライン応答率0.39・α=0.05): 大きな効果(+0.35)なら各腕30台で検出力0.81 / 控えめな効果(+0.15)を狙うなら"
    "各腕80〜120台(0.46〜0.64) / 各腕120台×効果+0.25で0.976 ※実験結果ではなく必要台数の見積り",
    "状態: 介入記録・BIOS/EC/電池FWバージョンはデータに存在しない(NOT AVAILABLE、availability probeで確認済み)"],
    size=11.5, first=False)
tf = _tb(s, Inches(0.55), Inches(3.85), Inches(12.3), Inches(2.3))
_set(tf.paragraphs[0], "将来⑤ 機種非依存スクリーニング + version 局在", size=14, bold=True, color=NAVY)
bullets(tf, [
    "内容: 分類は行動特徴のみ(機種名/ベンダー/FRUを学習に使わない)。分類確定後に、記述統計としてハードウェア偏在を集計し、"
    "FWバージョン確認の優先順位付けに使う二段構成。",
    "実データの例(記述のみ・分類には不使用): ある電池FRUでFW確認候補が26台中6台(経験ベイズ縮約後19.7%、BH補正q=0.0001)に偏在"
    "— versionデータがあれば『どの版で起きるか』の局在まで進める",
    "状態: スクリーニングは実装済(エビデンスMEDIUM)。version局在はBIOS/EC/FW列が NOT AVAILABLE → PROSPECTIVE"],
    size=11.5, first=False)
takeaway(s, "『診断→対処→本当に直ったか』の輪を閉じるのが次の宿題。実験計画と必要台数の試算まで済ませてあり、あとはデータ収集の判断待ち。")
footer(s)

# =========================================================================== #
# 33. 正直な開示① 敵対的レビュー
# =========================================================================== #
s = section("正直な開示① 敵対的レビューと必須開示事項", RED, "§6 開示")
tf = _tb(s, Inches(0.55), Inches(1.0), Inches(12.3), Inches(1.35))
_set(tf.paragraphs[0],
     "出願前に社内で『悪魔の代弁者』役のAIレビュアー8観点がエビデンスを66回攻撃 → 52件は防御成功(コード正当)、"
     "『言い過ぎ』5件+『説明不足』5件を確定し全て修正、真のロジック欠陥0件、PENDING 2件(FTO調査等)。",
     size=12.5, bold=True, color=NAVY)
q = tf.add_paragraph()
_set(q, "修正はすべて文書・集計レイヤのみで、解析ロジック・数値・ラベルは不変。修正後テスト108件pass。以下は資料全体を拘束する必須開示。",
     size=11.5, color=GREY)
make_table(
    s, ["#", "必須開示事項"],
    [["1", "候補②(IC2)の非対称リセットは本解析より前にproductionコード実装済み。v4は『特徴付け・検証』であり『着想』ではない。新規性は着想日に依存(法務確認事項)"],
     ["2", "『技術エビデンスが強い』≠『特許が取れる』。エビデンス強度と先行技術リスクは独立の軸として管理"],
     ["3", "候補①(IC1)は広義(『機会後の無応答検出』一般)では出さない。80/20/80帯・72h窓・50mWh・censor除外を明示した狭め/中位で出願準備"],
     ["4", "50mWhは事前指定の設計値。データ(GMM谷35.2、CI[26.3,54.1])は独立の裏付けだが一意の最適証明ではない。具体値はobviousness-to-tryに脆弱"],
     ["5", "先行技術(US7610172ほか)は全てAIサーベイ由来で未検証(UNVERIFIED)。正式なFTO/特許性調査+弁理士意見が出願前必須"],
     ["6", "介入・BIOS/EC/電池FWバージョンのデータは存在しない(NOT AVAILABLE)。将来④⑤に効果実証データは無い"],
     ["7", "MLはnear-random(規範AUC 0.56)であり進歩性の根拠に不使用。発明の核は決定論カウンタ+状態機械(=検算可能・非ブラックボックス)"],
     ["8", "数値不整合1件(保持検証MAE 0.0375/0.05の混在)を検出し同一run由来に統一・修正。結論への影響なし"]],
    Inches(0.5), Inches(2.32), [0.45, 11.9], row_h=0.44, cell_size=10,
    styler=lambda i, j, v, b: (b, j == 0, RED if j == 0 else INK))
takeaway(s, "66回攻撃して計算のバグは0件 — ただし『言い方が強すぎた』10箇所は直した。この開示リストごと信じてもらうのが本資料の流儀。")
footer(s)

# =========================================================================== #
# 34. 正直な開示② 限界と残課題
# =========================================================================== #
s = section("正直な開示② 限界と出願前の残課題", RED, "§6 開示")
tf = _tb(s, Inches(0.55), Inches(1.02), Inches(12.3), Inches(2.9))
_set(tf.paragraphs[0], "限界（この資料が主張できないこと）", size=14, bold=True, color=NAVY)
bullets(tf, [
    "proxyラベルは本番システムの出力であり地上真実ではない。precision/recall等は全てproxy基準(『FW不良を検出』ではなく『FW確認対象を抽出』)",
    "A6(production自身との一致 precision/recall=1.0)は同義反復のため実証に使わない。使うのはアブレーションA0→A5(proxy精度0.33→0.89)",
    "END汚染0は『実証された発見』ではなく『定義による構造的性質』(設計の正しさの確認)",
    "極端な欠測(MCAR 50%)では機会回収率が約0.84まで低下(現実的なギャップ・打ち切りregimeでは頑健)",
    "bit単位の完全再現は非保証(seeded RNG+リプレイのテストでアルゴリズム的再現を担保)"], size=11.5, first=False)
tf = _tb(s, Inches(0.55), Inches(4.0), Inches(12.3), Inches(2.2))
_set(tf.paragraphs[0], "出願前の残課題（6/18特許性評価が『時期尚早』とした3点+PENDING）", size=14, bold=True, color=RED)
bullets(tf, [
    "(a) 介入→FCC回復のクローズドループ実証データが未取得(最大のギャップ。将来④のプロトコルで収集開始が可能)",
    "(b) 主要先行特許(バッテリフリート監視系 US12061240 等)の独立クレーム未精読 — AIサーベイのリストは下調べ段階",
    "(c) 社内外の公開有無(論文・製品ドキュメント・本テレメトリ機能自体の公開状況)の確認 — 新規性喪失の例外適用可否も含め法務確認",
    "正式なFTO/特許性調査(弁理士)・請求項チャートの作成"], size=11.5, first=False)
takeaway(s, "技術の証拠集めは完了。ここから先(先行技術の精査・公開有無・介入実証)は、reviewerと弁理士の判断と追加データが必要な領域。")
footer(s)

# =========================================================================== #
# 35. 知財戦略オプション
# =========================================================================== #
s = section("知財戦略オプション — 2つの評価の突き合わせ（判断材料）", PURPLE, "§6 開示")
tf = _tb(s, Inches(0.55), Inches(0.98), Inches(12.3), Inches(0.75))
_set(tf.paragraphs[0],
     "6/18の特許性評価(独立視点・v4のPENDING解析完了前)と、v4エビデンス完了後の整理は一部で結論が異なる。両方を並べて示す — どちらを採るかはreviewer判断。",
     size=12, bold=True, color=NAVY)
make_table(
    s, ["対象", "6/18 特許性評価の提言", "v4エビデンス完了後の位置づけ"],
    [["IC1（候補①の核）", "中・最有望 — 出願の核に", "STRONG / リスクMEDIUM-HIGH — 狭め/中位で出願候補"],
     ["IC5（候補③）", "中・最も堅い — 出願の核に", "STRONG(v3のMEDIUMから昇格) / リスクMEDIUM-HIGH — 最小状態を核に"],
     ["IC6（候補①内の品質ゲート）", "低〜中 — 従属クレームで", "STRONG — 段階ティア×censor-awareが差別化点"],
     ["IC2（候補②）", "低 — 防御的公開へ", "STRONG / リスクHIGH(着想日依存) — 出願か防御的公開か要判断"],
     ["IC7（統計計算）", "出願不可(純粋な数学的方法) — 復活させない", "同左(PROSPECTIVE)"],
     ["IC8（トリアージ/局在）", "低", "スクリーニングMEDIUM / version局在PROSPECTIVE"],
     ["微調整閾値(係数・定数類)", "営業秘密として保持", "同左(クレームに含めない)"]],
    Inches(0.5), Inches(1.72), [2.9, 4.4, 5.0], row_h=0.46, cell_size=10,
    styler=lambda i, j, v, b: (b, j == 0, NAVY if j == 0 else INK))
tf = _tb(s, Inches(0.55), Inches(5.55), Inches(12.3), Inches(0.85))
bullets(tf, [
    "6/18評価の推奨出願ルート: JP先願 → PCT → US/EP（総合判定『条件付きで出願に値する』）",
    "相違が残る主点はIC2の扱い(防御的公開 vs 着想日確認のうえ出願) — 着想日・公開有無の事実確認が判断の前提"], size=11.5)
takeaway(s, "選択肢は『出願する』だけではない。出願(IC1+IC5核)・営業秘密(閾値詳細)・防御的公開(IC2ほか)の3枚のカードをどう切るかがreviewerへの問い。")
footer(s)

# =========================================================================== #
# 36. まとめ
# =========================================================================== #
s = section("まとめ — 社内reviewer向け判断材料と依頼事項", BLUE, "§6 Conclusion")
make_table(
    s, ["特許候補（集約後）", "技術エビデンス", "新規性リスク(UNVERIFIED)", "出願準備度の目安"],
    [["① 機会条件付き無応答監査(中核)", "STRONG", "MEDIUM-HIGH", "出願候補(狭め/中位スコープで)"],
     ["② 非対称リセット状態機械", "STRONG", "HIGH", "出願候補(着想日要確認・規則明示) or 防御的公開"],
     ["③ 有界保持の因果証拠台帳", "STRONG", "MEDIUM-HIGH", "出願候補(最小状態を核に)"],
     ["(将来④) クローズドループ介入", "PROSPECTIVE", "—", "継続/将来開示(データ収集の判断待ち)"],
     ["(将来⑤) 機種非依存+version局在", "MEDIUM/PROSPECTIVE", "MEDIUM", "スクリーニングは候補/局在は継続"]],
    Inches(0.45), Inches(1.05), [4.5, 2.3, 2.9, 3.1], row_h=0.5, cell_size=10.5,
    styler=status_styler({1, 2}))
tf = _tb(s, Inches(0.5), Inches(4.0), Inches(12.4), Inches(2.3))
_set(tf.paragraphs[0], "reviewerへの依頼（4点）", size=14, bold=True, color=NAVY)
bullets(tf, [
    "① 候補①〜③の特許化の粒度・範囲の判断（狭め/中位/防御的公開の選択。判断材料は本資料の各候補スライドと知財戦略スライド）",
    "② 候補②(IC2)の着想日と社内外公開有無の事実確認（新規性の成否を左右する法的前提）",
    "③ 正式なFTO/特許性調査（弁理士）の発注判断 — 本資料の先行技術リストは全てUNVERIFIED",
    "④ 将来④の介入データ収集（プロトコル・スキーマ・必要台数試算は準備済み）を開始するかの判断"], size=12, first=False)
takeaway(s, "本題の判定技術はテレメトリのみ・欠測/保持制約下でも誤判定を抑えて機能する(5本柱で実証済み)。特許化の粒度・範囲の最終判断は社内reviewerに委ねる。")
footer(s)

# =========================================================================== #
# 37. 付録A 特徴量一覧
# =========================================================================== #
s = features("付録A 特徴量エンジニアリング一覧 — 各変数の作成方法", "§7 付録", [
    ("RSOC / FCC / cycleCount / DesignCapacity", RAW, "テレメトリ生データ。SoH=FCC×100/DesignCapacity"),
    ("design_capacity(復元)", ENG, "median(FCC×100/soh_design_pct)(単位曖昧回避): recover_design_mwh"),
    ("FCCステップ/abs_step/is_effective", ENG, "時刻順FCC差分の非ゼロ値、|Δ|≥50mWhで有効: fcc_step_indicator"),
    ("量子化単位", ENG, "観測された最小の正の|step|(=10mWh)"),
    ("学習機会(start/low/end)/episode_id", ENG, "RSOC high→low→high 状態機械 / user|band|start|end"),
    ("max_gap / quality_tier / coverage", ENG, "エピソード内ギャップ→段階品質、観測カバレッジ率: online_gap_quality"),
    ("応答ステータス/応答遅延", ENG, "END起点窓の有効step有無+窓完全性 / end→初回有効stepの時間"),
    ("直近any/effective変化, pending, seen_ids, FSM", ENG, "因果リプレイの永続最小状態: online_step_state"),
    ("累積無応答/機会カウンタ", ENG, "帯別・品質別の累積(FW/ゲージ二分岐を駆動)"),
    ("適応閾値", ENG, "max(k·quantization, α·DesignCapacity, noise percentile): C3"),
    ("行動特徴(ac_time_ratio, cycles_per_year, min_rsoc ...)", ENG, "RSOC/使用挙動由来33特徴(機種非依存): build_rsoc_features"),
    ("BIOS/EC/FW version, 介入結果", NA, "現データに存在しない。将来スキーマのみ(捏造なし)"),
], row_h=0.4)
takeaway(s, "判定に効いているのは生データ4列からの『数え方の設計』。ハードウェア識別子はどの段でも学習に使わない。")

# =========================================================================== #
# 38. 付録B データセットと再現性
# =========================================================================== #
s = section("付録B データセットと再現性", GREY, "§7 付録")
tf = _tb(s, Inches(0.55), Inches(1.02), Inches(6.2), Inches(5.2))
_set(tf.paragraphs[0], "データセット", size=14, bold=True, color=NAVY)
bullets(tf, [
    "取得元: S3フリートテレメトリ(1,808コレクションを走査)",
    "実バッテリ履歴のある752ユーザーを全数解析(≥20KBで空のプリロード/テスト機を除外)",
    "時系列 3,130,394行×22列(PWMログ、約30分間隔ロガー)",
    "変数定義はPower Manager PWM Log仕様PDF準拠(RSOC=remaining/FCC×100を実データで確認 corr 0.9999)",
    "パック交換 0/752(シリアル不変) → FCC変化は真のゲージ再学習",
    "学習機会エピソード 24,711件(うち大ギャップ20,314 / OK品質4,397)"], size=11.5, first=False)
tf = _tb(s, Inches(7.0), Inches(1.02), Inches(5.9), Inches(5.2))
_set(tf.paragraphs[0], "再現性・完全性", size=14, bold=True, color=NAVY)
bullets(tf, [
    "ベースライン再現ゲート: 全期間7指標+rolling-v2 9ラベルを期待値と完全照合(16/16 PASS)。不一致なら結論を出さない設計",
    "公開40アーティファクト全てにSHA-256チェックサム(manifest)",
    "全解析モジュールで決定的シード。修正後テスト108件pass",
    "PII: 公開CSV26+parquet9+図19枚にuser_id/serial/UUIDなし。IDは一方向ハッシュ(12桁)に置換",
    "介入・BIOS/EC/電池FWバージョン列の不在をavailability probeで機械的に確認(NOT AVAILABLEの根拠)"], size=11.5, first=False)
takeaway(s, "『数字を出し直したら変わる』ことがないよう、再現ゲート・チェックサム・決定的シードで固定。個人特定情報は全成果物から除去済み。")
footer(s)

# =========================================================================== #
# 39. 付録C 代替実施形態
# =========================================================================== #
s = section("付録C 代替実施形態の範囲（クレームの広がりの材料）", GREY, "§7 付録")
make_table(
    s, ["パラメータ", "主実施形態", "代替範囲"],
    [["RSOC帯(機会定義)", "80/20/80", "70/30/70 〜 90/10/90(strict)"],
     ["応答窓", "72h", "24 / 48 / 120 / 168h"],
     ["応答アンカー", "END(再充電完了)", "start / low(比較でENDの優位を実証済み)"],
     ["有効ステップ閾値", "固定50mWh", "固定10〜100mWh / DesignCapacity比0.05〜0.5% / 適応 max(k·量子化, α·設計容量, ノイズ百分位) / 混合分布・変化点・ユーザー別ノイズ推定"],
     ["リセット規則", "非対称(microはanyのみ)", "any単独 / effective単独 / 対称×2 / 適応型非対称(D5)"],
     ["生データ保持", "30日・stride 1日", "7〜90日・full / stride 1・7日 / 整列オフセット0〜6"],
     ["ギャップ/censor処理", "段階ティア+censor-aware", "naive / 二値ゲート / 段階ティアのみ"],
     ["永続状態", "最小7要素", "部分集合は等価性が破綻(アブレーションで実証)→ 下限として機能"]],
    Inches(0.45), Inches(1.05), [2.5, 2.6, 7.3], row_h=0.56, cell_size=10,
    styler=lambda i, j, v, b: (b, j == 1, DGREEN if j == 1 else INK))
tf = _tb(s, Inches(0.5), Inches(5.85), Inches(12.4), Inches(0.55))
_set(tf.paragraphs[0],
     "per-sample電流テーパ/電圧/温度/休止を使う実施形態は『その列が存在する実装で追加可能』の位置づけ(本データにはNOT AVAILABLE)。",
     size=11, color=GREY)
takeaway(s, "主実施形態の数値は一点張りではない。各パラメータに実証済みの代替範囲があり、クレームの広さはこの表から選べる。")
footer(s)

# =========================================================================== #
# 40. 付録D 先行技術差別化マトリクス
# =========================================================================== #
s = section("付録D 先行技術との差別化（全件 UNVERIFIED・弁理士確認前）", GREY, "§7 付録")
make_table(
    s, ["先行技術(UNVERIFIED)", "教示していること", "教示していないこと（差別化点の候補）"],
    [["US7610172ほか(非発生イベント監視)", "『起きるはずのイベントが起きない』ことの監視一般", "燃料計の物理的学習機会の抽出・FCC有効ステップ・censor-aware除外"],
     ["TI US6832171(Impedance Track)", "ゲージ内部での適格放電によるFCC/Qmax学習", "上位テレメトリ層での機会反復×無応答の検出・原因二分岐"],
     ["汎用デッドバンド/ヒステリシス", "小変化を無視する閾値処理", "microがany系統のみリセットしeffective証拠を保持する非対称規則"],
     ["US20130085715 / US9218527(ストリーミング異常検知)", "スライディング窓・オンライン処理", "窓越し未解決エピソードの確認・complete<reset<deadline順序・exactly-once計数"],
     ["Qualcomm US9330257", "他分野での識別子除外", "分類後の記述的偏在集計としての限定利用"],
     ["Song 2007(条件付き異常検知)", "規範比較の枠組み", "本件ではnear-randomと正直に開示し進歩性の根拠に不使用"],
     ["US12061240(バッテリフリート監視)", "※独立クレーム未精読(スキャンPDF)", "— 出願前に精読必須(6/18評価の残課題(b))"]],
    Inches(0.45), Inches(1.05), [3.3, 4.0, 5.1], row_h=0.56, cell_size=9.5,
    styler=lambda i, j, v, b: (b, j == 0, NAVY if j == 0 else INK))
tf = _tb(s, Inches(0.5), Inches(5.85), Inches(12.4), Inches(0.55))
_set(tf.paragraphs[0],
     "この表はAIサーベイによる下調べ(候補出し)であり、原文確認・FTO/特許性調査は未実施。『似た特許はない』という主張ではない。",
     size=11, bold=True, color=RED)
takeaway(s, "『これが先行技術です』ではなく『弁理士に最初に見てもらう候補リスト』。差別化欄はクレームドラフトの叩き台。")
footer(s)

# =========================================================================== #
# 41-44. 付録E 補足図
# =========================================================================== #
APX = "§7 付録 ─ 補足図"
figure("補足図1 A2の帰無分布 — 5種のシャッフルすべてで実測に届かない", APX,
       "negative_control_randomization_distribution.png",
       "根拠①の分布版。5種類のランダム化(シャッフル)を何百回も繰り返して得た『偶然ならこの範囲』の分布と、実測値の位置。",
       "横軸=対照リプリケートの72h応答確率。縦軸=度数。色付きの山=各ランダム化手法の帰無分布(0.22〜0.33)。赤の縦線=実測値(約0.39)。",
       "どの帰無分布も0.33以下に収まり、実測0.39はすべての山の右外側に明確なギャップを持って位置する。"
       "厳しめの対照(matched_pseudo・rsoc_phase_shift)でも届かない。",
       "山と赤線の間に重なりが無いことを見る。紫(rsoc_phase_shift)の山が低いのはリプリケート数が少ないため。"
       "赤線の上端は凡例ボックスに一部隠れる。",
       "『偶然ならこうなる』を5通り×数百回試しても、実測には一度も届かなかった — それが山と赤線の隙間の意味。")
figure("補足図2 リセット規則6案のトレードオフ（C2）", APX,
       "dual_track_reset_semantics.png",
       "6つのリセット規則(D0〜D5)を『証拠をどれだけ残せるか』と『ユーザーへの強い操作要求がどれだけ出るか』の2軸で比較。",
       "横軸=規則(D0〜D5)。緑棒=確定無応答の証拠を保持できた件数(高いほど良い)。赤棒=ハードアクション指示数(低いほど良い)。",
       f"対称系(D0/D2)は赤が少ないが緑も少ない(≈360)。effective単独(D1/D3)は緑が多い(≈475)が赤が倍増(≈210)。"
       f"非対称D4(提案)は緑≈475のまま赤≈97 — 両立を達成。適応型D5はさらに良い(緑≈600/赤≈83)。",
       "緑が高く赤が低い方式ほど良い。D4が主実装、D5(適応閾値)は代替実施形態として確保(付録C)。"
       "数値は棒に印字されておらず軸からの読み取り。",
       "『証拠を残す』と『ユーザーを煩わせない』は普通トレードオフ — 非対称リセットはその両取りができることを示す図。")
figure("補足図3 ストレージ比と等価性 — statelessは増やしても届かない（D）", APX,
       "storage_vs_equivalence.png",
       "保存データ量(全生データ比)を横軸に、全履歴版との判定一致率を縦軸に取った、台帳の費用対効果の図。",
       "横軸=ストレージ比(0.04〜0.52)。縦軸=応答判定の一致率(1.00=完全一致)。橙=stateful(台帳あり)、青=stateless(台帳なし)。",
       f"statefulはストレージ比≈{D['min_stateful_equivalent_storage_ratio']:.2f}(約4%)でも一致率1.00の水平線。"
       "statelessは保存を増やしても≈0.986で頭打ちし、1.00に到達しない。",
       "橙が最初から1.00に張り付いていること、青が右へ行っても届かないことの2点を見る。"
       "『データを増やす』では解決せず『何を覚えるか(台帳)』が本質であることを示す。",
       "力業(保存を増やす)では完全一致に届かない。賢いメモ(台帳)なら4%の容量で足りる — 発明の価値がこの1枚に出ている。")
figure("補足図4 打ち切り注入の安全性（E）", APX,
       "censor_injection_safety.png",
       "観測を人工的に途中で打ち切った(truncate)ときに、『本当は保留にすべきケース』を誤って無応答と確定してしまう件数の比較。",
       "横軸=打ち切りregime(72h/168h時点で切断)。縦軸=誤確定無応答の件数。4方式(naive/graded/binary gate/proposed)。",
       "naive(緑)だけが両regimeで約650件を誤確定。graded/binary gate/proposed(提案)は0件(棒が見えない=ゼロ)。",
       "棒が『無い』ことが良い結果。censored(打ち切り)を無応答に数えないという規律の効果が最も端的に出る検証。"
       "y軸ラベルの末尾が画像上端で切れている点はご容赦。",
       "『データがそこで終わっただけ』を『応答しなかった』と混同しない — 冤罪防止の仕組みが機能している証拠。")

# =========================================================================== #
# page numbers
# =========================================================================== #
for idx, slide in enumerate(prs.slides, 1):
    tf = _tb(slide, SW - Inches(0.75), SH - Inches(0.38), Inches(0.55), Inches(0.3))
    _set(tf.paragraphs[0], str(idx), size=9, color=GREY, align=2)

prs.save(str(OUT))
print(f"pptx -> {OUT.relative_to(pc.REPO)}  ({len(prs.slides._sldIdLst)} slides)")
