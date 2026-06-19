#!/usr/bin/env python
"""FCC更新停止者の判定技術 — サマリ deck (.pptx).

主役は『テレメトリのみから FCC更新停止（ゲージ凍結）者をどう判定するか』。
その実現過程で生まれた特許候補アイデアを少数(3+将来2)に集約して提示する。
特許化の粒度判断は社内reviewerに委ねる前提（判断材料を併載）。

Part1 本題: 目的と難しさ → 失敗した素朴手法(発想転換) → 提案手法の全体像 →
          判定法が機能する根拠(結果グラフを4ブロック解説)。
Part2 過程で生まれた特許候補: 候補①〜③ + 将来候補。各候補に背景/問題/課題/解決法/
          特許性、必要特徴量、根拠図。
最後: 特徴量エンジニアリング派生表 + reviewer向け判断材料。

数値は produced v4 results から引用(traceable)。図は匿名PNG。NOT a legal opinion.
"""
from __future__ import annotations

import json
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from battery_usage import patent_common_v4 as pc

R = json.load(open(pc.V4_DIR / "_v4_results_summary.json", encoding="utf-8"))
A2, A3, B, C2, C3, D, DM, E = (R["A2"], R["A3"], R["B"], R["C2"], R["C3"],
                               R["D"], R["Dmin"], R["E"])
FIG = pc.FIG_DIR
OUT = pc.REPORTS / "fcc_patent_summary_slides_v4.pptx"

NAVY = RGBColor(0x1F, 0x37, 0x5B); BLUE = RGBColor(0x1F, 0x77, 0xB4)
GREY = RGBColor(0x55, 0x55, 0x55); RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2C, 0xA0, 0x2C); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7); ORANGE = RGBColor(0xB0, 0x5A, 0x00)
TEAL = RGBColor(0x2E, 0x6E, 0x6A)

DISC = ("技術的特許性エビデンス（NOT a legal opinion）。先行技術は UNVERIFIED。"
        "特許化の粒度・範囲の最終判断は社内reviewer/弁理士。捏造（地上真実/介入/FW/因果結論）なし。")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size=14, bold=False, color=None, align=None):
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size); run.font.bold = bold
        if color:
            run.font.color.rgb = color
    if align:
        p.alignment = align


def footer(slide):
    tf = _tb(slide, Inches(0.3), SH - Inches(0.4), SW - Inches(0.6), Inches(0.33))
    _set(tf.paragraphs[0], DISC, size=8, color=GREY)


def band(slide, color, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def section(title, accent=BLUE, tag=None):
    s = prs.slides.add_slide(BLANK)
    band(s, accent, Inches(0.9))
    tf = _tb(s, Inches(0.5), Inches(0.14), SW - Inches(2.6), Inches(0.66))
    _set(tf.paragraphs[0], title, size=22, bold=True, color=WHITE)
    if tag:
        tt = _tb(s, SW - Inches(2.5), Inches(0.2), Inches(2.2), Inches(0.5))
        _set(tt.paragraphs[0], tag, size=12, bold=True, color=RGBColor(0xDD, 0xE8, 0xF2), align=2)
    return s


def box(slide, left, top, w, h, text, fill, fg=WHITE, size=12, sub=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.color.rgb = NAVY
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    _set(tf.paragraphs[0], text, size=size, bold=True, color=fg, align=2)
    if sub:
        p = tf.add_paragraph(); _set(p, sub, size=9.5, color=fg, align=2)
    return shp


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = NAVY; c.line.width = Pt(2.25)


# -------- reusable content slides --------
def glossary(title, items, tag="用語定義"):
    s = section(title, RGBColor(0x33, 0x66, 0x99), tag)
    tf = _tb(s, Inches(0.55), Inches(1.02), SW - Inches(1.1), Inches(5.95))
    first = True
    for term, defn in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(4)
        r = p.add_run(); r.text = f"{term}　"; r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = NAVY
        r2 = p.add_run(); r2.text = defn; r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    footer(s)


def concept(title, tag, strength, risk, bg, prob, chal, sol, novelty, accent=TEAL):
    s = section(title, accent, tag)
    tf = _tb(s, Inches(0.5), Inches(0.96), SW - Inches(1.0), Inches(0.42))
    p = tf.paragraphs[0]; p.text = ""
    for txt, col in [("技術エビデンス強度: ", NAVY),
                     (strength + "    ", GREEN if "STRONG" in strength else (GREY if "PROSPECTIVE" in strength else BLUE)),
                     ("先行技術リスク(UNVERIFIED): ", NAVY), (risk, RED if "HIGH" in risk else GREY)]:
        r = p.add_run(); r.text = txt; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = col
    tf = _tb(s, Inches(0.5), Inches(1.45), SW - Inches(1.0), Inches(5.55))
    first = True
    for head, body in [("背景", bg), ("問題", prob), ("課題", chal), ("解決法", sol), ("特許性 / 新規性", novelty)]:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        _set(p, head, size=14, bold=True, color=NAVY); p.space_before = Pt(5)
        q = tf.add_paragraph(); _set(q, body, size=12.5)
    footer(s)


def features(title, tag, rows):
    s = section(title, RGBColor(0x2E, 0x5A, 0x88), tag)
    headers = ["変数 / 特徴量", "種別", "作成方法（特徴量エンジニアリング）"]
    n = len(rows) + 1
    tbl = s.shapes.add_table(n, 3, Inches(0.4), Inches(1.02), Inches(12.5), Inches(0.4) * n).table
    tbl.columns[0].width = Inches(3.4); tbl.columns[1].width = Inches(1.35); tbl.columns[2].width = Inches(7.75)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for pa in c.text_frame.paragraphs:
            for run in pa.runs:
                run.font.size = Pt(11.5); run.font.bold = True; run.font.color.rgb = WHITE
    for i, (name, kind, deriv) in enumerate(rows, 1):
        for j, v in enumerate((name, kind, deriv)):
            c = tbl.cell(i, j); c.text = v; c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            for pa in c.text_frame.paragraphs:
                for run in pa.runs:
                    run.font.size = Pt(10.5)
                    if j == 1:
                        run.font.bold = True
                        run.font.color.rgb = GREEN if v == "派生" else (RED if "NOT" in v else GREY)
    footer(s)


def figure(title, tag, img, what, axes, claim, interp, accent=RGBColor(0x37, 0x6B, 0x3A)):
    s = section(title, accent, tag)
    p = FIG / img
    if p.exists():
        from PIL import Image
        try:
            iw, ih = Image.open(p).size
        except Exception:
            iw, ih = 1600, 1000
        scale = min(Inches(7.1) / iw, Inches(5.2) / ih)
        s.shapes.add_picture(str(p), Inches(0.4), Inches(1.12),
                             width=Emu(int(iw * scale)), height=Emu(int(ih * scale)))
        cap = _tb(s, Inches(0.4), Inches(6.5), Inches(7.1), Inches(0.4))
        _set(cap.paragraphs[0], f"図: {img}（dpi=300・匿名）", size=9, color=GREY)
    tf = _tb(s, Inches(7.75), Inches(1.05), Inches(5.25), Inches(5.8))
    first = True
    for head, body, col in [("【何のグラフか】", what, BLUE), ("【軸・変数の定義】", axes, NAVY),
                            ("【主張（言いたいこと）】", claim, GREEN), ("【読み方・解釈】", interp, ORANGE)]:
        p1 = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        _set(p1, head, size=12.5, bold=True, color=col); p1.space_before = Pt(6)
        p2 = tf.add_paragraph(); _set(p2, body, size=11.5, color=RGBColor(0x22, 0x22, 0x22))
    footer(s)


def table_slide(title, headers, data, accent, colw, hl=None, tag=None, note=None):
    s = section(title, accent, tag)
    n = len(data) + 1
    tbl = s.shapes.add_table(n, len(headers), Inches(0.4), Inches(1.02),
                             Inches(sum(colw)), Inches(0.42) * n).table
    for j, w in enumerate(colw):
        tbl.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for pa in c.text_frame.paragraphs:
            for run in pa.runs:
                run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = WHITE
    for i, rowv in enumerate(data, 1):
        for j, v in enumerate(rowv):
            c = tbl.cell(i, j); c.text = v; c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            for pa in c.text_frame.paragraphs:
                for run in pa.runs:
                    run.font.size = Pt(10)
                    if hl is not None and j == hl:
                        run.font.bold = True
                        run.font.color.rgb = (GREEN if v == "STRONG" else RED if v == "HIGH"
                                              else GREY if v in ("PROSPECTIVE", "WEAK", "派生", "NOT AVAIL") else BLUE)
    if note:
        tf = _tb(s, Inches(0.4), Inches(1.05) + Inches(0.42) * n + Inches(0.15), Inches(12.5), Inches(0.9))
        _set(tf.paragraphs[0], note, size=12, bold=True, color=RED)
    footer(s)
    return s


RAW, ENG, NA = "raw", "派生", "NOT AVAIL"

# =========================================================================== #
# Title
# =========================================================================== #
s = prs.slides.add_slide(BLANK); band(s, NAVY, SH)
tf = _tb(s, Inches(0.8), Inches(1.7), SW - Inches(1.6), Inches(2.2))
_set(tf.paragraphs[0], "テレメトリからの FCC更新停止（ゲージ凍結）者の判定技術", size=34, bold=True, color=WHITE)
p = tf.add_paragraph()
_set(p, "— と、その実現過程で生まれた特許候補アイデア —", size=20, color=RGBColor(0xCF, 0xDD, 0xEE))
tf2 = _tb(s, Inches(0.8), Inches(4.0), SW - Inches(1.6), Inches(2.6))
for i, ln in enumerate([
        "本題: ノートPCバッテリのテレメトリ(RSOC/FCC/サイクル等)だけで『満充電容量(FCC)を再学習しなくなった個体』を判定する",
        "母集団: 実バッテリ履歴 752ユーザー / 全期間 24,711 学習機会(エピソード)",
        "特許候補は3つ(+将来2つ)に集約。どこをどの粒度で特許化するかの最終判断は社内reviewerに委ねる",
        "技術効果は proxy ラベルに依存しない独立指標で検証。法的結論は主張しない"]):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    _set(p, "•  " + ln, size=15, color=WHITE)
footer(s)

# =========================================================================== #
# Glossary (terms up front, 2 slides)
# =========================================================================== #
glossary("用語定義 (1/2) — バッテリと学習機会", [
    ("FCC（満充電容量）", "燃料計が学習する『今の満充電容量』(mWh)。劣化で減る。充放電で段階的に再学習・更新。"),
    ("FCC更新停止 / ゲージ凍結", "FCCが長期間更新されない状態。SoH(健全性=FCC×100/DesignCapacity)が動かなくなる。本技術の判定対象。"),
    ("RSOC（相対残量, %）", "= remainingCapacity/FCC×100。今の充電残量%。high=満充電付近、low=空付近。"),
    ("学習機会 / エピソード", "燃料計がFCCを再学習できる『深い放電→再充電』1サイクル。RSOCが high→low→high と動いた区間。"),
    ("START / LOW / END", "エピソードの 開始(高RSOC) / 谷(最深放電) / 終了(再び高RSOC=再充電完了)。END=機会の完了時刻。"),
    ("有効ステップ（≥50mWh）", "意味あるFCC更新=『学習応答』とみなす最小変化。量子化最小は10mWh、それ未満は micro-step(微小)。"),
    ("応答ステータス（END後72h窓）", "responded=有効step有 / no_response=窓を全観測したが無し / censored=窓未観測で未確定(無応答に数えない)。"),
])
glossary("用語定義 (2/2) — 検証手法と統計", [
    ("対照 / ヌル分布(null)", "効果が無い場合に期待される値。実データがその95%区間より外れれば『効果あり』と判断。"),
    ("負の対照検定(A2)", "機会END時刻をわざとズラした『ニセ機会』と比較し、応答が真の機会に特異かを検証。"),
    ("アンカー比較(A3)", "応答の起点を START/LOW/END で比較し『因果汚染(=再充電完了前の変化を応答と誤計上)』を定量化。"),
    ("応答ハザード(B)", "機会ENDからの経過時間に対する累積応答率(CIF)を生存解析で推定。"),
    ("欠測ストレス(E) / 保持グリッド(D)", "E=欠測・打ち切りを注入し誤判定耐性を検証 / D=生データ保持を短くしても全期間と等価か検証。"),
    ("user-bootstrap", "ユーザ単位で再標本し信頼区間を計算。同一ユーザの複数機会を独立扱いしない(過小評価を防ぐ)。"),
    ("dual-track / 非対称リセット", "any(全変化)とeffective(≥50mWh)を別系統で追跡。微小変化はany系統のみリセットしeffective証拠は保持。"),
    ("stateful / stateless", "永続状態あり/なしの処理。statefulは保持窓をまたぐ証拠を状態で回収できる。"),
])

# =========================================================================== #
# PART 1 — 本題: 目的・難しさ
# =========================================================================== #
s = section("本来の目的と、なぜ難しいか", NAVY, "Part 1 ─ 本題")
tf = _tb(s, Inches(0.6), Inches(1.15), Inches(6.0), Inches(5.6))
_set(tf.paragraphs[0], "目的", size=16, bold=True, color=BLUE)
q = tf.add_paragraph(); _set(q, "テレメトリ(RSOC・FCC・サイクル・時刻)だけから、FCCを再学習しなくなった個体(ゲージ凍結)をフリート規模で正しく判定する。", size=13)
p = tf.add_paragraph(); _set(p, "なぜ難しいか", size=16, bold=True, color=BLUE); p.space_before = Pt(12)
for ln in ["静的検査では区別不能: FCCが動かない理由が『浅い充放電で正常』か『要再較正』か『FW/HW起因』か判別できない",
           "欠測・睡眠ギャップ・記録打ち切り: データの穴を『無応答』と誤判定しやすい",
           "生データ保持が有界: 直近30日しか残せない等の制約で、過去の証拠が失われる",
           "機種非依存が必須: 機種名でハードコードすると過学習・汎化不能"]:
    q = tf.add_paragraph(); _set(q, "•  " + ln, size=12.5); q.space_after = Pt(4)
box(s, Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.0),
    "FCCが動かない = ?", BLUE, sub="静的検査では下の3つを区別できない")
for i, (txt, col) in enumerate([("正常\n(浅充放電)", GREEN), ("要再較正\n(ゲージ)", ORANGE), ("FW/HW\n起因", RED)]):
    box(s, Inches(7.0) + Inches(2.0) * i, Inches(4.0), Inches(1.85), Inches(1.3), txt, col, size=12)
footer(s)

# 失敗した素朴手法 → 発想転換
s = section("失敗した素朴アプローチ → 発想の転換", NAVY, "Part 1 ─ 本題")
tf = _tb(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.4))
_set(tf.paragraphs[0], "『使用挙動から凍結を予測』しようとしたが、解けなかった:", size=15, bold=True, color=NAVY)
for ln in ["教師あり予測(33特徴量, 公平比較領域): AUC ≈ 0.54 — ほぼランダム",
           "FCC履歴を除いた規範ML(異常スコア): AUC ≈ 0.56 — near-random",
           "→ 『予測』では解けない。最重要特徴 min_rsoc は『深放電ほど凍結』という反usage方向(交絡)。"]:
    q = tf.add_paragraph(); _set(q, "•  " + ln, size=13); q.space_after = Pt(3)
box(s, Inches(1.2), Inches(4.2), Inches(4.6), Inches(1.6),
    "予測アプローチ", GREY, sub="使用挙動→凍結を学習 (AUC≈0.54)")
arrow(s, Inches(5.9), Inches(5.0), Inches(7.3), Inches(5.0))
box(s, Inches(7.4), Inches(4.2), Inches(4.6), Inches(1.6),
    "機械的監査アプローチ（提案）", GREEN,
    sub="学習機会があったのに応答したか?を機会END起点で因果的に監査")
tf = _tb(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7))
_set(tf.paragraphs[0], "発想転換: 『凍結を当てる』のではなく『再学習する“機会”に対して実際に応答したかを機械的に監査する』。本技術の核。",
     size=13, bold=True, color=GREEN)
footer(s)

# 提案手法の全体像（パイプライン）
s = section("提案手法の全体像 — テレメトリから判定までの流れ", NAVY, "Part 1 ─ 本題")
steps = [
    ("① 学習機会の抽出", "RSOC high→low→high\n(深放電→再充電)"),
    ("② END起点で応答監査", "再充電完了後72h窓の\n有効ステップ≥50mWh"),
    ("③ 欠測/打ち切り耐性", "段階品質ティア+\ncensored除外"),
    ("④ 二分岐トリアージ", "機会反復×無応答→FW\n機会皆無→ゲージ再較正"),
    ("⑤ 有界保持で証拠保全", "最小状態の因果台帳\n(30日保持で全期間等価)"),
]
bw = Inches(2.36); gap = Inches(0.12); x0 = Inches(0.35); y = Inches(2.7)
for i, (t, sub) in enumerate(steps):
    x = x0 + (bw + gap) * i
    box(s, x, y, bw, Inches(1.7), t, [BLUE, TEAL, ORANGE, RGBColor(0x6A, 0x3D, 0x8A), NAVY][i], sub=sub, size=12)
    if i < 4:
        arrow(s, x + bw, y + Inches(0.85), x + bw + gap, y + Inches(0.85))
tf = _tb(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(1.4))
_set(tf.paragraphs[0], "入力: テレメトリ(RSOC/FCC/cycleCount/timestamp) のみ。ハードウェア識別子は判定に使わない(機種非依存)。",
     size=13, bold=True, color=NAVY)
q = tf.add_paragraph()
_set(q, "出力: 各ユーザを NORMAL / ゲージ再較正候補 / FW確認候補 / 判定保留(データ品質) に振り分け。"
     "①〜⑤の各段が後述の特許候補に対応する。", size=12.5)
tf = _tb(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(1.6))
for i, ln in enumerate([
        "①+③+④ → 特許候補①（機会条件付き無応答監査・中核判定法）",
        "②(微小変化の扱い) → 特許候補②（デュアルトラック非対称リセット状態機械）",
        "⑤(保持制約下の証拠保全) → 特許候補③（有界保持の因果証拠台帳・最小状態）"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    _set(p, "•  " + ln, size=12.5, color=GREY)
footer(s)

# 判定法が機能する根拠（結果図）
RES_TAG = "Part 1 ─ 判定法が機能する根拠"
figure("根拠① 応答は『真の機会』に特異（負の対照検定 A2）", RES_TAG, "negative_control_true_vs_null.png",
       "判定の前提『機会ENDの後に本当にFCCが応答する』を、END時刻をズラした5種のニセ機会(対照)と比較して確認。",
       "縦軸=応答確率(END後72h以内に有効ステップが起きた割合)。赤線=実データ。棒=各ニセ機会の応答率と95%区間。",
       f"実データの応答率({A2['true_resp_prob_72h']:.2f})はどのニセ機会より高い→応答は『真の機会END時刻』に特異。単なる時間経過・活動量・ユーザ差ではない。",
       f"赤線が棒(95%区間)より上ほど明確。{A2['n_controls_outside_null']}/{A2['n_controls_total']}の対照で区間外→判定の土台(機会-応答の因果)が成立。")
figure("根拠② 応答の起点はEND(再充電完了)が正しい（A3）", RES_TAG, "response_anchor_contamination.png",
       "応答を測る『起点』を START/LOW/END のどこに置くべきかを比較。誤った起点だと放電途中の変化を応答と誤計上する。",
       "縦軸=因果汚染率(応答とカウントしたステップのうち実は再充電完了前=放電途中だった割合)。横軸=3起点。",
       f"END起点なら汚染0。START起点だと{A3['start_contamination_frac_72h']:.2f}(=56%)が放電途中を誤計上→END起点が因果的に正しい。",
       "棒が低いほど良い。END=0が理想で、判定が『再充電完了後の応答』だけを見ていることを保証。")
figure("根拠③ 真の機会の後ほど速く応答する（応答ハザード B）", RES_TAG, "response_hazard_true_vs_pseudo.png",
       "機会ENDからの経過時間に対する累積応答率。真の機会 vs ランダムに置いたニセ機会(pseudo)。",
       "横軸=機会END後の経過時間(h)。縦軸=その時刻までに有効応答した累積割合(0〜1, 生存解析CIF)。",
       f"真の曲線がpseudoより上=真の機会後の方が速く・多く応答(72hで {B['true_cif_72h_50mwh']:.2f} vs {B['pseudo_cif_72h_50mwh']:.2f}、中央値約49h)。",
       "上の曲線ほど応答が早い/多い。差が機会の効果。応答が無いまま時間が経つ個体=更新停止候補。")
figure("根拠④ 欠測・打ち切りでも誤判定しない（E）", RES_TAG, "missingness_false_escalation.png",
       "データ欠測を人工注入したとき、各検出器が『誤って無応答と判定』した件数。実運用の穴に強いかを検証。",
       "縦軸=欠測regime(MCAR・連続ギャップ・末尾打ち切り・睡眠ギャップ等)。横軸=誤った確定無応答の件数。色=検出器。",
       f"naive(素朴)は欠測を無応答と誤判定し多発。提案(段階品質+censor考慮)は誤判定を {E['naive_mean_false_no_response']:.0f}→{E['proposed_mean_false_no_response']:.1f}に抑制。",
       "棒が短いほど良い。提案が全regimeで最短=穴のあるデータでも『見かけの無応答』に騙されない。")
figure("根拠⑤ 30日保持でも全期間と同じ判定（D）", RES_TAG, "retention_invariance_heatmap.png",
       "生データ保持を短くしても『全期間版と同じ結論』が出るか。フリート運用のストレージ制約下での妥当性。",
       "行=処理方式(stateful=状態あり/stateless=状態なし)。列=保持日数(7〜90)。セル=全期間版との一致率(緑=高,1.0=完全一致)。",
       f"状態あり方式は全保持期間で一致≈1.0。約{D['min_stateful_equivalent_storage_ratio']*100:.0f}%のストレージで全期間と同じ判定(重複0/誤差≈{D['stateful_verify_no_response_mae']})。",
       "statefulの行が全緑=有界保持でも判定が変わらない。statelessは短保持で劣化(状態を持つ価値)。")

# =========================================================================== #
# PART 2 — 過程で生まれた特許候補
# =========================================================================== #
s = section("過程で生まれた特許候補アイデア（位置づけ）", RGBColor(0x6A, 0x3D, 0x8A), "Part 2")
tf = _tb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(2.0))
_set(tf.paragraphs[0], "本技術(FCC更新停止者の判定)は単一目的の一連の手法。その実現過程で、単独でも特許になりうる中核アイデアが生まれた。",
     size=14, bold=True, color=NAVY)
q = tf.add_paragraph()
_set(q, "細分化しすぎると評価しづらいため、ここでは3つ(+将来2つ)に集約して提示する。"
     "どこをどの粒度で特許化するかの最終判断は社内reviewerに委ねる。各候補に技術エビデンス強度と新規性リスクを併記し判断材料を提供。",
     size=13, color=GREY); q.space_before = Pt(4)
table_slide  # noqa (placeholder to keep flake quiet)
# candidate overview table on same slide
data = [
    ["① 機会条件付き無応答監査（中核判定法）", "①+③+④", "STRONG", "MEDIUM-HIGH"],
    ["② デュアルトラック非対称リセット状態機械", "②(微小変化)", "STRONG", "HIGH"],
    ["③ 有界保持の因果証拠台帳（最小状態）", "⑤(保持制約)", "STRONG", "MEDIUM-HIGH"],
    ["(将来④) 診断依存クローズドループ介入検証", "介入後検証", "PROSPECTIVE", "—"],
    ["(将来⑤) 機種非依存スクリーニング+version局在", "原因局在", "MEDIUM", "MEDIUM"],
]
tbl = s.shapes.add_table(len(data) + 1, 4, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5) * (len(data) + 1)).table
for j, w in enumerate([5.4, 2.0, 2.4, 2.5]):
    tbl.columns[j].width = Inches(w)
for j, h in enumerate(["特許候補（集約後）", "対応する手法段", "技術エビデンス", "新規性リスク(UNVERIFIED)"]):
    c = tbl.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for pa in c.text_frame.paragraphs:
        for run in pa.runs:
            run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = WHITE
for i, rowv in enumerate(data, 1):
    for j, v in enumerate(rowv):
        c = tbl.cell(i, j); c.text = v; c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        for pa in c.text_frame.paragraphs:
            for run in pa.runs:
                run.font.size = Pt(10.5)
                if j == 2:
                    run.font.bold = True
                    run.font.color.rgb = GREEN if v == "STRONG" else (GREY if v in ("PROSPECTIVE",) else BLUE)
                if j == 3 and v == "HIGH":
                    run.font.bold = True; run.font.color.rgb = RED
footer(s)

# ----- 候補① -----
concept("特許候補① 機会条件付き無応答監査（中核判定法）", "Part 2 ─ 候補①  [IC1+IC6+二分岐 相当]",
        "STRONG", "MEDIUM-HIGH",
        "FCC再学習は『深い充放電(学習機会)』でしか起きない。だが機会の有無と応答の有無は静的検査では分離できない。",
        "FCC凍結を『機会が無いだけ(正常/再較正)』なのか『機会はあるのに応答しない(FW/HW疑い)』なのか判別できない。",
        "テレメトリだけで、機会の有無と応答の有無を因果的に分離し、欠測・打ち切りに騙されずに判定する。",
        "RSOC high→low→high 機会を抽出→END起点72h窓で有効ステップを responded/no_response/censored に分類→段階品質ティアとcensor除外で欠測を無応答に数えない→機会反復×無応答=FW候補/機会皆無=ゲージ再較正に二分岐。",
        "A2で機会-応答の特異性、A3でEND起点の因果汚染0、Eで欠測耐性を実証(根拠①②④)。判定の中核。【リスク】『非発生イベント監視』は広い先行技術→狭め(80/20/80・72h・50mWh・censor除外)で出願推奨。UNVERIFIED。")
features("特許候補① 必要な特徴量（と作成方法）", "Part 2 ─ 候補①", [
    ("RSOC / FCC / cycleCount / timestamp", RAW, "テレメトリ生データ（PWMログ）。識別子は使わない"),
    ("学習機会 episode (start/low/end)", ENG, "RSOC high→low→high 状態機械(80/20/80等): extract_high_low_high_episodes"),
    ("有効FCCステップ(≥50mWh)", ENG, "時刻順FCC差分の絶対値≥50mWh: fcc_step_indicator(FCC,50)"),
    ("応答ステータス response_status_72h", ENG, "END起点[end,end+72h]の有効step有無+窓完全性で responded/no_response/censored"),
    ("段階品質ティア quality_tier", ENG, "エピソード最大ギャップ→HIGH_OK/MEDIUM_GAP/LOW_LARGE_GAP: online_gap_quality"),
    ("二分岐スコア(FW/ゲージ)", ENG, "機会反復×無応答 vs 機会皆無 のカウンタから振り分け"),
    ("(学習除外) 機種名/ベンダー/FRU", "—", "ハードウェア識別子は判定に一切使わない(機種非依存)"),
])

# ----- 候補② -----
concept("特許候補② デュアルトラック非対称リセット状態機械", "Part 2 ─ 候補②  [IC2 + 有効閾値C3 相当]",
        "STRONG", "HIGH",
        f"燃料計は整数mWh量子化(最小{C3['quantization_unit_mwh']:.0f}mWh)。微小変化(micro)と意味ある再学習(effective≥50mWh)が混在(二峰: {C3['gmm_micro_mode_mwh']:.0f}/{C3['gmm_effective_mode_mwh']:.0f}mWh)。",
        "状態を1系統で持つと、microステップが『まだ未解決の無応答証拠』をリセットして消し、FW疑いを見逃す。",
        "微小ゆらぎで証拠を失わず、かつ『ソフト較正』と『ハードリセット』を区別する(過剰なハード指示を避ける)。",
        "any系統(任意変化でリセット)とeffective系統(≥50mWhのみリセット)を分離。microはany系統だけリセットしeffective状態/未解決証拠は保持(非対称)。同時刻は complete<reset<deadline の決定的順序。",
        f"C2で対称リセットが消す無応答{C2['d2_no_response_erased']}件を非対称が温存(+{C2['evidence_preserved_vs_symmetric']})、ハード指示{C2['hard_prompts_d1_effective_only']}→{C2['hard_prompts_d4_proposed']}。【リスク高】この設計は既にproduction実装済+deadband系は一般的先行技術→新規性は着想日依存(法的)。非対称規則を明示してクレーム。")
features("特許候補② 必要な特徴量（と作成方法）", "Part 2 ─ 候補②", [
    ("FCCステップ(符号付き/絶対値)", ENG, "時刻順FCC差分の非ゼロ値: patent_dual_track.fcc_steps"),
    ("is_effective / is_micro", ENG, "絶対値 ≥ / < 50mWh(有効閾値): step_threshold_mwh"),
    ("量子化単位 quantization_unit", ENG, "観測された最小の正の|step|(=10mWh)"),
    ("any/effective 状態・未解決証拠カウンタ", ENG, "イベント列の状態機械リプレイ(micro/effで非対称リセット): online_step_state"),
    ("混合分布 micro/effective モード・谷", ENG, "log10(|step|)に2成分Gaussian Mixtureを適合: patent_effective_threshold"),
    ("持続/反転(persistence/reversal)", ENG, "各stepの次step時刻差・符号反転(microがノイズか否かの検証)"),
    ("適応閾値(代替)", ENG, "max(k·quantization, α·DesignCapacity, noise percentile)"),
])
figure("候補②の根拠: 微小変化で証拠が消える vs 守る（C2）", "Part 2 ─ 候補②", "dual_track_erased_evidence.png",
       "microステップによって『消されてしまう』未解決証拠の件数を、リセット方式別に比較。",
       "縦軸=消去された pending(保留中の機会)+no_response(確定無応答)件数。横軸=リセット方式(D0〜D5)。",
       f"対称リセット(D0/D2)はmicroで pending {C2['d2_pending_erased']}・no_response {C2['d2_no_response_erased']}件を消去。非対称(D4,提案)は0。",
       "棒が高い=証拠を失う方式。D4=0が理想。これが『非対称リセットが証拠を守る』直接証拠。")
figure("候補②の根拠: 有効閾値はデータの谷で裏付け（C3）", "Part 2 ─ 候補②", "effective_threshold_mixture_fit.png",
       "FCCステップ大きさ(対数)の分布。微小ゆらぎと意味ある再学習を分ける閾値の妥当性。",
       f"横軸=log10(|FCCステップ|,mWh)。縦軸=密度。点線=候補閾値、黒線=50mWh。",
       f"分布は二峰(micro≈{C3['gmm_micro_mode_mwh']:.0f}/effective≈{C3['gmm_effective_mode_mwh']:.0f}mWh)。谷(≈{C3['gmm_valley_mwh']:.0f}mWh)の上に50mWh→閾値として妥当。",
       "2つの山の谷が境界。固定50mWhは一実施形態、広い概念は量子化・ノイズ帯を超える適応閾値(粒度はreviewer判断)。")

# ----- 候補③ -----
concept("特許候補③ 有界保持下の因果証拠台帳（最小十分状態）", "Part 2 ─ 候補③  [IC5 相当]",
        "STRONG", "MEDIUM-HIGH",
        "フリート監視では生テレメトリの保持期間が有界(例:直近30日のみ)。",
        f"保持窓をまたぐ機会の証拠が生データ破棄で失われる/重複する(状態なしは7日保持で recall {D['stateless_7d_recall_72h']:.2f}・重複率 {D['stateless_7d_dup_rate_72h']:.1f})。",
        "有界保持でも『全期間版と等価』な判定を、できるだけ小さい永続状態で実現する。",
        "部分FSM・pending期限キュー・確定済みID集合・直近any/effective変化を永続化し、生データ破棄後も因果リプレイで復元。期限は窓観測時のみ発火(未来情報リーク無)。",
        f"Dで状態あり方式が recall=1/重複=0/誤差≈{D['stateful_verify_no_response_mae']} を ストレージ {D['min_stateful_equivalent_storage_ratio']*100:.0f}% で達成(根拠⑤)。最小状態アブレーションで各構成要素が必要と実証。【リスク】streaming+cachingは既知の組合せ→『最小状態構造』をクレーム核に。")
features("特許候補③ 必要な特徴量（と作成方法）", "Part 2 ─ 候補③", [
    ("episode_id(正準ID)", ENG, "user|band|start|end(同一機会を一意識別、重複防止。公開時user部はhash)"),
    ("応答期限 / 直近変化 ts+cycle", ENG, "end+72h / 状態機械が逐次更新(経過日数・サイクルを駆動)"),
    ("確定済みID集合 seen_ids", ENG, "確定済み機会IDの集合(重複カウント防止)"),
    ("pending 期限キュー", ENG, "未解決機会の期限待ち(生データ破棄後も証拠保持)"),
    ("部分FSM 状態", ENG, "WAIT_HIGH/LOW/HIGH_AGAIN の途中状態(窓またぎ機会検出)"),
    ("最小十分状態(全体)", ENG, "上記の集合。アブレーションで各要素の必要性を実証: patent_state_minimality"),
])
figure("候補③の根拠: どの状態も削れない（最小性 D）", "Part 2 ─ 候補③", "minimal_state_necessity.png",
       "永続状態の各構成要素を1つずつ外したとき、どの不変量(等価性)が壊れるか。",
       "横軸=外した構成要素(fsm/pending/last_eff等)。縦軸=不変量の破綻量(recall低下=赤, no_response誤差=橙)。",
       "各要素を外すと等価性が崩れる→これらは『最小十分状態』であり削れない(=クレームの核)。",
       "棒が高い要素ほど必須。例: fsmを外すと窓またぎ機会を取り逃しrecall急落。")

# ----- 将来候補 -----
s = section("将来の特許候補（PROSPECTIVE）", GREY, "Part 2 ─ 将来候補")
tf = _tb(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.0))
_set(tf.paragraphs[0], "将来④ 診断依存クローズドループ介入検証", size=15, bold=True, color=NAVY)
for ln in ["背景/課題: ラベル付け後の介入(較正/FW更新)が本当にゲージを回復させたかを検証したい。",
           "解決法: 介入後の『次の良質な機会』での有効応答を観測しラベルを検証(介入を因果台帳に記録)。",
           "状態: 介入・BIOS/EC/FWバージョンが NOT AVAILABLE → 将来スキーマ・プロトコル・検出力シミュのみ(捏造なし)。"]:
    q = tf.add_paragraph(); _set(q, "•  " + ln, size=12.5)
p = tf.add_paragraph(); _set(p, "将来⑤ 機種非依存スクリーニング + version 局在", size=15, bold=True, color=NAVY); p.space_before = Pt(12)
for ln in ["背景/課題: 識別子を使わず振り分け、後段で原因をversionに局在させる二段構成。",
           "解決法: 分類は行動特徴のみ(機種名/ベンダー/FRUを学習に使わない)。versionは記述的層別としてのみ事後利用。",
           "状態: スクリーニングは実装済(MEDIUM)。version局在はBIOS/EC/FW列が NOT AVAILABLE → PROSPECTIVE。"]:
    q = tf.add_paragraph(); _set(q, "•  " + ln, size=12.5)
footer(s)

# =========================================================================== #
# 特徴量エンジニアリング 派生表 (judging材料)
# =========================================================================== #
FE = [
    ("RSOC / FCC / cycleCount / DesignCapacity", RAW, "テレメトリ生データ。SoH=FCC×100/DesignCapacity"),
    ("design_capacity(復元)", ENG, "median(FCC×100/soh_design_pct)(単位曖昧回避): recover_design_mwh"),
    ("FCCステップ/abs_step/is_effective", ENG, "時刻順FCC差分の非ゼロ値、|Δ|≥50mWhで有効: fcc_step_indicator"),
    ("量子化単位", ENG, "観測された最小の正の|step|(=10mWh)"),
    ("学習機会(start/low/end)/episode_id", ENG, "RSOC high→low→high 状態機械 / user|band|start|end"),
    ("max_gap / quality_tier / coverage", ENG, "エピソード内ギャップ→段階品質、観測カバレッジ率: online_gap_quality"),
    ("応答ステータス/応答遅延", ENG, "END起点窓の有効step有無+窓完全性 / end→初回有効stepの時間"),
    ("直近any/effective変化, pending, seen_ids, FSM", ENG, "因果リプレイの永続最小状態: online_step_state"),
    ("累積無応答/機会カウンタ", ENG, "帯別・品質別の累積(FW/ゲージ二分岐を駆動)"),
    ("適応閾値", ENG, "max(k·quantization, α·DesignCapacity, noise percentile): C3"),
    ("行動特徴(ac_time_ratio, cycles_per_year, min_rsoc ...)", ENG, "RSOC/使用挙動由来33特徴(機種非依存): build_rsoc_features"),
    ("BIOS/EC/FW version, 介入結果", NA, "現データに存在しない。将来スキーマのみ(捏造なし)"),
]
table_slide("特徴量エンジニアリング 一覧 — 各変数の作成方法（判断材料）", ["変数", "種別", "作成方法 / 由来"],
            FE, RGBColor(0x6A, 0x3D, 0x8A), [4.0, 1.2, 7.6], hl=1, tag="付録")

# =========================================================================== #
# まとめ（reviewer向け）
# =========================================================================== #
table_slide("まとめ — 社内reviewer向け判断材料",
    ["特許候補（集約後）", "技術エビデンス", "新規性リスク(UNVERIFIED)", "出願準備度の目安"],
    [["① 機会条件付き無応答監査(中核)", "STRONG", "MEDIUM-HIGH", "出願候補(狭め/中位で)"],
     ["② 非対称リセット状態機械", "STRONG", "HIGH", "出願候補(着想日要確認・規則明示)"],
     ["③ 有界保持の因果証拠台帳", "STRONG", "MEDIUM-HIGH", "出願候補(最小状態を核に)"],
     ["(将来④) クローズドループ介入", "PROSPECTIVE", "—", "継続/将来開示"],
     ["(将来⑤) 機種非依存+version局在", "MEDIUM", "MEDIUM", "スクリーニングは候補/局在は継続"]],
    BLUE, [4.4, 2.4, 3.0, 3.0], hl=1, tag="Conclusion",
    note="判定技術(本題)はテレメトリのみ・欠測/保持制約下でも誤判定を抑えて FCC更新停止者を判定できる。"
         "どこをどの粒度で特許化するかの最終判断は社内reviewerに委ねる。技術強度≠新規性、先行技術はUNVERIFIED。")

prs.save(str(OUT))
print(f"pptx -> {OUT.relative_to(pc.REPO)}  ({len(prs.slides._sldIdLst)} slides)")
