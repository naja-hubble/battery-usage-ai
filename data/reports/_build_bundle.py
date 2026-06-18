# -*- coding: utf-8 -*-
"""Curate key figures + build FIGURE_INDEX.md + generate patent-review PPTX.
Run from data/reports/. data/ is git-ignored so this stays local."""
import os, shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, "figures")
BUNDLE = os.path.join(BASE, "patent_review_bundle")
OUTFIG = os.path.join(BUNDLE, "04_figures")
JP = "Meiryo"

# slug -> (src under figures/, theme folder, caption, IC tag)
CUR = [
 # background / problem
 ("soh_overlay_by_class","soh_overlay_by_class.png","00_problem","原因クラス別 SoH 軌跡。凍結群は末尾でFCCが張り付く（静的には正常と判別困難）","背景"),
 ("very_stale_xgb_shap","very_stale_xgb_shap.png","00_problem","very_stale を使用挙動で予測したXGBoostのSHAP。最重要 min_rsoc は『深放電ほど凍結』の反usage方向","背景"),
 ("very_stale_tree","very_stale_tree.png","00_problem","解釈可能な決定木。公平領域AUC≈0.54でランダム同然＝行動で説明不能","背景"),
 ("cohort_soh_vs_cycles","cohort_soh_vs_cycles.png","00_problem","SoH vs サイクル数の散布（コホート）","背景"),
 ("soh_reason_trends","soh_reason_trends.png","00_problem","凍結トレンド要約（ベンダ別・HW疑い群機種構成・サイクル散布・クラス内訳）","背景"),
 # IC1 core
 ("opportunity_vs_response","fcc_action/opportunity_vs_response.png","01_core_ic1","【核】学習機会(opportunity) × FCC応答(response) の関係。機会あり×応答なしが検出対象","IC1"),
 ("final_funnel_counts","fcc_final_thresholds/final_funnel_counts.png","01_core_ic1","検出パイプラインのファネル（各段の絞り込み数）","IC1"),
 ("final_label_counts","fcc_final_thresholds/final_label_counts.png","01_core_ic1","二分岐ラベル数（gauge-recalibration / firmware-suspected / 他）","IC1"),
 ("tail_unresponded_vs_cycles","fcc_final_thresholds/tail_unresponded_opportunities_vs_cycles_final.png","01_core_ic1","末尾の無応答機会 × サイクル。FW疑い群は深くサイクルするのに無応答","IC1"),
 ("tail_opps_vs_flat_tail","fcc_final_thresholds/tail_opportunities_vs_flat_tail_final.png","01_core_ic1","末尾機会数 × 平坦尾部日数","IC1"),
 ("surrogate_decision_tree","fcc_final_thresholds/surrogate_decision_tree.png","01_core_ic1","二分岐ロジックの代理決定木（解釈可能・非ブラックボックス）","IC1"),
 ("label_transition_heatmap","fcc_final_thresholds/label_transition_baseline_to_final_heatmap.png","01_core_ic1","ベースライン→確定ラベルの遷移ヒートマップ","IC1"),
 # thresholds
 ("effective_fcc_step_sensitivity","fcc_final_thresholds/effective_fcc_step_sensitivity.png","02_thresholds","有効FCCステップ閾値の感度（≥50mWh で micro-wobble を除外）","IC1/IC2"),
 ("response_delay_cdf","fcc_final_thresholds/response_delay_cdf_24_72_168.png","02_thresholds","FCC応答遅延CDF。72hで約95%(0.9513)をカバー＝応答窓72hの根拠","IC1"),
 ("no_response_probability_by_k","fcc_final_thresholds/no_response_probability_by_k.png","02_thresholds","健全応答確率下での連続k回無応答の確率（k=2で0.013）＝計数閾値の根拠","IC1"),
 ("response_window_sensitivity","fcc_final_thresholds/response_window_sensitivity_counts.png","02_thresholds","応答窓 24/72/168h 摂動の集合安定性（Jaccard=1.0＝閾値非恣意の実証）","IC1"),
 ("flat_tail_distribution","fcc_final_thresholds/flat_tail_distribution_with_thresholds.png","02_thresholds","平坦尾部日数分布としきい値（≥180d FW / ≥120d gauge）","IC1"),
 ("tail_cycle_delta_distribution","fcc_final_thresholds/tail_cycle_delta_distribution_with_thresholds.png","02_thresholds","末尾サイクル増分分布としきい値（≥30）","IC1"),
 # IC6 gap quality
 ("large_gap_opportunity_audit","fcc_final_thresholds/large_gap_opportunity_audit.png","03_gap_quality_ic6","large-gap機会の監査。ロガー休眠/打ち切りを無応答証拠から構造的に除外","IC6"),
 ("large_gap_quality_distribution","fcc_online_v2/large_gap_quality_distribution.png","03_gap_quality_ic6","ギャップ品質ティア分布（HIGH_OK/MEDIUM_GAP/LOW_LARGE_GAP）","IC6"),
 ("gap_rule_sensitivity_counts","fcc_online_v2/gap_rule_sensitivity_counts.png","03_gap_quality_ic6","ギャップルール感度（ラベル数の頑健性）","IC6"),
 # IC5 stateful
 ("stateful_vs_stateless_counts","fcc_online_v2/stateful_vs_stateless_counts.png","04_stateful_ic5","【最堅】stateful vs stateless 検出数。永続状態で窓外証拠を回収（gain=29）","IC5"),
 ("stateful_only_evidence_examples","fcc_online_v2/stateful_only_evidence_examples.png","04_stateful_ic5","30日窓をまたぐエピソードの回収実例。episode_idキーで時刻順リプレイ","IC5"),
 # IC2 dual track
 ("any_vs_effective_state_scatter","fcc_online_v2/any_vs_effective_state_scatter.png","05_dualtrack_ic2","any-change(≥1mWh) vs effective(≥50mWh) の二系統状態","IC2"),
 ("micro_wobble_step_distribution","fcc_online_v2/micro_wobble_step_distribution.png","05_dualtrack_ic2","micro-wobble ステップ分布（軟较正 soft-calibration への分離）","IC2"),
 # IC4 normative vs personalized
 ("personalized_vs_normative_roc_pr","fcc_online_v2/personalized_vs_normative_roc_pr.png","06_normative_ic4","個別(AUC≈0.82)vs規範(AUC≈0.56)のROC/PR。規範はリーク回避の代償でnear-random","IC4"),
 ("personalized_vs_normative_calibration","fcc_online_v2/personalized_vs_normative_calibration.png","06_normative_ic4","個別/規範モデルの較正曲線","IC4"),
 ("normative_feature_importance","fcc_online_v2/normative_feature_importance.png","06_normative_ic4","規範モデルの特徴量重要度（FCC履歴を構造的に除外）","IC4"),
 # v2 results
 ("v2_label_counts","fcc_online_v2/v2_label_counts.png","07_v2_results","v2 トリアージラベル数（FW Core5/Watch43/Gauge Core4/Soft22/Review325…）","結果"),
 ("v2_policy_matrix_heatmap","fcc_online_v2/v2_policy_matrix_heatmap.png","07_v2_results","9段単一ラベル方策マトリクス","IC8"),
 ("v2_funnel_counts","fcc_online_v2/v2_funnel_counts.png","07_v2_results","v2 ファネル","結果"),
 ("v2_final_proxy_cross_tab","fcc_online_v2/v2_final_proxy_cross_tab_heatmap.png","07_v2_results","v2 ラベル × バッチ確定(proxy真値)のクロス集計","結果"),
 ("fw_topn_yield_curve","fcc_online_v2/fw_topn_yield_curve.png","07_v2_results","FW top-N 収量曲線（top50 recall=1.0）","結果"),
 ("active_false_alert_dual_basis","fcc_online_v2/active_false_alert_dual_basis.png","07_v2_results","誤警報率（any-change基準0.71 vs effective基準0）","結果"),
 ("v2_transition_v1_to_v2","fcc_online_v2/v2_transition_v1_to_v2_heatmap.png","07_v2_results","v1→v2 ラベル遷移","結果"),
 ("lead_time_by_proxy_label","fcc_online_v2/lead_time_by_proxy_label.png","07_v2_results","proxyラベル別の先行検知リードタイム","結果"),
 # IC3 HW enrichment
 ("hardware_enrichment_fw_core","fcc_online_v2/hardware_enrichment_fw_core.png","08_hw_enrichment_ic3","FW Core のHW偏在（分類後の記述的富化。判定には不使用）","IC3"),
 ("eb_enrichment_fw_check","fcc_final_thresholds/hardware_enrichment_empirical_bayes_fw_check.png","08_hw_enrichment_ic3","経験ベイズ収縮によるFRU/機種別FW疑い率（Fisher+BH-FDR）","IC3"),
 # examples (PII)
 ("example_fw_core","fcc_online_v2/example_fw_core_top20/MHAGI-PF3YCK6H_mhagi.png","09_examples_PII","FW core 実例：機会反復×FCC無応答（※ファイル名にPII）","IC1"),
 ("example_gauge_core","fcc_online_v2/example_gauge_core_top20/X1YADL_KMK.png","09_examples_PII","Gauge core 実例：適格機会が皆無（※ファイル名にPII）","IC1"),
 ("example_normal_responding","fcc_online/examples_normal_responding/AINOW_UX.png","09_examples_PII","正常応答の対照例（機会後にFCCが応答）（※ファイル名にPII）","対照"),
]

# copy + build index
by_slug = {}
idx = ["# 図版インデックス（FIGURE_INDEX）",
       "",
       "特許出願レビュー発表パッケージの厳選図版。テーマ別フォルダに格納。各図のキャプションと対応する発明的要素(IC)を示す。",
       "",
       "> ⚠️ `09_examples_PII/` の個別端末パネルはファイル名に端末名・ユーザー名（仮名化されていない生ID）を含む。**社外配布・公開前に匿名化／除外**すること。",
       "",
       "| テーマ | ファイル | 対応IC | キャプション |",
       "|---|---|---|---|"]
missing = []
for slug, src, theme, cap, ic in CUR:
    s = os.path.join(FIG, src)
    if not os.path.exists(s):
        missing.append(src); continue
    d_dir = os.path.join(OUTFIG, theme); os.makedirs(d_dir, exist_ok=True)
    ext = os.path.splitext(src)[1]
    dst = os.path.join(d_dir, slug + ext)
    shutil.copy2(s, dst)
    by_slug[slug] = dst
    idx.append(f"| {theme} | `{theme}/{slug}{ext}` | {ic} | {cap} |")
idx.append("")
idx.append(f"合計 {len(by_slug)} 図 / 欠落 {len(missing)}: {missing}")
open(os.path.join(OUTFIG, "FIGURE_INDEX.md"), "w", encoding="utf-8").write("\n".join(idx))
print("figures copied:", len(by_slug), "missing:", missing)

# ---------------- PPTX ----------------
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]
NAVY = RGBColor(0x1F,0x33,0x55); GREY = RGBColor(0x55,0x55,0x55); ACC = RGBColor(0xB0,0x1F,0x1F)

def _tb(slide, l,t,w,h):
    tb = slide.shapes.add_textbox(l,t,w,h); return tb.text_frame

def add_title_slide(title, sub, foot):
    s = prs.slides.add_slide(blank)
    tf = _tb(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(2.0)); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=title
    r.font.size=Pt(32); r.font.bold=True; r.font.name=JP; r.font.color.rgb=NAVY
    p2=tf.add_paragraph(); r=p2.add_run(); r.text=sub; r.font.size=Pt(18); r.font.name=JP; r.font.color.rgb=GREY
    tf2=_tb(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
    r=tf2.paragraphs[0].add_run(); r.text=foot; r.font.size=Pt(12); r.font.name=JP; r.font.color.rgb=GREY
    return s

def add_header(s, title):
    tf=_tb(s, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9)); tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run(); r.text=title
    r.font.size=Pt(24); r.font.bold=True; r.font.name=JP; r.font.color.rgb=NAVY

def add_image_slide(title, slug, caption):
    s=prs.slides.add_slide(blank); add_header(s,title)
    if slug not in by_slug:
        tf=_tb(s,Inches(1),Inches(3),Inches(11),Inches(1)); tf.paragraphs[0].add_run().text=f"[missing figure: {slug}]"
        return s
    path=by_slug[slug]
    iw,ih=Image.open(path).size
    win,hin=iw/96.0, ih/96.0
    maxw,maxh=12.0, 5.3
    sc=min(maxw/win, maxh/hin)
    w=Inches(win*sc); h=Inches(hin*sc)
    l=int((SW-w)/2); t=Inches(1.25)
    s.shapes.add_picture(path,l,t,width=w,height=h)
    tf=_tb(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7)); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=caption; r.font.size=Pt(12); r.font.name=JP; r.font.color.rgb=GREY
    return s

def add_two_image_slide(title, slugA, slugB, caption):
    s=prs.slides.add_slide(blank); add_header(s,title)
    cols=[(slugA, Inches(0.4)), (slugB, Inches(6.95))]
    for slug,left in cols:
        if slug not in by_slug: continue
        path=by_slug[slug]; iw,ih=Image.open(path).size; win,hin=iw/96.0,ih/96.0
        maxw,maxh=5.9,5.0; sc=min(maxw/win,maxh/hin)
        w=Inches(win*sc); h=Inches(hin*sc)
        s.shapes.add_picture(path,left,Inches(1.3),width=w,height=h)
    tf=_tb(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7)); tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=caption; r.font.size=Pt(12); r.font.name=JP; r.font.color.rgb=GREY
    return s

def add_bullets_slide(title, bullets):
    s=prs.slides.add_slide(blank); add_header(s,title)
    tf=_tb(s, Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.6)); tf.word_wrap=True
    first=True
    for lvl,txt,bold in bullets:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level=lvl; r=p.add_run(); r.text=txt
        r.font.size=Pt(20 if lvl==0 else 16); r.font.bold=bold; r.font.name=JP
        r.font.color.rgb=NAVY if bold and lvl==0 else RGBColor(0x22,0x22,0x22)
        p.space_after=Pt(6)
    return s

def add_table_slide(title, headers, rows, note=None):
    s=prs.slides.add_slide(blank); add_header(s,title)
    nr=len(rows)+1; nc=len(headers)
    gt=s.shapes.add_table(nr,nc,Inches(0.5),Inches(1.3),Inches(12.3),Inches(0.4*nr)).table
    for j,hh in enumerate(headers):
        c=gt.cell(0,j); c.text=hh
        for para in c.text_frame.paragraphs:
            for r in para.runs: r.font.size=Pt(13); r.font.bold=True; r.font.name=JP; r.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
        c.fill.solid(); c.fill.fore_color.rgb=NAVY
    for i,row in enumerate(rows,1):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.text=val
            for para in c.text_frame.paragraphs:
                for r in para.runs: r.font.size=Pt(12); r.font.name=JP
    if note:
        tf=_tb(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5)); tf.word_wrap=True
        r=tf.paragraphs[0].add_run(); r.text=note; r.font.size=Pt(11); r.font.name=JP; r.font.color.rgb=GREY
    return s

# ---- build deck ----
add_title_slide("バッテリ燃料計 SoH凍結／FCC無応答の検出・原因切り分けアルゴリズム",
    "特許出願レビュー — 特許性評価  （整理番号(仮) BUA-FCC-FREEZE-2026）",
    "2026-06-18 ／ 技術的特許性評価（法的助言ではありません・先行特許番号は未検証・出願前に弁理士レビュー必須）")
add_image_slide("S2. 課題 ─ SoH凍結は静的に見えない","soh_overlay_by_class",
    "SoH=FCC×100/DesignCapacity。SoHはFCCのステップ更新時のみ動く。凍結は健全な浅充放電とも外形が同じ＝静的に判別不能。")
add_image_slide("S3. 凍結は使用挙動から予測できない（=HW/FW起因）","very_stale_xgb_shap",
    "公平領域AUC≈0.54。min_rsoc は『深放電ほど凍結』の反usage方向 → 行動で説明不能。これが『予測』から『機会条件付き無応答の監査』への転換動機。")
add_image_slide("S4. 発明の核アイデア（機会 × 応答 の直積）","opportunity_vs_response",
    "機会=RSOCの high→low→high 遠足(80→20→80)。応答=機会後72h窓内のFCC有効ステップ(≥50mWh)。『機会あり×応答なし』を検出。censored/unknown/休眠は無応答に算入しない。")
add_image_slide("S5. 検出パイプライン（ファネル）","final_funnel_counts",
    "全観測→平坦尾部→機会検出→応答判定→品質ゲート→二分岐。各段の絞り込み。")
add_image_slide("S6. 二分岐：gauge-recalibration vs firmware-suspected","final_label_counts",
    "機会反復×無応答→FW疑い。機会皆無→再較正要。機種名・ベンダ名は判定に一切使わない（モデル非依存）。")
add_image_slide("S7. FW証拠：無応答機会 × サイクル","tail_unresponded_vs_cycles",
    "FW疑い群は active より多くサイクル・深放電するのにFCC無応答。『使えば直る』前提が崩れた群。")
add_image_slide("S8. 閾値正当化① 有効ステップ ≥50mWh","effective_fcc_step_sensitivity",
    "1mWh(any-change)は量子化ノイズ(micro-wobble)を拾う。50mWhを『有効再学習』と定義。感度カーブで妥当性。")
add_image_slide("S9. 閾値正当化② 応答窓 72h","response_delay_cdf",
    "FCC応答遅延CDFは72hで約95%(0.9513)をカバー → 応答窓=72hの根拠。24/72/168h併記。")
add_image_slide("S10. 閾値正当化③ 無応答計数しきい k","no_response_probability_by_k",
    "健全応答確率下で連続k回無応答の確率（k=2で0.013）→ 計数閾値の統計的根拠。")
add_image_slide("S11. ロバスト性：応答窓感度（Jaccard=1.0）","response_window_sensitivity",
    "24/72/168h 摂動でFW/GAUGE集合のJaccard=1.0 → 閾値が恣意的でない実証（進歩性補強・反『ルーチン最適化』）。")
add_image_slide("S12. 右打ち切り／ギャップ品質ゲート","large_gap_opportunity_audit",
    "ロガー休眠(large-gap)・観測打ち切りを無応答証拠から構造的に除外 → active→actionable 誤分類0件。large-gap=機会なしと誤結論しない安全策。")
add_image_slide("S13. 解釈可能な判定ルール（代理決定木）","surrogate_decision_tree",
    "二分岐ロジックを可視化。非ブラックボックス＝説明性・適格性に有利。")
add_image_slide("S14. IC5 状態永続化：窓外証拠の回収","stateful_vs_stateless_counts",
    "オンライン制約=直近30日のみ可視。stateless は窓先頭で開始highが窓前の機会を取りこぼす。永続状態で回収→stateful-only gain=29。")
add_image_slide("S15. IC5 実例：30日窓をまたぐ証拠回収【最も堅い要素】","stateful_only_evidence_examples",
    "episode_idキーで解決済イベントを時刻順リプレイ、complete<reset<deadline 順序意味論で確定、物理エピソード一度限り計数。審査官も『設計事項で潰しにくい』と評価。")
add_image_slide("S16. IC2 デュアルトラック","any_vs_effective_state_scatter",
    "any-change(≥1mWh)とeffective(≥50mWh)を並列追跡。micro-wobbleのみ=soft-calibrationに分離し誤検出抑制。")
add_image_slide("S17. IC4 規範 vs 個別モデル（リーク回避の正直な開示）","personalized_vs_normative_roc_pr",
    "個別AUC≈0.82だが自分の無応答を『劣化なら当然』と学習＝リーク。規範(FCC履歴全除外)AUC≈0.56でnear-random。→中核は決定論カウンタで構成、MLは独立クレームから除外。")
add_two_image_slide("S18. v2 結果：トリアージ階層（9段単一ラベル）","v2_label_counts","v2_policy_matrix_heatmap",
    "FW Core5/Watch43/engineering top50/Gauge Core4/Soft22/Review325…。高信頼の確定アクション=FW5+Gauge4=9台に厳格化。")
add_two_image_slide("S19. v2 精度：proxy照合 と top-N 収量","v2_final_proxy_cross_tab","fw_topn_yield_curve",
    "バッチ確定版(fcc_final)をproxy真値に：FW Core precision1.0 / top50 recall1.0 / Gauge Core precision1.0。")
add_image_slide("S20. 誤警報ゼロ（effective基準）","active_false_alert_dual_basis",
    "any-change基準では0.71に見えるが effective基準では0。差はmicro-wobbleの定義差で誤判定ではない。")
add_image_slide("S21. HW富化（記述的・分類後／判定には不使用）","hardware_enrichment_fw_core",
    "分類確定後にのみFRU/機種偏在を経験ベイズ(Beta事前・Fisher・BH-FDR)で集計。判定には逆流させない（モデル非依存の徹底）。")
add_two_image_slide("S22. 実例パネル（※ファイル名にPII）","example_fw_core","example_gauge_core",
    "左=FW core（機会反復×無応答）／右=Gauge core（適格機会が皆無）。視覚的な『動かぬ証拠』。社外配布前に匿名化。")
add_table_slide("S23. 特許性マトリクス",
    ["発明的要素(IC)","新規性","進歩性","適格性","総合"],
    [["IC1 機会条件付き無応答＋右打ち切り＋二分岐","中","中","中","中（最有望）"],
     ["IC5 状態永続化＋イベント順序意味論","中","中","中","中（最も堅い）"],
     ["IC6 ギャップ品質ゲート","低〜中","低","低","低〜中"],
     ["IC4 規範/個別ツイン・リーク回避","低〜中","低","弱","低〜中"],
     ["IC2 デュアルトラック","低","低","中","低"],
     ["IC3 モデル非依存分類＋EB富化","低","低","弱","低"],
     ["IC8 トリアージラダー＋アラート制御","低","低","中","低"],
     ["IC7 Poisson-binomial異常","低","低","最弱","最低（出願不可）"]],
    note="主軸=IC1+IC5。IC4/IC6は補強従属。IC2/IC3/IC7/IC8は防御的公開／営業秘密へ。IC7は純数学で適格性不可。")
add_bullets_slide("S24. 推奨クレーム要旨", [
    (0,"独立クレーム2件（中核に集中）",True),
    (1,"Claim 1（方法）= IC1：機会条件付きFCC無応答＋右打ち切り除外＋二分岐 ＋ ラベル依存の具体的物理介入（深放電プロンプト／FWエスカレーション）＋ 介入→FCC回復のクローズドループ検証",False),
    (1,"Claim 2（システム）= IC5：30日窓制約下の状態永続化による窓外証拠回収＋complete<reset<deadline 順序意味論＋episode一度限り計数",False),
    (0,"主要従属クレーム",True),
    (1,"閾値詳細(≥180d/≥30cyc/≥50mWh/72h) ／ ギャップ品質ゲート ／ デュアルトラック ／ 規範モデルを『基準線検証』に機能再定義 ／ アラート・クールダウン(30日) ／ 経験ベイズHW富化(分類後descriptive)",False),
    (0,"Poisson-binomial（IC7）は独立から全面削除（純数学で適格性不可）",True),
])
add_bullets_slide("S25. リスク・未解決事項", [
    (0,"規範モデルAUC≈0.56（ML空洞化・的中hit）",True),
    (1,"異常スコアと無応答カウントの相関0.993＝実質同一物 → 中核を決定論カウンタで構成しMLを独立から除外して回避",False),
    (0,"BIOS/FWバージョン・介入結果データの欠如（最大の実証ギャップ）",True),
    (1,"『介入→FCC回復』のクローズドループ証拠が未取得 → 審査官は出願『時期尚早』と認定",False),
    (0,"その他",True),
    (1,"proxyラベル依存（真のFW不具合確定ラベルでの再検証が望ましい）／ クレーム抽象性（保守アクションの具体化で対応済）",False),
])
add_bullets_slide("S26. 出願戦略・次アクション ／ Disclaimer", [
    (0,"法域：JP先願 → PCT → US/EP（US/EPはBMS・制御・クローズドループ強紐付け版で適格性確保）",True),
    (0,"補強の最優先：介入→FCC回復のクローズドループ実証データ取得",True),
    (0,"出願前必須：主要先行特許の独立クレーム精読／既存の社内・外部公開有無の確認（新規性喪失の時限）",True),
    (1,"特許化=IC1・IC5（＋IC4/IC6補強）／ 営業秘密=具体的微調整値 ／ 防御的公開=IC2/IC3/IC7/IC8",False),
    (0,"Disclaimer：本資料は技術的特許性評価であり法的助言ではありません。先行特許番号は未検証。出願前に登録弁理士の正式レビューを受けてください。",True),
])

out=os.path.join(BUNDLE,"03_presentation","patent_review_slides.pptx")
prs.save(out)
print("pptx slides:", len(prs.slides._sldIdLst), "->", out)
