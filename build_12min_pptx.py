#!/usr/bin/env python
"""FCC更新停止（ゲージ凍結）判定技術 — 社内発表 12分版 deck (.pptx).

fcc_patent_summary_slides_v5.md（44枚版）を 12分口頭発表用に 13枚へ凝縮した版。
話者ノート（話す内容）を各スライドのノートペインに埋込。
テキスト版・ノート原稿: data/reports/fcc_patent_12min_presentation.md

数値は v5 と同一（_v4_results_summary.json / online_latest_snapshot_v2.csv /
fcc_final_action_labels.csv で一次データ照合済みの値のみ）。
候補③の7要素は minimal_state_ablation.csv の必要構成
(fsm / pending / seen_ids / last_eff_ts / eff_cycle / gap_censor / ordering) と一致。
NOT a legal opinion / 先行技術 UNVERIFIED / 介入・FWバージョン NOT AVAILABLE / 捏造なし。
"""
from __future__ import annotations

import sys
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
REPORTS = ROOT / "data" / "reports"
FIG = REPORTS / "figures" / "fcc_patent_evidence_v4"
OUT = REPORTS / "fcc_patent_12min_presentation.pptx"

# ---- palette (build_v5_pptx.py と同一) ------------------------------------ #
NAVY = RGBColor(0x1F, 0x37, 0x5B); BLUE = RGBColor(0x1F, 0x77, 0xB4)
STEEL = RGBColor(0x33, 0x66, 0x99); GREY = RGBColor(0x55, 0x55, 0x55)
RED = RGBColor(0xC0, 0x39, 0x2B); GREEN = RGBColor(0x2C, 0xA0, 0x2C)
DGREEN = RGBColor(0x37, 0x6B, 0x3A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7); ORANGE = RGBColor(0xB0, 0x5A, 0x00)
TEAL = RGBColor(0x2E, 0x6E, 0x6A); PURPLE = RGBColor(0x6A, 0x3D, 0x8A)
INK = RGBColor(0x22, 0x22, 0x22)
YELLOW_BG = RGBColor(0xFF, 0xF6, 0xDA); YELLOW_LN = RGBColor(0xD9, 0xA4, 0x2A)
LGREEN_BG = RGBColor(0xE4, 0xF0, 0xE4); LBLUE_BG = RGBColor(0xE2, 0xEC, 0xF7)
LRED_BG = RGBColor(0xF7, 0xE4, 0xE2); LGREY_BG = RGBColor(0xEC, 0xEC, 0xEC)

FONT = "Meiryo UI"
DISC = ("技術的特許性エビデンス（NOT a legal opinion）。先行技術は UNVERIFIED。"
        "特許化の粒度・範囲の最終判断は社内reviewer/弁理士。捏造（地上真実/介入/FW/因果結論）なし。")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def EM(v):
    return Emu(int(round(v)))


def _cjk(run):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", FONT)


def _style(run, size=14, bold=False, color=None):
    run.font.size = Pt(size); run.font.bold = bold; run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    _cjk(run)


def _tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(EM(left), EM(top), EM(width), EM(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size=14, bold=False, color=None, align=None):
    p.text = text
    for run in p.runs:
        _style(run, size, bold, color)
    if align is not None:
        p.alignment = align


def footer(slide):
    tf = _tb(slide, Inches(0.3), SH - Inches(0.36), SW - Inches(0.6), Inches(0.32))
    _set(tf.paragraphs[0], DISC, size=8, color=GREY)


def band(slide, color, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EM(SW), EM(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def rect(slide, left, top, width, height, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, EM(left), EM(top), EM(width), EM(height))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1.2)
    shp.shadow.inherit = False
    return shp


def hint(slide, text):
    """非エンジニア向け 💡 一言要約の黄色帯。"""
    top = SH - Inches(1.02)
    shp = rect(slide, Inches(0.3), top, SW - Inches(0.6), Inches(0.56),
               YELLOW_BG, YELLOW_LN)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], "💡 " + text, size=12.5, bold=True, color=INK)


def content_slide(title, tag=None, accent=BLUE):
    s = prs.slides.add_slide(BLANK)
    band(s, accent, Inches(0.86))
    tf = _tb(s, Inches(0.5), Inches(0.12), SW - Inches(2.6), Inches(0.64))
    _set(tf.paragraphs[0], title, size=21, bold=True, color=WHITE)
    if tag:
        tf2 = _tb(s, SW - Inches(2.15), Inches(0.2), Inches(1.9), Inches(0.5))
        _set(tf2.paragraphs[0], tag, size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    footer(s)
    return s


def bullets(slide, items, left, top, width, height, size=13, gap=6):
    tf = _tb(slide, left, top, width, height)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, kw = item
        else:
            text, kw = item, {}
        _set(p, text, size=kw.get("size", size), bold=kw.get("bold", False),
             color=kw.get("color", INK))
        p.space_after = Pt(kw.get("gap", gap))
        if kw.get("level"):
            p.level = kw["level"]
    return tf


def table(slide, data, left, top, width, height, col_w=None, font=12,
          header_fill=NAVY, align_right_cols=()):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, EM(left), EM(top), EM(width), EM(height))
    t = gt.table
    t.first_row = False; t.horz_banding = False
    if col_w:
        for j, w in enumerate(col_w):
            t.columns[j].width = EM(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = t.cell(i, j)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                color, bold = WHITE, True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
                color, bold = INK, False
            txt = str(val)
            if txt.startswith("**") and txt.endswith("**"):
                txt, bold = txt[2:-2], True
            p = cell.text_frame.paragraphs[0]
            _set(p, txt, size=font, bold=bold, color=color)
            if j in align_right_cols and i > 0:
                p.alignment = PP_ALIGN.RIGHT
    return t


def notes(slide, budget, text):
    ns = slide.notes_slide
    ns.notes_text_frame.text = f"[目安 {budget}]\n{text}"


def line(slide, x1, y1, x2, y2, color, width=2.2, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, EM(x1), EM(y1), EM(x2), EM(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    if dash is not None:
        cn.line.dash_style = dash
    cn.shadow.inherit = False
    return cn


def label(slide, x, y, w, h, text, size=11, bold=True, color=INK, align=PP_ALIGN.CENTER):
    tf = _tb(slide, x, y, w, h)
    _set(tf.paragraphs[0], text, size=size, bold=bold, color=color, align=align)
    return tf


# ========================================================================== #
# Slide 1 — タイトル
# ========================================================================== #
s = prs.slides.add_slide(BLANK)
band(s, NAVY, Inches(2.6))
tf = _tb(s, Inches(0.7), Inches(0.55), SW - Inches(1.4), Inches(1.9))
_set(tf.paragraphs[0], "テレメトリからの FCC更新停止（ゲージ凍結）判定技術",
     size=30, bold=True, color=WHITE)
p = tf.add_paragraph()
_set(p, "— と、特許候補3件（+将来2件） —  社内発表 12分版", size=19, bold=True,
     color=RGBColor(0xCF, 0xDE, 0xF2))
bullets(s, [
    "本題: テレメトリ（RSOC / FCC / サイクル / 時刻）だけで『FCCを再学習しなくなった個体』を判定する",
    "母集団: 実バッテリ履歴 752ユーザー / 3,130,394 サンプル / 24,711 学習機会で検証",
    "本資料は 44枚のフル資料（v5）の要約。技術エビデンスの報告であり、特許化の判断は本日の審議事項",
], Inches(0.8), Inches(3.1), SW - Inches(1.6), Inches(2.2), size=15, gap=10)
footer(s)
notes(s, "0:35",
      "本日は、ノートPCバッテリのテレメトリだけから『燃料計の学習が止まった個体』、"
      "いわゆるゲージ凍結を判定する技術と、その開発過程で生まれた特許候補3件をご報告します。"
      "実フリート752台・約313万サンプルで検証済みです。本資料は法的判断ではなく技術エビデンスの"
      "報告で、特許化の判断は皆さまにお諮りします。12分お付き合いください。")

# ========================================================================== #
# Slide 2 — 全体サマリ
# ========================================================================== #
s = content_slide("全体サマリ — この1枚で結論", tag="2/13")
rows = [
    ("課題", "SoH表示を駆動するFCCが更新されない『ゲージ凍結』は静的検査で判別不能。"
             "752台中114台が凍結、うち70台は使用で説明不能", LRED_BG),
    ("解決", "『凍結を予測』は失敗(AUC≈0.54=コイン投げ)→『再学習の機会の直後72時間に"
             "応答したか』を機械的に監査する方式へ転換", LBLUE_BG),
    ("実績", "全履歴監査でゲージ再較正18台 / FW確認14台に自動振り分け・active層からの"
             "誤エスカレーション0件。オンライン版は Core判定 精度1.0（proxy=推定ラベル基準）・能動誤報0件", LGREEN_BG),
    ("依頼", "(1) 出願の粒度・範囲  (2) 候補②の着想日・公開有無  "
             "(3) 正式FTO/特許性調査の発注  (4) 介入データ収集開始", LGREY_BG),
]
y = Inches(1.08)
for head, body, bg in rows:
    hshp = rect(s, Inches(0.45), y, Inches(1.25), Inches(1.02), NAVY)
    htf = hshp.text_frame; htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(htf.paragraphs[0], head, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bshp = rect(s, Inches(1.78), y, SW - Inches(2.3), Inches(1.02), bg)
    btf = bshp.text_frame; btf.word_wrap = True
    btf.margin_left = Inches(0.12); btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(btf.paragraphs[0], body, size=13.5, color=INK, align=PP_ALIGN.LEFT)
    y += Inches(1.12)
hint(s, "『動くべき機会に動かなかった』事実だけを、欠測に騙されず数える監査技術。"
        "エビデンスは強いが、新規性の判断はここから先が本番。")
notes(s, "0:55",
      "まず全体像です。課題は、バッテリの健康度表示の元になるFCCが更新されない『ゲージ凍結』を、"
      "従来の静的検査では正常か要再較正かファームウェア起因か区別できないことです。実フリートでは"
      "752台中114台が凍結し、うち70台は使い方では説明できませんでした。当初はAIで凍結を予測しよう"
      "として、コイン投げ並みの精度で失敗しました。そこで発想を変え、『再学習の機会が来た直後に"
      "実際に応答したか』を機械的に監査する方式にしたところ、ゲージ再較正対象18台・FW確認対象14台"
      "まで自動で絞り込め、正常稼働層からの誤検知はゼロでした。ここから特許候補が3件生まれました。"
      "お願いは最後にまとめる4点です。")

# ========================================================================== #
# Slide 3 — 基礎: 燃料計とFCC学習
# ========================================================================== #
s = content_slide("基礎: バッテリの『燃料計』と FCC 学習", tag="3/13", accent=TEAL)
table(s, [
    ("用語", "意味", "身近なたとえ"),
    ("FCC（満充電容量）", "『いま満タンでどれだけ入るか』の学習値（mWh=ミリワット時）。劣化で減る",
     "計り直して覚える『いまの満タン量』"),
    ("RSOC（相対残量）", "いまの残量% = 残容量 ÷ FCC × 100", "ガソリンメーターの針"),
    ("SoH（健全性）", "健康度% = FCC × 100 ÷ 設計容量", "新品比の容量%"),
    ("FCC学習（再較正）", "『満充電→深く放電→再び満充電』のときだけ計り直せる", "体重計のゼロ点合わせ"),
    ("ゲージ凍結", "FCCが長期間更新されない。SoH表示も止まる", "校正が止まり表示が古いまま"),
], Inches(0.45), Inches(1.1), SW - Inches(0.9), Inches(3.6),
      col_w=[Inches(2.6), Inches(5.7), Inches(4.1)], font=13)
bullets(s, [
    "実容量は直接測れない。ゲージICが充放電から推定する学習値であり、学習の機会がないと更新されない",
    "SoH・残り時間・保証判断は、すべてこの学習値に依存する",
], Inches(0.5), Inches(4.95), SW - Inches(1.0), Inches(1.0), size=13)
hint(s, "電池の残量・健康度表示は『実測』ではなく燃料計チップの学習値。学習が止まると表示だけが古いまま取り残される。")
notes(s, "0:45",
      "前提を1枚だけ。PCの電池残量や健康度は実測値ではなく、電池内の『燃料計チップ』が学習した"
      "満充電容量FCCから計算されます。そしてこのFCCは、満充電から深く放電して再び満充電に戻る、"
      "いわば体重計のゼロ点合わせのような機会にしか計り直せません。この学習が止まるのがゲージ凍結で、"
      "SoHも残り時間も保証判断も、古い値のまま固定されます。難しいのは、凍結イコール故障ではない"
      "ことです。浅い使い方で機会が無いだけかもしれません。")

# ========================================================================== #
# Slide 4 — 課題: 実フリートの実態
# ========================================================================== #
s = content_slide("課題: 実フリート752台の実態", tag="4/13", accent=RED)
table(s, [
    ("区分", "台数", "内訳・特徴"),
    ("凍結（全体）", "**114台 (15.2%)**", "stale 59 + very_stale 55"),
    ("うち使用で説明可", "44台", "常時AC 22 / 低サイクル 16 / 浅放電 6 — ユーザー起因"),
    ("**説明不能（FW/HW疑い）**", "**70台**",
     "activeより多サイクル（96.9 vs 65.7 cyc/yr）・残量1%までの深放電もあるのに凍結"),
], Inches(0.45), Inches(1.15), SW - Inches(0.9), Inches(2.3),
      col_w=[Inches(3.3), Inches(1.9), Inches(7.2)], font=13, align_right_cols=(1,))
bullets(s, [
    "※『疑い』であり確定診断ではない（FWバージョンデータは NOT AVAILABLE）",
    "放置リスク: 劣化見逃し→突然のシャットダウン / 健全電池の誤交換→無駄コスト / 保証・リース判断の歪み",
    ("　（定性評価。金額効果は未算定）", {"color": GREY, "size": 12}),
], Inches(0.5), Inches(3.75), SW - Inches(1.0), Inches(1.7), size=13)
hint(s, "原因はユーザー起因〜FW疑いまで様々で対処も違う。見分ける仕組みがないと保守が空回りする。")
notes(s, "0:50",
      "実態です。752台のうち114台、15.2%が凍結していました。このうち44台は常時AC接続や低サイクル"
      "など、使い方で説明がつきます。問題は残る70台です。この70台は、正常に学習している群より年間"
      "サイクル数が多く、残量1%まで使い切る深い放電もしているのに、FCCが動いていません。つまり機会は"
      "あるのに応答していない、ファームウェアやハードウェア起因が疑われる群です。ただしFWバージョンの"
      "データが無いため、あくまで『疑い』です。放置すれば、劣化の見逃しや健全な電池の誤交換に"
      "つながります。")

# ========================================================================== #
# Slide 5 — 失敗 → 発想の転換
# ========================================================================== #
s = content_slide("失敗した素朴アプローチ → 発想の転換", tag="5/13", accent=ORANGE)
lshp = rect(s, Inches(0.45), Inches(1.15), Inches(6.0), Inches(3.3), LRED_BG, RED)
ltf = lshp.text_frame; ltf.word_wrap = True
ltf.margin_left = Inches(0.16); ltf.margin_top = Inches(0.12); ltf.vertical_anchor = MSO_ANCHOR.TOP
_set(ltf.paragraphs[0], "✗ 『凍結を予測する』（教師ありML）", size=15, bold=True, color=RED, align=PP_ALIGN.LEFT)
for txt in [
    "33特徴量・公平比較領域: AUC 0.535 / 0.540 ≈ コイン投げ",
    "最重要特徴 min_rsoc は『深放電ほど凍結』という反usage方向（交絡）",
    "『壊れそうか』の予言は当たらない",
]:
    p = ltf.add_paragraph(); _set(p, "・" + txt, size=13, color=INK); p.space_before = Pt(6)
rshp = rect(s, Inches(6.9), Inches(1.15), Inches(6.0), Inches(3.3), LGREEN_BG, DGREEN)
rtf = rshp.text_frame; rtf.word_wrap = True
rtf.margin_left = Inches(0.16); rtf.margin_top = Inches(0.12); rtf.vertical_anchor = MSO_ANCHOR.TOP
_set(rtf.paragraphs[0], "○ 『機会への応答を監査する』（決定論カウンタ）", size=15, bold=True, color=DGREEN, align=PP_ALIGN.LEFT)
for txt in [
    "『再学習の機会』が来たときに実際に応答したかを機械的に数える",
    "MLではない: 説明可能・再現可能・機種非依存",
    "この転換自体が本技術の核（特許候補の出発点）",
]:
    p = rtf.add_paragraph(); _set(p, "・" + txt, size=13, color=INK); p.space_before = Pt(6)
rect(s, Inches(6.42), Inches(2.5), Inches(0.52), Inches(0.6), NAVY,
     shape=MSO_SHAPE.RIGHT_ARROW)
hint(s, "壊れそうかをAIに予言させるのは失敗（コイン投げ並み）。『チャンスに応えたか』を出席簿のように数える方式に切り替えたら解けた。")
notes(s, "0:45",
      "最初のアプローチは教師あり機械学習でした。33特徴量で凍結を予測させたところ、AUCは0.54前後、"
      "ほぼコイン投げで失敗です。しかも最重要特徴量は『深く放電するユーザーほど凍結する』という、"
      "直感と逆向きの交絡を拾っていました。ここで発想を転換しました。壊れそうかを予言させるのでは"
      "なく、『再学習のチャンスが来たときに応えたか』を出席簿のように数える。機械学習ではなく決定論的"
      "なカウンタなので、説明可能で再現可能です。この転換が本技術の核です。")

# ========================================================================== #
# Slide 6 — 手法: END起点72h監査（図解）
# ========================================================================== #
s = content_slide("手法: 学習機会の抽出と『END起点72h監査』", tag="6/13")
# --- RSOC waveform schematic --------------------------------------------- #
x0, x1, x2, x3, x4, x5 = (Inches(0.9), Inches(2.3), Inches(4.6), Inches(5.6),
                          Inches(7.6), Inches(9.2))
y_hi, y_lo = Inches(1.55), Inches(3.05)
label(s, Inches(0.28), y_hi - Inches(0.14), Inches(0.6), Inches(0.3), "80%",
      size=10, color=GREY, align=PP_ALIGN.RIGHT)
label(s, Inches(0.28), y_lo - Inches(0.14), Inches(0.6), Inches(0.3), "20%",
      size=10, color=GREY, align=PP_ALIGN.RIGHT)
# audit window (draw first so the waveform stays visible on top)
win = rect(s, x4, Inches(1.25), Inches(1.6), Inches(2.1), LBLUE_BG, STEEL)
wtf = win.text_frame; wtf.vertical_anchor = MSO_ANCHOR.BOTTOM
_set(wtf.paragraphs[0], "監査窓 72h", size=12, bold=True, color=STEEL, align=PP_ALIGN.CENTER)
for xa, ya, xb, yb in [(x0, y_hi, x1, y_hi), (x1, y_hi, x2, y_lo),
                       (x2, y_lo, x3, y_lo), (x3, y_lo, x4, y_hi),
                       (x4, y_hi, x5, y_hi)]:
    line(s, xa, ya, xb, yb, BLUE, 2.6)
for xx, yy, txt, col in [(x1, y_hi, "START", GREY), (x2, y_lo, "LOW", GREY),
                         (x4, y_hi, "END", RED)]:
    rect(s, xx - Inches(0.07), yy - Inches(0.07), Inches(0.14), Inches(0.14),
         col, shape=MSO_SHAPE.OVAL)
    label(s, xx - Inches(0.55), yy + (Inches(0.12) if yy == y_lo else -Inches(0.44)),
          Inches(1.1), Inches(0.3), txt, size=11, color=col)
label(s, Inches(0.9), Inches(3.35), Inches(8.0), Inches(0.35),
      "学習機会 = RSOC 高(≈80%) → 低(≈20%) → 高(≈80%) の往復1サイクル",
      size=11.5, bold=False, color=GREY, align=PP_ALIGN.LEFT)
for i, (head, body, bg, ln) in enumerate([
        ("responded（応答あり）", "≥50mWhの有効ステップを観測", LGREEN_BG, DGREEN),
        ("no_response（無応答）", "窓を全部観測して変化なし", LRED_BG, RED),
        ("censored（保留）", "観測が途中終了 → 無応答に数えない", LGREY_BG, GREY)]):
    chip = rect(s, Inches(9.55), Inches(1.15 + i * 0.78), Inches(3.45), Inches(0.68),
                bg, ln, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ctf = chip.text_frame; ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(ctf.paragraphs[0], head, size=12, bold=True, color=ln, align=PP_ALIGN.CENTER)
    p = ctf.add_paragraph(); _set(p, body, size=10, color=INK, align=PP_ALIGN.CENTER)
bullets(s, [
    "起点をSTARTにすると放電途中の変化まで誤計上（56%汚染）。ENDなら構造的に0%",
    "保留（censored）を無応答に数えないことで、データの穴を故障の冤罪にしない",
    "二分岐トリアージ: 機会反復×無応答 → FW疑い ／ 機会皆無 → ゲージ再較正案内",
], Inches(0.5), Inches(4.25), SW - Inches(1.0), Inches(1.9), size=14, gap=9)
hint(s, "採点は試験が終わって（END）から72時間だけ。途中退室（censored=保留）は不合格に数えない。")
notes(s, "1:00",
      "RSOCが高い状態から低い状態を経て再び高い状態へ戻る往復を『学習機会』として抽出します。"
      "機会の終わりENDから72時間の監査窓を置き、容量全体から見ればごく小さな50ミリワット時以上の"
      "FCC変化があれば『応答あり』、窓を最後まで観測して変化が無ければ『無応答』、観測が途中で"
      "切れたら『保留』として無応答には数えません。データの穴を故障の冤罪にしない設計です。"
      "起点をENDに置くのも本質で、START起点では放電途中の変化を誤って数える汚染が56%起きるのに"
      "対し、ENDでは構造的にゼロです。機会が繰り返しあるのに無応答ならFW疑い、機会が無いなら"
      "再較正案内、という二分岐で振り分けます。")

# ========================================================================== #
# Slide 7 — 実フリート結果
# ========================================================================== #
s = content_slide("実フリート752台での判定結果", tag="7/13", accent=DGREEN)
label(s, Inches(0.45), Inches(1.0), Inches(6.5), Inches(0.35),
      "A. 全履歴監査（v2.0-final）", size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
table(s, [
    ("最終ラベル", "件数", "割合"),
    ("正常応答（NORMAL_OR_RESPONDING）", "327", "43.5%"),
    ("データ不足・保留（REVIEW_INSUFFICIENT_DATA）", "338", "44.9%"),
    ("経過観察（WATCH）", "55", "7.3%"),
    ("**ゲージ再較正対象（ACTIONABLE_GAUGE_RESET）**", "**18**", "2.4%"),
    ("**FW確認対象（ACTIONABLE_FW_CHECK）**", "**14**", "1.9%"),
], Inches(0.45), Inches(1.4), Inches(6.35), Inches(3.2),
      col_w=[Inches(4.25), Inches(1.0), Inches(1.1)], font=11, align_right_cols=(1, 2))
bullets(s, [
    "候補96台 → 二分岐でゲージ再較正18 / FW確認14 に振り分け",
    ("active層からの誤エスカレーション 0件", {"bold": True}),
], Inches(0.5), Inches(4.75), Inches(6.3), Inches(1.0), size=12.5, gap=4)
rshp = rect(s, Inches(7.1), Inches(1.4), Inches(5.85), Inches(4.35), LBLUE_BG, STEEL)
rtf = rshp.text_frame; rtf.word_wrap = True
rtf.margin_left = Inches(0.16); rtf.margin_top = Inches(0.12); rtf.vertical_anchor = MSO_ANCHOR.TOP
_set(rtf.paragraphs[0], "B. 30日オンライン運用版 v2", size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
for txt, bold in [
    ("Core判定 precision 1.0（proxy基準 = 本番システムの推定ラベルを仮の正解とした基準）", True),
    ("能動誤報 0件（3種の活動基準すべて）", True),
    ("FWエンジニアリング上位50件で proxy FW 14台を全カバー（recall 1.0）", False),
    ("直近30日の生データしか持たなくても、全期間監査と同じ結論に到達（→ 特許候補③）", False),
]:
    p = rtf.add_paragraph(); _set(p, "・" + txt, size=13, bold=bold, color=INK)
    p.space_before = Pt(10)
hint(s, "約半数は『データ不足で保留』と正直に言う設計 — 誤指示より保留が安全。")
notes(s, "0:55",
      "結果です。全履歴監査では752台を、正常応答327・データ不足の保留338・経過観察55と仕分けし、"
      "最終的にゲージ再較正対象18台、FW確認対象14台まで自動で絞り込みました。正常稼働層からの"
      "誤エスカレーションはゼロです。約半数を『データ不足で保留』と正直に言う設計にしているのは、"
      "保守現場では誤った指示より保留のほうが安全だからです。さらに直近30日のデータだけで動く"
      "オンライン版でも、コア判定の精度は、本番システムの推定ラベルを仮の正解としたproxy基準で1.0、"
      "アクティブなユーザーへの誤報はゼロ、FW調査候補の上位50件でproxyのFW疑い14台を全てカバー"
      "しています。")

# ========================================================================== #
# Slide 8 — 根拠5本柱
# ========================================================================== #
s = content_slide("判定法が機能する根拠 — 5本柱の独立検証", tag="8/13", accent=PURPLE)
table(s, [
    ("柱", "検証", "結果", "ひとことで"),
    ("① (A2)", "負の対照（ニセ機会5種）", "真0.39が5/5対照のヌル分布外", "プラセボでは効かず本物でだけ効く"),
    ("② (A3)", "応答起点の比較", "END汚染0 vs START 0.56", "採点はテスト後から"),
    ("③ (B)", "応答ハザード", "72hで0.39 vs 疑似0.29", "本物の機会の後ほど早く計り直す"),
    ("④ (E)", "欠測・観測途切れの注入", "誤無応答 naive 643→4.1（99%減）", "データの穴を冤罪にしない"),
    ("⑤ (D)", "保持グリッド7〜90日", "stateful一致1.00・容量比4.2%", "30日のメモで全履歴と同じ結論"),
], Inches(0.45), Inches(1.1), Inches(7.85), Inches(3.7),
      col_w=[Inches(0.85), Inches(2.1), Inches(2.7), Inches(2.2)], font=11)
bullets(s, [
    "いずれも proxyラベルに依存しない独立検証",
    "9項目すべて支持・ベースライン再現ゲート 16/16 通過",
], Inches(0.5), Inches(4.95), Inches(7.6), Inches(1.0), size=12.5, gap=4)
img_a2 = FIG / "negative_control_true_vs_null.png"
img_d = FIG / "retention_invariance_heatmap.png"
if img_a2.exists():
    s.shapes.add_picture(str(img_a2), EM(Inches(8.5)), EM(Inches(1.1)),
                         height=EM(Inches(2.35)))
    label(s, Inches(8.5), Inches(3.42), Inches(4.4), Inches(0.3),
          "①A2: 本物の機会だけがヌル分布の外", size=10, color=GREY, align=PP_ALIGN.LEFT)
if img_d.exists():
    s.shapes.add_picture(str(img_d), EM(Inches(8.5)), EM(Inches(3.75)),
                         height=EM(Inches(2.2)))
    label(s, Inches(8.5), Inches(5.98), Inches(4.4), Inches(0.3),
          "⑤D: statefulなら保持7〜90日で一致1.00", size=10, color=GREY, align=PP_ALIGN.LEFT)
hint(s, "プラセボ対照・採点タイミング・欠測耐性・30日等価 — 治験並みの検証を通した。")
notes(s, "1:05",
      "この判定を信じてよい根拠が5本柱の検証です。第一に負の対照。ニセの機会5種類と比べ、本物の"
      "機会の後でだけ応答率が高い。プラセボでは効かず本物で効く、治験と同じ構図です。第二・第三に、"
      "採点の起点をENDに置けば汚染が構造的にゼロであること、そして本物の機会の後ほど早く計り直しが"
      "起きること。第四に欠測ストレス。欠測を18パターン人工注入しても、誤確定は素朴な方式の平均"
      "643件から4.1件、99%減に抑えます。第五に保持グリッド。生データを30日しか持たなくても、7要素の"
      "状態台帳があれば全期間監査と結論が完全一致します。これらはすべてproxyラベルに依存しない独立"
      "指標で、9項目すべて支持、再現ゲートも16件全て通過しています。")

# ========================================================================== #
# Slide 9 — 特許候補の全体像
# ========================================================================== #
s = content_slide("特許候補の全体像 — 2軸で見る", tag="9/13", accent=NAVY)
table(s, [
    ("特許候補", "技術エビデンス", "新規性リスク（UNVERIFIED）"),
    ("① 機会条件付き無応答監査（中核）", "STRONG", "MEDIUM-HIGH"),
    ("② デュアルトラック非対称リセット", "STRONG", "**HIGH（着想日依存）**"),
    ("③ 有界保持の因果証拠台帳", "STRONG", "MEDIUM-HIGH"),
    ("(将来④) クローズドループ介入検証", "PROSPECTIVE", "—"),
    ("(将来⑤) 機種非依存+version局在", "MEDIUM / PROSPECTIVE", "MEDIUM"),
], Inches(0.45), Inches(1.2), SW - Inches(0.9), Inches(3.6),
      col_w=[Inches(6.0), Inches(3.2), Inches(3.2)], font=13)
bullets(s, [
    "左列 = データで実証した強さ（済） / 右列 = 先行技術と重なる危険性（未検証・別軸）",
], Inches(0.5), Inches(5.0), SW - Inches(1.0), Inches(0.5), size=12.5)
hint(s, "『技術エビデンスが強い』と『特許が取れる』は別物 — 左右の列を分けて見るのがこの表の読み方。")
notes(s, "0:35",
      "ここからが特許の話です。候補は3件、将来候補が2件。この表で大事なのは2軸を分けて見ることです。"
      "左の『技術エビデンス』は我々がデータで実証した強さ、右の『新規性リスク』は先行技術と重なる危険性"
      "で、独立の軸です。3候補ともエビデンスはSTRONGですが、先行技術調査は未実施のため新規性は未検証、"
      "特に候補②はHIGHリスクです。順に見ていきます。")

# ========================================================================== #
# Slide 10 — 候補①
# ========================================================================== #
s = content_slide("特許候補① 機会条件付き無応答監査（中核）", tag="10/13", accent=BLUE)
bullets(s, [
    ("発明の骨格", {"bold": True, "color": NAVY, "size": 14}),
    "機会抽出 → END起点72h窓で 応答あり/無応答/保留 に分類 → 品質ティア（観測データの信頼度の段階分け）＋保留除外 → 二分岐トリアージ",
    ("裏付け", {"bold": True, "color": NAVY, "size": 14}),
    "①A2特異性（プラセボ対照） / ②A3汚染0（END起点） / ④E欠測耐性 — 5本柱のうち3本が直接支持",
    ("リスクと対策", {"bold": True, "color": NAVY, "size": 14}),
    "『非発生イベントの監視』という広い枠では先行技術多数 → 具体構成（RSOC 80/20/80・72h窓・50mWh閾値・保留除外明示）で狭め/中位クレームの出願を推奨",
    ("適用範囲", {"bold": True, "color": NAVY, "size": 14}),
    "必要特徴量は raw 4列のみ（RSOC/FCC/cycleCount/timestamp。HW識別子不使用）→ 機種非依存で適用範囲が広い",
], Inches(0.5), Inches(1.15), SW - Inches(1.0), Inches(4.9), size=13.5, gap=7)
notes(s, "0:50",
      "候補①が中核です。機会を抽出し、END起点72時間窓で応答・無応答・保留の3値に分類し、データ信頼度の"
      "品質ティアと二分岐トリアージまで含めた監査パイプライン全体を権利化する案です。入力は生テレメトリ"
      "4列だけで、機種非依存です。弱点は、『起きるはずのイベントが起きないことの監視』という広い枠では"
      "先行技術が多いこと。そこで80/20/80のRSOC帯、72時間窓、50ミリワット時閾値、保留除外の明示といった"
      "具体構成で狭めた、中位クレームでの出願を推奨します。")

# ========================================================================== #
# Slide 11 — 候補②③
# ========================================================================== #
s = content_slide("特許候補② 非対称リセット / 候補③ 因果証拠台帳", tag="11/13", accent=BLUE)
lshp = rect(s, Inches(0.45), Inches(1.1), Inches(6.15), Inches(4.8), WHITE, STEEL)
ltf = lshp.text_frame; ltf.word_wrap = True
ltf.margin_left = Inches(0.16); ltf.margin_top = Inches(0.12); ltf.vertical_anchor = MSO_ANCHOR.TOP
_set(ltf.paragraphs[0], "② デュアルトラック非対称リセット", size=14.5, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
for txt, kw in [
    ("FCC変化の記録を『全変化(any)』と『有効変化(effective, ≥50mWh)』の2冊の台帳に分離", {}),
    ("微小変化（<50mWh。全ステップの58.1%）では『全変化』側だけをリセットする非対称設計", {}),
    ("対称リセットは281ユーザーで保留1802+確定無応答462件の証拠を消去。非対称は0件", {"bold": True}),
    ("【リスク高】production実装済 → 新規性は着想日・公開有無に依存（本日の依頼②）", {"color": RED, "bold": True}),
]:
    p = ltf.add_paragraph(); _set(p, "・" + txt, size=12.5, bold=kw.get("bold", False),
                                  color=kw.get("color", INK))
    p.space_before = Pt(7)
rshp = rect(s, Inches(6.75), Inches(1.1), Inches(6.15), Inches(4.8), WHITE, STEEL)
rtf = rshp.text_frame; rtf.word_wrap = True
rtf.margin_left = Inches(0.16); rtf.margin_top = Inches(0.12); rtf.vertical_anchor = MSO_ANCHOR.TOP
_set(rtf.paragraphs[0], "③ 有界保持の因果証拠台帳", size=14.5, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
for txt, kw in [
    ("生データ30日保持でも、7要素の最小状態（部分FSM・pending期限・seen_ids・直近有効変化時刻・有効変化後サイクル数・gap/censorカウンタ・順序規則）を持ち回れば全期間と同一結論", {}),
    ("台帳なし7日保持: 一致率0.62 ／ 台帳あり: recall 1.0・重複0・ストレージ比4.2%（1176構成で検証）", {"bold": True}),
    ("全7要素の必要性をアブレーションで実証 — 『どの7項目か』の特定が発明", {}),
    ("【リスク】streaming+cachingは既知の組合せ（UNVERIFIED）→『最小状態構造』をクレーム核に", {"color": RED}),
]:
    p = rtf.add_paragraph(); _set(p, "・" + txt, size=12, bold=kw.get("bold", False),
                                  color=kw.get("color", INK))
    p.space_before = Pt(7)
hint(s, "②は家計簿の2冊分け、③は『レシートを30日で捨てても7項目の家計簿があれば同じ結論』。")
notes(s, "0:55",
      "候補②と③は中核を支える2つの工夫です。②は微小なFCC変化の扱いで、変化の記録を『全変化』と"
      "『有効変化』の2冊の台帳で別々に管理し、微小変化では『全変化』側だけをリセットする非対称設計"
      "です。対称にすると、281ユーザーで保留1802件と確定無応答462件、あわせて2200件を超える証拠が"
      "消えます。ただし②はすでにproduction実装済みのため、新規性が着想日と公開有無に依存する点が"
      "最大のリスクです。③は、生データを30日しか保持できない制約の下で、7項目の最小状態だけを"
      "持ち回れば全期間監査と同一の結論になる証拠台帳です。どの7項目が必要十分かを特定し、必要性まで"
      "実証した点が発明の核です。")

# ========================================================================== #
# Slide 12 — 正直な開示
# ========================================================================== #
s = content_slide("正直な開示 — 敵対的レビューと限界", tag="12/13", accent=GREY)
tshp = rect(s, Inches(0.45), Inches(1.1), SW - Inches(0.9), Inches(1.0), LGREEN_BG, DGREEN)
ttf = tshp.text_frame; ttf.word_wrap = True
ttf.margin_left = Inches(0.14); ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
_set(ttf.paragraphs[0],
     "検証の徹底: 8観点66指摘の敵対的レビュー → 52件は棄却（コード正当）/ 過大主張5+文書課題5を修正 / "
     "真のロジック欠陥0・PENDING 2件（追跡中） / 修正後テスト108件pass", size=13, bold=True, color=INK)
bullets(s, [
    ("それでも残る限界(隠さない):", {"bold": True, "color": NAVY, "size": 14}),
    "proxyラベルは地上真実ではない（『FW確認対象の抽出』であり『FW不良の検出』ではない）",
    "先行技術は全てUNVERIFIED。主要先行特許 US12061240 の独立クレームは未精読",
    "介入データ・BIOS/EC/FWバージョンは NOT AVAILABLE（存在しないものは無いと言う）",
    "候補②の新規性は着想日依存（v4作業は『特徴付け・検証』であり『着想』ではない）",
    "50mWhは事前指定値。データ（GMM谷35.2mWh・CI[26.3, 54.1]）は『矛盾しない』ことの裏付けで、最適性の証明ではない",
], Inches(0.5), Inches(2.35), SW - Inches(1.0), Inches(3.6), size=13, gap=7)
hint(s, "弱点は先に自分たちで潰した。それでも新規性の判断はここから先が本番。")
notes(s, "0:55",
      "正直な開示です。この結果は8観点66項目の敵対的レビューにかけ、52件は検証の結果棄却、過大主張5件と"
      "文書課題5件は修正済み、ロジックの欠陥はゼロでした。ただし未解決の指摘も2件残っています。修正後の"
      "テストは108件すべて通っています。その上で残る限界を申し上げます。先ほど来のproxyラベルは本番"
      "システムの推定値であり、地上真実ではありません。先行技術は全て未検証で、特に米国特許12061240の"
      "独立クレームは未精読です。介入データとFWバージョンのデータは存在しません。無いものは無いと"
      "申し上げた上で、判断をお願いするのが本資料の立場です。")

# ========================================================================== #
# Slide 13 — 知財戦略と依頼
# ========================================================================== #
s = content_slide("知財戦略と依頼事項 — 本日決めていただきたいこと", tag="13/13", accent=NAVY)
lshp = rect(s, Inches(0.45), Inches(1.1), Inches(6.3), Inches(4.7), WHITE, STEEL)
ltf = lshp.text_frame; ltf.word_wrap = True
ltf.margin_left = Inches(0.16); ltf.margin_top = Inches(0.12); ltf.vertical_anchor = MSO_ANCHOR.TOP
_set(ltf.paragraphs[0], "3枚のカード: 出願 / 防御的公開 / 営業秘密", size=14.5, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
for txt, kw in [
    ("防御的公開 = 自社では出願せず内容を公開して他社の特許化を防ぐ手", {"color": GREY, "size": 11.5}),
    ("推奨ルート（6月18日の特許性評価）: JP先願 → PCT → US/EP（総合判定『条件付きで出願に値する』）", {"bold": True}),
    ("候補①（狭め/中位）と候補③（最小状態を核に）が出願の核。品質ティア・保留除外などの細部構成は従属クレームで担保", {}),
    ("候補②は着想日次第で出願 or 防御的公開", {}),
    ("運用の微調整閾値は営業秘密として保持", {}),
]:
    p = ltf.add_paragraph()
    _set(p, "・" + txt, size=kw.get("size", 13), bold=kw.get("bold", False),
         color=kw.get("color", INK))
    p.space_before = Pt(8)
rshp = rect(s, Inches(6.9), Inches(1.1), Inches(6.0), Inches(4.7), LBLUE_BG, NAVY)
rtf = rshp.text_frame; rtf.word_wrap = True
rtf.margin_left = Inches(0.16); rtf.margin_top = Inches(0.12); rtf.vertical_anchor = MSO_ANCHOR.TOP
_set(rtf.paragraphs[0], "依頼（4点）", size=15, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
for i, txt in enumerate([
    "出願の粒度・範囲のご判断",
    "候補②の着想日・公開有無の確認",
    "正式FTO（他社特許を侵害しないかの調査）/ 特許性調査の発注判断",
    "クローズドループ実証のための介入データ収集開始の判断",
]):
    p = rtf.add_paragraph()
    _set(p, f"{i + 1}. {txt}", size=14, bold=True, color=INK)
    p.space_before = Pt(12)
hint(s, "出願・営業秘密・防御的公開の3枚のカードをどう切るかが、本日reviewerの皆さまへの問い。")
notes(s, "1:00",
      "最後に知財戦略と依頼です。取れる手は3つ。出願、防御的公開――自社では出願せず内容を公開して"
      "他社の特許化を防ぐ方法――、そして営業秘密です。推奨は、候補①を狭め・中位クレームで、候補③の"
      "最小状態台帳を核に、日本先願から国際出願PCT、米欧へ進むルートです。候補②は着想日と公開有無の"
      "確認結果次第で、出願か防御的公開かを選びます。微調整閾値は営業秘密として社内に留めます。"
      "お願いは4点です。第一に出願の粒度と範囲のご判断。第二に候補②の着想日と公開有無の確認。第三に、"
      "他社特許を侵害しないかを調べる正式なFTO調査と特許性調査の発注判断。第四に介入データ収集開始の"
      "ご判断です。ご審議をお願いいたします。")

# ========================================================================== #
prs.save(OUT)
n_notes = sum(1 for sl in prs.slides if sl.has_notes_slide)
total_chars = sum(
    len(sl.notes_slide.notes_text_frame.text.split("\n", 1)[1])
    for sl in prs.slides if sl.has_notes_slide)
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides._sldIdLst)} / notes: {n_notes} / "
      f"note body chars: {total_chars} (~{total_chars / 300:.1f} min @300字/分)")
