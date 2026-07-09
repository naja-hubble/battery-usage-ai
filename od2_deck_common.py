"""Shared slide primitives for the OD2 patent-review decks (mirrors build_v5_pptx.py).

A ``Deck`` wraps a python-pptx Presentation and exposes the exact visual grammar of the
OD1 v5 deck — section band, 💡 takeaway band, boxes, connectors, tables, glossary, the
patent-candidate ``concept`` slide (背景/問題/課題/解決法/特許性), the 4-block ``figure``
slide (何のグラフ/軸/主張/読み方), and the honesty footer — so the OD2 decks look identical
to the originals. 16:9, Meiryo UI with CJK run wiring.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.oxml.ns import qn

# ---- palette (identical to build_v5_pptx.py) ----
NAVY = RGBColor(0x1F, 0x37, 0x5B); BLUE = RGBColor(0x1F, 0x77, 0xB4)
STEEL = RGBColor(0x33, 0x66, 0x99); GREY = RGBColor(0x55, 0x55, 0x55)
RED = RGBColor(0xC0, 0x39, 0x2B); GREEN = RGBColor(0x2C, 0xA0, 0x2C)
DGREEN = RGBColor(0x37, 0x6B, 0x3A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7); ORANGE = RGBColor(0xB0, 0x5A, 0x00)
TEAL = RGBColor(0x2E, 0x6E, 0x6A); PURPLE = RGBColor(0x6A, 0x3D, 0x8A)
BROWN = RGBColor(0x6B, 0x4A, 0x1F); INK = RGBColor(0x22, 0x22, 0x22)
YELLOW_BG = RGBColor(0xFF, 0xF6, 0xDA); YELLOW_LN = RGBColor(0xD9, 0xA4, 0x2A)
LGREEN_BG = RGBColor(0xE4, 0xF0, 0xE4); LBLUE_BG = RGBColor(0xE2, 0xEC, 0xF7)
LRED_BG = RGBColor(0xF7, 0xE4, 0xE2); LGREY_BG = RGBColor(0xEC, 0xEC, 0xEC)

FONT = "Meiryo UI"
DISC = ("技術的特許性エビデンス（NOT a legal opinion）。先行技術は UNVERIFIED。"
        "特許化の粒度・範囲の最終判断は社内reviewer/弁理士。捏造（地上真実/介入/FW/因果結論）なし。")


def EM(v):
    return Emu(int(round(v)))


def _cjk(run):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", FONT)


def _style(run, size=14, bold=False, color=None):
    run.font.size = Pt(size); run.font.bold = bold; run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    _cjk(run)


def _tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(EM(left), EM(top), EM(width), EM(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size=14, bold=False, color=None, align=None):
    p.text = text
    for run in p.runs:
        _style(run, size, bold, color)
    if align is not None:
        p.alignment = align


class Deck:
    def __init__(self, out_path, fig_dir, disc=DISC):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333); self.prs.slide_height = Inches(7.5)
        self.SW, self.SH = self.prs.slide_width, self.prs.slide_height
        self.BLANK = self.prs.slide_layouts[6]
        self.FIG = fig_dir; self.OUT = out_path; self.DISC = disc

    def add(self):
        return self.prs.slides.add_slide(self.BLANK)

    def band(self, slide, color, height):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EM(self.SW), EM(height))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def footer(self, slide):
        tf = _tb(slide, Inches(0.3), self.SH - Inches(0.38), self.SW - Inches(1.1), Inches(0.32))
        _set(tf.paragraphs[0], self.DISC, size=8, color=GREY)

    def section(self, title, accent=BLUE, tag=None):
        s = self.add()
        self.band(s, accent, Inches(0.86))
        tf = _tb(s, Inches(0.5), Inches(0.12), self.SW - Inches(3.2), Inches(0.64))
        _set(tf.paragraphs[0], title, size=21, bold=True, color=WHITE)
        if tag:
            tt = _tb(s, self.SW - Inches(3.2), Inches(0.2), Inches(2.9), Inches(0.5))
            _set(tt.paragraphs[0], tag, size=11.5, bold=True,
                 color=RGBColor(0xDD, 0xE8, 0xF2), align=2)
        return s

    def takeaway(self, slide, text, y=None):
        y = Inches(6.5) if y is None else y
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, EM(Inches(0.35)), EM(y),
                                     EM(self.SW - Inches(0.7)), EM(Inches(0.52)))
        shp.fill.solid(); shp.fill.fore_color.rgb = YELLOW_BG
        shp.line.color.rgb = YELLOW_LN; shp.shadow.inherit = False
        tf = shp.text_frame; tf.word_wrap = True
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        _set(tf.paragraphs[0], "💡 " + text, size=11.5, bold=True, color=BROWN)

    def box(self, slide, left, top, w, h, text, fill, fg=WHITE, size=12, sub=None,
            sub_size=9.5, line=NAVY):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, EM(left), EM(top), EM(w), EM(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.color.rgb = line
        shp.shadow.inherit = False
        tf = shp.text_frame; tf.word_wrap = True
        _set(tf.paragraphs[0], text, size=size, bold=True, color=fg, align=2)
        if sub:
            p = tf.add_paragraph(); _set(p, sub, size=sub_size, color=fg, align=2)
        return shp

    def line_seg(self, slide, x1, y1, x2, y2, color=NAVY, width=2.25, dash=None, arrow_head=False):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, EM(x1), EM(y1), EM(x2), EM(y2))
        c.line.color.rgb = color; c.line.width = Pt(width)
        if dash is not None:
            c.line.dash_style = dash
        if arrow_head:
            ln = c.line._get_or_add_ln()
            tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
            ln.append(tail)
        return c

    def label(self, slide, x, y, w, h, text, size=10.5, bold=False, color=INK, align=None):
        tf = _tb(slide, x, y, w, h)
        _set(tf.paragraphs[0], text, size=size, bold=bold, color=color, align=align)
        return tf

    def bullets(self, tf, items, size=12.5, gap=4, color=INK, first=True, marker="•  "):
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            _set(p, (marker + it) if marker else it, size=size, color=color)
            p.space_after = Pt(gap)
        return tf

    def _fill_cell(self, c, text, size=10, bold=False, color=INK, fill=None):
        c.text = text
        if fill is not None:
            c.fill.solid(); c.fill.fore_color.rgb = fill
        for pa in c.text_frame.paragraphs:
            for run in pa.runs:
                _style(run, size, bold, color)

    def make_table(self, slide, headers, data, x, y, colw, row_h=0.42, hdr_size=11,
                   cell_size=10, styler=None):
        n = len(data) + 1
        tbl = slide.shapes.add_table(n, len(headers), EM(x), EM(y), EM(Inches(sum(colw))),
                                     EM(Inches(row_h) * n)).table
        for j, w in enumerate(colw):
            tbl.columns[j].width = EM(Inches(w))
        for j, h in enumerate(headers):
            self._fill_cell(tbl.cell(0, j), h, size=hdr_size, bold=True, color=WHITE, fill=NAVY)
        for i, rowv in enumerate(data, 1):
            for j, v in enumerate(rowv):
                size, bold, color = cell_size, False, INK
                if styler:
                    size, bold, color = styler(i, j, v, cell_size)
                self._fill_cell(tbl.cell(i, j), v, size=size, bold=bold, color=color,
                                fill=LIGHT if i % 2 else WHITE)
        return tbl

    def glossary(self, title, items, tag="用語定義"):
        s = self.section(title, STEEL, tag)
        tf = _tb(s, Inches(0.55), Inches(1.0), self.SW - Inches(1.1), Inches(5.3))
        first = True
        for term, defn in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.space_after = Pt(4)
            r = p.add_run(); r.text = f"{term}　"; _style(r, 12.5, True, NAVY)
            r2 = p.add_run(); r2.text = defn; _style(r2, 11.5, False, INK)
        self.footer(s)
        return s

    def concept(self, title, tag, strength, risk, bg, prob, chal, sol, novelty, tw, accent=TEAL):
        s = self.section(title, accent, tag)
        tf = _tb(s, Inches(0.5), Inches(0.94), self.SW - Inches(1.0), Inches(0.4))
        p = tf.paragraphs[0]; p.text = ""
        for txt, col in [("技術エビデンス強度: ", NAVY),
                         (strength + "    ", GREEN if "STRONG" in strength
                          else (GREY if "PROSPECTIVE" in strength else BLUE)),
                         ("先行技術リスク(UNVERIFIED): ", NAVY),
                         (risk, RED if "HIGH" in risk else GREY)]:
            r = p.add_run(); r.text = txt; _style(r, 12, True, col)
        tf = _tb(s, Inches(0.5), Inches(1.36), self.SW - Inches(1.0), Inches(5.0))
        first = True
        for head, body in [("背景", bg), ("問題", prob), ("課題", chal),
                           ("解決法", sol), ("特許性 / 新規性", novelty)]:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            _set(p, head, size=13, bold=True, color=NAVY); p.space_before = Pt(4)
            q = tf.add_paragraph(); _set(q, body, size=11.8)
        self.takeaway(s, tw); self.footer(s)
        return s

    def features(self, title, tag, rows, row_h=0.44):
        s = self.section(title, RGBColor(0x2E, 0x5A, 0x88), tag)
        self.make_table(s, ["変数 / 特徴量", "種別", "作成方法（特徴量エンジニアリング）"],
                        rows, Inches(0.4), Inches(1.02), [4.6, 1.5, 6.4], row_h=row_h,
                        styler=lambda i, j, v, b: (
                            b, j == 1,
                            (GREEN if v == "派生" else RED if "NOT" in v else GREY)
                            if j == 1 else INK))
        self.footer(s)
        return s

    def figure(self, title, tag, img, what, axes, claim, interp, tw=None, accent=DGREEN,
               max_w=7.1, max_h=4.95):
        s = self.section(title, accent, tag)
        p = self.FIG / img
        img_bottom = Inches(1.1)
        if p.exists():
            from PIL import Image
            try:
                iw, ih = Image.open(p).size
            except Exception:
                iw, ih = 1600, 1000
            scale = min(Inches(max_w) / iw, Inches(max_h) / ih)
            pic_w, pic_h = EM(iw * scale), EM(ih * scale)
            s.shapes.add_picture(str(p), EM(Inches(0.4)), EM(Inches(1.08)), width=pic_w, height=pic_h)
            img_bottom = Inches(1.08) + pic_h
            cap = _tb(s, Inches(0.4), img_bottom + Inches(0.02), Inches(7.1), Inches(0.3))
            _set(cap.paragraphs[0], f"図: {img}（dpi=300・匿名化済み）", size=8.5, color=GREY)
        tf = _tb(s, Inches(7.7), Inches(1.0), Inches(5.3), Inches(5.45))
        first = True
        for head, body, col in [("【何のグラフか】", what, BLUE),
                                ("【軸・変数の定義】", axes, NAVY),
                                ("【主張（言いたいこと）】", claim, GREEN),
                                ("【読み方・解釈】", interp, ORANGE)]:
            p1 = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            _set(p1, head, size=12, bold=True, color=col); p1.space_before = Pt(5)
            p2 = tf.add_paragraph(); _set(p2, body, size=11, color=INK)
        if tw:
            self.takeaway(s, tw)
        self.footer(s)
        return s

    def title_slide(self, title, subtitle, bullets_list, size=33):
        s = self.add(); self.band(s, NAVY, self.SH)
        tf = _tb(s, Inches(0.8), Inches(1.2), self.SW - Inches(1.6), Inches(2.2))
        _set(tf.paragraphs[0], title, size=size, bold=True, color=WHITE)
        p = tf.add_paragraph(); _set(p, subtitle, size=18, color=RGBColor(0xCF, 0xDD, 0xEE))
        tf2 = _tb(s, Inches(0.8), Inches(3.5), self.SW - Inches(1.6), Inches(3.4))
        for i, ln in enumerate(bullets_list):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            _set(p, "•  " + ln, size=14, color=WHITE); p.space_after = Pt(5)
        self.footer(s)
        return s

    def flowchart(self, title="処理フロー — データ取得から判定まで", tag="§3 Part1"):
        """Pipeline flowchart. The triage node ⑥ genuinely BIFURCATES the candidates
        into FW-check vs gauge-reset (exactly 2 arrows). NORMAL/REVIEW/WATCH are
        non-actionable/hold states emitted by the pre-triage gate, drawn separately."""
        s = self.section(title, NAVY, tag)
        GOLD = RGBColor(0xC8, 0x9B, 0x2A)
        xs = [0.35, 2.85, 5.35, 7.85, 10.35]
        W = 2.35
        top = [
            ("① データ取得", "S3テレメトリ: RSOC/FCC/サイクル/chargeStatus/時刻（HW識別子不使用）", STEEL),
            ("② 前処理・品質ゲート", "時刻整列・重複除去・OK/SPARSE/COUNTER_RESET判定", STEEL),
            ("③ 学習機会の抽出（2機構）", "Type A 満充電→≤6%→満充電 / Type B 充電中60-80%通過→満充電（END=満充電）", BLUE),
            ("④ END起点168h 応答監査", "有効ステップ≥50mWh → responded / no_response / censored", TEAL),
            ("⑤ 品質ティア＋censored除外", "HIGH_OK/MEDIUM_GAP/LOW_LARGE_GAP・打切りは無応答に数えない", DGREEN),
        ]
        ty, th = 1.05, 0.92
        for i, (t, sub, col) in enumerate(top):
            self.box(s, Inches(xs[i]), Inches(ty), Inches(W), Inches(th), t, col,
                     size=10.3, sub=sub, sub_size=8.1)
            if i < 4:
                self.line_seg(s, Inches(xs[i] + W), Inches(ty + th / 2),
                              Inches(xs[i + 1]), Inches(ty + th / 2),
                              color=GREY, width=2.0, arrow_head=True)
        # --- counter-operation annotation for ④ (increment / reset / hold) ---
        self.box(s, Inches(0.4), Inches(2.08), Inches(3.2), Inches(1.12), " ", LIGHT, fg=NAVY, line=NAVY)
        tfa = _tb(s, Inches(0.55), Inches(2.13), Inches(2.95), Inches(1.02))
        _set(tfa.paragraphs[0], "④ 応答判定 → カウンタ操作", size=9.8, bold=True, color=NAVY)
        for txt, col in [("no_response → カウンタ +1", RED),
                         ("responded(≥50mWh) → 0 にリセット", GREEN),
                         ("censored → pending（保留・数えない）", GREY)]:
            q = tfa.add_paragraph(); _set(q, txt, size=9, color=col); q.space_after = Pt(1)
        # --- L-route from ⑤ into the pre-triage GATE ---
        cx5 = xs[4] + W / 2; gx = 6.65
        self.line_seg(s, Inches(cx5), Inches(ty + th), Inches(cx5), Inches(2.12), color=NAVY, width=2.0)
        self.line_seg(s, Inches(cx5), Inches(2.12), Inches(gx), Inches(2.12), color=NAVY, width=2.0)
        self.line_seg(s, Inches(gx), Inches(2.12), Inches(gx), Inches(2.32), color=NAVY, width=2.0, arrow_head=True)
        # --- pre-triage gate (decides candidate vs hold states) ---
        self.box(s, Inches(3.7), Inches(2.32), Inches(5.9), Inches(0.6),
                 "前段ゲート：データ十分か / FCC更新継続(候補外)か / 境界か", STEEL, size=11)
        # --- gate -> HOLD states (right column; non-actionable) ---
        self.label(s, Inches(9.75), Inches(2.44), Inches(1.4), Inches(0.3), "→ 非アクション/保留", size=9, color=GREY)
        self.line_seg(s, Inches(9.6), Inches(2.62), Inches(10.35), Inches(3.05), color=GREY, width=1.6, arrow_head=True)
        holds = [("NORMAL", "候補外＝FCC健全群並みに更新", STEEL, WHITE),
                 ("REVIEW", "obs<120d / n<200 / カウンタ・パック異常", GREY, WHITE)]
        for i, (t, sub, col, fg) in enumerate(holds):
            self.box(s, Inches(10.35), Inches(3.05 + i * 0.66), Inches(2.5), Inches(0.58),
                     t, col, fg=fg, size=10.5, sub=sub, sub_size=7.8)
        # --- gate -> ⑥ TWO-WAY triage (candidates) ---
        self.label(s, Inches(4.35), Inches(3.02), Inches(3.2), Inches(0.28),
                   "no/low-change候補", size=9.5, color=NAVY)
        self.line_seg(s, Inches(5.4), Inches(2.92), Inches(3.65), Inches(3.5), color=NAVY, width=2.0, arrow_head=True)
        self.box(s, Inches(1.4), Inches(3.5), Inches(4.5), Inches(0.62),
                 "⑥ 二分岐トリアージ（機構別k: A≥3 / B≥5）", NAVY, size=11.5)
        # exactly TWO branches from ⑥
        self.line_seg(s, Inches(3.65), Inches(4.12), Inches(2.1), Inches(4.55), color=NAVY, width=2.0, arrow_head=True)
        self.line_seg(s, Inches(3.65), Inches(4.12), Inches(4.7), Inches(4.55), color=NAVY, width=2.0, arrow_head=True)
        self.box(s, Inches(0.9), Inches(4.55), Inches(2.4), Inches(0.95), "FW確認候補", RED,
                 size=11.5, sub="機会あり×無応答", sub_size=8.6)
        self.box(s, Inches(3.5), Inches(4.55), Inches(2.4), Inches(0.95), "ゲージ再較正候補", ORANGE,
                 size=11.5, sub="学習機会が皆無", sub_size=8.6)
        # residual of ⑥: candidates that fit NEITHER FW nor gauge -> WATCH (hold)
        self.label(s, Inches(6.05), Inches(4.02), Inches(3.0), Inches(0.28),
                   "どちらも確信できず→保留", size=9, color=GREY)
        self.line_seg(s, Inches(5.9), Inches(3.9), Inches(7.8), Inches(4.55), color=GREY, width=1.6, arrow_head=True)
        self.box(s, Inches(6.6), Inches(4.55), Inches(2.4), Inches(0.95), "WATCH", GOLD, fg=INK,
                 size=11.5, sub="境界/large-gap/完全窓不足", sub_size=8.4)
        # --- ⑦ bounded-retention ledger strip ---
        self.box(s, Inches(0.35), Inches(5.72), Inches(12.35), Inches(0.6),
                 "⑦ 有界保持の因果証拠台帳（最小状態: FSM/pending期限/seen_ids/直近有効変化/gap-censor/順序）"
                 "→ 直近30日の生データだけで全期間と等価な判定（30日オンライン運用版 = 9段ラベル）",
                 LIGHT, fg=NAVY, size=10.3, line=NAVY)
        self.takeaway(s, "前段ゲートで REVIEW(データ不足)/NORMAL(候補外) を除外 → 候補は⑥で『FW確認 ⇄ ゲージ再較正』に二分岐。どちらも確信できない残差は WATCH で保留。")
        self.footer(s)
        return s

    def k_basis(self, tag="§3 Part1"):
        """The empirical justification for the mechanism-specific FW threshold k."""
        s = self.section("FW閾値 k の根拠 — 機構別に『誤警報5%以下』で較正", DGREEN, tag)
        tf = _tb(s, Inches(0.55), Inches(1.05), self.SW - Inches(1.1), Inches(0.9))
        _set(tf.paragraphs[0],
             "設計方針: 健全なゲージが偶然 k 回連続で無応答になる確率 (1−p)^k を 5% 以下にする（p=健全機の応答率）。"
             "弱い/高頻度のトリガほど誤って無応答が出やすいので、確信に必要な k は大きくなる。",
             size=12.5, bold=True, color=NAVY)
        self.make_table(s, ["機構", "健全応答率 p@168h", "(1−p)^k の計算", "採用 k"],
            [["Type A（深放電, 強いトリガ）", "0.74", "0.26² = 6.8%(NG) / 0.26³ = 1.8%(OK)", "3"],
             ["Type B（充電側, 弱い・高頻度）", "0.45", "0.55⁴ = 9.2%(NG) / 0.55⁵ = 5.0%(OK)", "5"]],
            Inches(0.5), Inches(2.1), [3.7, 2.6, 4.6, 1.4], row_h=0.62, cell_size=11,
            styler=lambda i, j, v, b: (b, j == 3, DGREEN if j == 3 else INK))
        tf2 = _tb(s, Inches(0.55), Inches(3.9), self.SW - Inches(1.1), Inches(2.4))
        self.bullets(tf2, [
            "健全応答率 p は active-reference（obs≥180d・cycle≥20・凍結<60日・品質OK）から実測（Type A 0.74 / Type B 0.45）",
            "Type B は健全でも55%が無応答なので、k=4では誤警報9.2%と過大 → k=5で5.0%に収める",
            "k は固定値ではなく調整可能（§7 代替実施形態）。窓・有効ステップ閾値と同じく感度解析対象",
            "運用の確定(FW_CORE)は k 単独ではなく『凍結日数・サイクル数・データ品質OK』も併せて要件化"], size=12)
        self.takeaway(s, "kは恣意的でなく『健全機の誤警報5%以下』から機構別に導出。弱いトリガ(Type B)ほど k を大きくする。")
        self.footer(s)
        return s

    def thresholds_slide(self, tag="§7 付録"):
        """One-page reference of every decision threshold (gate / candidate / FW / gauge)."""
        s = self.section("付録 — 判定閾値一覧（前段ゲート・候補判定・トリアージ）", GREY, tag)
        self.make_table(s, ["ラベル / 段", "判定条件（各行内はAND。『いずれか』はOR）"],
            [["REVIEW（保留）", "obs<120日 ／ n_samples<200 ／ cycleCount減少・serialNumber変化 の いずれか"],
             ["候補判定", "fcc_changes=0&obs≥120 ／ flat_tail≥180日 ／ cycle≥50&更新率≤4.24 ／ obs≥180&更新率≤1.41 の いずれか"],
             ["NORMAL", "候補条件に一つも該当しない（FCCが健全群並みに更新）"],
             ["FW確認", "候補 & flat_tail≥180日 & tail_cycle≥30 & 品質OK & (Type B無応答≥5 or Type A無応答≥3)"],
             ["ゲージ再較正", "候補 & flat_tail≥120日 & 学習機会ゼロ(A・Bとも) & 使用ゲート & 品質OK"],
             ["WATCH（保留）", "候補だがFW・ゲージのどちらにも確信できない残差（境界 / large-gap / 完全窓不足）"]],
            Inches(0.4), Inches(1.0), [2.3, 10.5], row_h=0.56, cell_size=9.6)
        py = 4.9
        self.box(s, Inches(0.4), Inches(py), Inches(12.5), Inches(1.35), " ", LIGHT, fg=NAVY, line=NAVY)
        self.label(s, Inches(0.6), Inches(py + 0.06), Inches(6), Inches(0.3), "■ 共通パラメータ",
                   size=11, bold=True, color=NAVY)
        tf = _tb(s, Inches(0.6), Inches(py + 0.38), self.SW - Inches(1.2), Inches(0.92))
        self.bullets(tf, [
            "満充電(END): RSOC≥99% ／ Type A 深放電: RSOC≤6% ／ Type B: 充電中に60-80%を通過(abort<60)",
            "有効ステップ: |ΔFCC|≥50mWh ／ 応答窓: 168h(主, 24/72hは副) ／ 品質ティア: HIGH_OK≤12h・MEDIUM_GAP≤24h",
            "FW閾値k: A=3 / B=5（健全応答0.74/0.45から誤警報5%以下で導出）｜ p05(健全群): 4.24(/100cyc)・1.41(/100d), active-reference n=214",
            "※ 更新率の単位は /100cyc・/100d。窓・各閾値は代替実施形態として調整可能（一部の微調整値は営業秘密として非公開）"],
            size=9.8, color=INK, gap=2)
        self.takeaway(s, "全ラベルは透明なしきい値ルール（機種名は不使用）。数値は本解析の運用値で、窓・閾値は調整可能な代替実施形態。")
        self.footer(s)
        return s

    def quality_tier_slide(self, tag="§7 付録"):
        """Appendix detail of step ⑤: graded quality tier + censored exclusion."""
        s = self.section("付録 — 品質ティア＋censored除外（⑤の判定詳細）", GREY, tag)
        tf = _tb(s, Inches(0.55), Inches(1.02), self.SW - Inches(1.1), Inches(0.95))
        self.bullets(tf, [
            "目的: データの穴(大ギャップ)・打切り(未観測窓)を『無応答』と誤計上しない（根拠E: 誤無応答 204→5, 約40x削減）",
            "品質スコア = 0.45×(最大ギャップ成分) + 0.35×(観測カバレッジ) + 0.20×(端点ギャップ成分)（重み和=1.0）"
            "｜ 最大ギャップ成分: ≤12h→1.0 / 12-24hで1.0→0.5 / 24-48hで0.5→0.0 / >48h→0"], size=11, gap=3)
        self.make_table(s, ["品質ティア", "条件", "no_response 計上"],
            [["HIGH_OK", "max_gap ≤ 12h かつ score ≥ 0.80", "✅ 可（FW_CORE/GAUGE_CORE も支持）"],
             ["MEDIUM_GAP", "max_gap ≤ 24h かつ score ≥ 0.50", "✅ 可（FW_WATCH。FW_CORE単独は不可）"],
             ["LOW_LARGE_GAP", "それ以外", "❌ 曖昧のみ・数えない"],
             ["INVALID", "FCC/RSOC欠損・順序異常", "❌ 除外"]],
            Inches(0.5), Inches(2.05), [2.6, 4.9, 5.3], row_h=0.56, cell_size=10.5,
            styler=lambda i, j, v, b: (b, j == 2,
                                       (GREEN if "✅" in v else RED if "❌" in v else INK) if j == 2 else INK))
        tf2 = _tb(s, Inches(0.55), Inches(4.85), self.SW - Inches(1.1), Inches(1.5))
        self.bullets(tf2, [
            "censored除外: 応答窓 END+168h が最終観測サンプルより後 → censored → pending(保留)・数えない（END+168h ≤ 最終観測 なら no_response 確定可）",
            "計上ルール: no_response(+1) は『(HIGH_OK または MEDIUM_GAP) かつ 非censored』の機会だけ。それ以外は保留（加算しない）",
            "※ 全履歴監査版は簡易に最大ギャップ 12h/24h のみでティア分け（オンライン版はスコア併用）。結論は同じ"], size=10.5, gap=3)
        self.takeaway(s, "データの穴と打切りを冤罪にしない二重の関門。数えるのは『十分サンプリング(≤12/24h)＋完全観測(非censored)』の機会のみ。")
        self.footer(s)
        return s

    def detail_flow_audit(self, title="詳細フロー① ④⑤ — 機会ごとの判定と『Type A / Type B 別カウント』", tag="§3 Part1"):
        """Per-opportunity flow for ④(response)+⑤(quality/censored) — emphasises that the
        no-response count is kept SEPARATELY per mechanism (Type A vs Type B)."""
        s = self.section(title, NAVY, tag)
        GOLD = RGBColor(0xC8, 0x9B, 0x2A)
        # header
        self.box(s, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.58),
                 "学習機会の END（満充電到達）— この機会は Type A（深放電サイクル）か Type B（充電側部分再学習）のどちらか（③で確定済み）",
                 STEEL, size=11)
        self.line_seg(s, Inches(5.0), Inches(1.66), Inches(5.0), Inches(1.92), color=NAVY, width=2.0, arrow_head=True)
        # shared judgment (same for A and B)
        self.box(s, Inches(0.5), Inches(1.92), Inches(9.0), Inches(0.92),
                 "共通判定（A・B 同じロジック）  ⑤品質ティア(HIGH_OK / MEDIUM_GAP のみ有効) → "
                 "⑤完全観測(END+168h ≤ 最終観測サンプル = 非censored) → ④ 窓[END, END+168h] に 有効ステップ |ΔFCC|≥50mWh か？",
                 DGREEN, size=9.6)
        self.box(s, Inches(9.65), Inches(1.98), Inches(3.15), Inches(0.8),
                 "LOW_LARGE_GAP / INVALID / censored → 保留（どちらのカウンタにも数えない）", GOLD, fg=INK, size=9)
        # split into two mechanism lanes
        self.line_seg(s, Inches(2.72), Inches(2.84), Inches(2.72), Inches(3.32), color=NAVY, width=2.0, arrow_head=True)
        self.line_seg(s, Inches(7.32), Inches(2.84), Inches(7.32), Inches(3.32), color=NAVY, width=2.0, arrow_head=True)
        lanes = [
            (0.55, "Type A（深放電サイクル）の機会", BLUE, "Type A 無応答カウンタ", "→ FW判定に寄与: Type A 無応答 ≥ 3"),
            (5.15, "Type B（充電側部分再学習）の機会", TEAL, "Type B 無応答カウンタ", "→ FW判定に寄与: Type B 無応答 ≥ 5"),
        ]
        lw = 4.35
        for lx, hdr, hcol, ctr, fwnote in lanes:
            self.box(s, Inches(lx), Inches(3.32), Inches(lw), Inches(0.5), hdr, hcol, size=10.5)
            self.box(s, Inches(lx), Inches(3.95), Inches(lw), Inches(0.56), f"④ no_response → {ctr} +1", RED, size=10.5)
            self.box(s, Inches(lx), Inches(4.63), Inches(lw), Inches(0.56), "④ responded (≥50mWh) → 同カウンタ 0 にリセット", GREEN, size=9.8)
            self.box(s, Inches(lx), Inches(5.31), Inches(lw), Inches(0.5), fwnote, LBLUE_BG, fg=NAVY, size=9.8)
        self.takeaway(s, "★各機会は Type A / Type B のどちらか。共通の判定(⑤④)を通し、無応答なら『その機構』のカウンタに +1・応答なら 0 リセット。"
                       "A と B は別カウンタ・別閾値（FW: A≥3 / B≥5）で数える。")
        self.footer(s)
        return s

    def detail_flow_triage(self, title="詳細フロー② ⑥ — ユーザーごとの二分岐トリアージ（候補のみ）", tag="§3 Part1"):
        """Per-user triage cascade for step ⑥ (first matching rule wins), with thresholds."""
        s = self.section(title, NAVY, tag)
        GOLD = RGBColor(0xC8, 0x9B, 0x2A)
        mx, mw, ox, ow = 0.9, 7.0, 8.3, 4.6
        cx = mx + mw / 2

        def down(y1, y2, label="No ↓"):
            self.line_seg(s, Inches(cx), Inches(y1), Inches(cx), Inches(y2), color=NAVY, width=2.0, arrow_head=True)
            if label:
                self.label(s, Inches(cx - 1.15), Inches((y1 + y2) / 2 - 0.14), Inches(1.0), Inches(0.28),
                           label, size=9, color=GREY)

        def yes(box_bottom_y, out_y, text, col):
            self.line_seg(s, Inches(mx + mw), Inches(box_bottom_y - 0.34), Inches(ox), Inches(out_y + 0.28),
                          color=GREEN, width=1.6, arrow_head=True)
            self.box(s, Inches(ox), Inches(out_y), Inches(ow), Inches(0.56), text, col, size=10.5)

        self.box(s, Inches(mx), Inches(1.1), Inches(mw), Inches(0.52), "候補（前段ゲートで REVIEW/NORMAL は除外済）", STEEL, size=10.5)
        down(1.62, 1.85, "")
        ys = [1.85, 2.68, 3.51, 4.34]
        dh = 0.72
        decs = [
            ("① FW確認(high): flat_tail≥180日 & tail_cyc≥30 & 品質OK & (Type B無応答≥5 or Type A無応答≥3)", "FW確認候補（high）", RED),
            ("② ゲージ(high): flat_tail≥120日 & 学習機会ゼロ(A・Bとも) & 使用ゲート & 品質OK", "ゲージ再較正候補（high）", ORANGE),
            ("③ FW確認(medium): flat_tail≥120日 & tail_cyc≥20 & Type B無応答≥4", "FW確認候補（medium）", RED),
            ("④ ゲージ(medium): flat_tail≥60日 & 学習機会ゼロ & 使用ゲート", "ゲージ再較正候補（medium）", ORANGE),
        ]
        for i, (dtext, otext, ocol) in enumerate(decs):
            self.box(s, Inches(mx), Inches(ys[i]), Inches(mw), Inches(dh), dtext, NAVY, size=9.3)
            yes(ys[i] + dh, ys[i] + 0.08, otext, ocol)
            nexty = ys[i + 1] if i + 1 < len(ys) else 5.17
            down(ys[i] + dh, nexty, "No ↓")
        self.box(s, Inches(mx), Inches(5.17), Inches(mw), Inches(0.52), "上記いずれも不成立 → WATCH（保留）", GOLD, fg=INK, size=10.5)
        self.label(s, Inches(0.55), Inches(5.82), Inches(12.4), Inches(0.5),
                   "適用順: fw_high > gauge_high > fw_med > gauge_med > watch（上から順・最初の一致で確定）｜ "
                   "使用ゲート = tail_cycle<20 / tail_min_rsoc>20 / tail_rsoc_swing<60 / tail_ac≥0.80 のいずれか（mediumは <30 / >25 / <50 / ≥0.75）",
                   size=8.8, color=GREY)
        self.takeaway(s, "候補を上から順に判定し、最初に一致した規則で確定。FW=機会ありに無応答(機構別≥k)、ゲージ=学習機会ゼロ、どれも不成立ならWATCH。")
        self.footer(s)
        return s

    def notes(self, slide, budget, text):
        ns = slide.notes_slide
        ns.notes_text_frame.text = f"[目安 {budget}]\n{text}"

    def save(self):
        self.prs.save(str(self.OUT))
        n_notes = sum(1 for sl in self.prs.slides if sl.has_notes_slide)
        total = sum(len(sl.notes_slide.notes_text_frame.text.split("\n", 1)[-1])
                    for sl in self.prs.slides if sl.has_notes_slide)
        print(f"saved: {self.OUT}")
        print(f"slides: {len(self.prs.slides._sldIdLst)} / notes: {n_notes} / "
              f"note chars: {total} (~{total/300:.1f} min @300字/分)")
        return len(self.prs.slides._sldIdLst)


def status_styler(hl_cols):
    def f(i, j, v, base):
        bold, color = False, INK
        if j in hl_cols:
            bold = True
            if "STRONG" in v or "SUPPORTED" in v or "✅" in v:
                color = GREEN
            elif v.strip().startswith("HIGH") or "NOT" in v:
                color = RED
            elif "PROSPECTIVE" in v or "WEAK" in v or v == "—":
                color = GREY
            else:
                color = BLUE
        return base, bold, color
    return f
